#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT美化器模块
处理未填充占位符的清理和重新排版
"""

import re
from typing import Dict, List, Any, Tuple, TYPE_CHECKING
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from logger import get_logger
from config import get_config

if TYPE_CHECKING:
    from pptx.presentation import Presentation
else:
    from pptx import Presentation

class PPTBeautifier:
    """PPT美化器"""

    def __init__(self, presentation: Presentation):
        self.presentation = presentation
        self.logger = get_logger()
        self.config = get_config()

    def cleanup_and_beautify(self, filled_placeholders: Dict[str, Any]) -> Dict[str, Any]:
        """
        清理未填充的占位符并美化布局
        
        Args:
            filled_placeholders: 已填充的占位符信息
            
        Returns:
            Dict: 清理和美化结果
        """
        results = {
            'removed_placeholders': [],
            'reorganized_slides': [],
            'layout_changes': []
        }
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_result = self._process_slide(slide, slide_idx, filled_placeholders)
            
            if slide_result['removed_count'] > 0:
                results['removed_placeholders'].append({
                    'slide_index': slide_idx,
                    'removed_count': slide_result['removed_count'],
                    'removed_placeholders': slide_result['removed_placeholders']
                })
            
            if slide_result['reorganized']:
                results['reorganized_slides'].append({
                    'slide_index': slide_idx,
                    'layout_change': slide_result['layout_change']
                })
                results['layout_changes'].append(slide_result['layout_change'])
        
        return results
    
    def _process_slide(self, slide, slide_idx: int, filled_placeholders: Dict[str, Any]) -> Dict[str, Any]:
        """
        处理单个幻灯片
        
        Args:
            slide: 幻灯片对象
            slide_idx: 幻灯片索引
            filled_placeholders: 已填充的占位符信息
            
        Returns:
            Dict: 处理结果
        """
        result = {
            'removed_count': 0,
            'removed_placeholders': [],
            'reorganized': False,
            'layout_change': None
        }
        
        # 找到所有包含占位符的文本框 - 识别所有{}格式的占位符
        placeholder_shapes = []
        filled_shapes = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                placeholder_matches = re.findall(r'\{([^}]+)\}', shape.text)
                if placeholder_matches:
                    # 检查是否已被填充
                    is_filled = any(
                        placeholder in filled_placeholders.get(str(slide_idx), {}) 
                        for placeholder in placeholder_matches
                    )
                    
                    if is_filled:
                        filled_shapes.append(shape)
                    else:
                        placeholder_shapes.append((shape, placeholder_matches))
        
        # 删除未填充的占位符
        for shape, placeholders in placeholder_shapes:
            try:
                # 记录删除的占位符
                result['removed_placeholders'].extend(placeholders)
                result['removed_count'] += 1
                
                # 从幻灯片中删除形状
                slide.shapes.element.remove(shape.element)
                
                self.logger.info(f"删除幻灯片 {slide_idx+1} 中的未填充占位符: {placeholders}")
                
            except Exception as e:
                self.logger.error(f"删除占位符时出错: {e}")
        
        # 重新排版剩余的形状
        if filled_shapes and result['removed_count'] > 0:
            layout_change = self._reorganize_shapes(slide, filled_shapes)
            if layout_change:
                result['reorganized'] = True
                result['layout_change'] = layout_change
        
        return result
    
    def _reorganize_shapes(self, slide, shapes: List) -> Dict[str, Any]:
        """
        重新组织形状的布局
        
        Args:
            slide: 幻灯片对象
            shapes: 需要重新排版的形状列表
            
        Returns:
            Dict: 布局变化信息
        """
        if not shapes:
            return {}
        
        try:
            # 获取幻灯片尺寸
            slide_width = float(self.presentation.slide_width or 0)
            slide_height = float(self.presentation.slide_height or 0)
            
            # 计算可用区域（排除标题区域）
            margins = self.config.layout_margins
            available_width = slide_width - Inches(margins['slide_margin_left'] + margins['slide_margin_right'])
            available_height = slide_height - Inches(margins['slide_margin_top'] + margins['slide_margin_bottom'])
            available_top = Inches(margins['slide_margin_top'])
            available_left = Inches(margins['slide_margin_left'])
            
            shape_count = len(shapes)
            
            if shape_count <= 4:
                # 4个或更少，使用2x2布局
                layout_info = self._arrange_2x2_layout(
                    shapes, available_left, available_top, available_width, available_height
                )
            elif shape_count <= 6:
                # 5-6个，使用2x3布局
                layout_info = self._arrange_2x3_layout(
                    shapes, available_left, available_top, available_width, available_height
                )
            else:
                # 更多的使用3x3布局
                layout_info = self._arrange_3x3_layout(
                    shapes, available_left, available_top, available_width, available_height
                )
            
            self.logger.info(f"重新排版 {shape_count} 个形状，使用 {layout_info['layout_type']} 布局")
            
            return layout_info
            
        except Exception as e:
            self.logger.error(f"重新排版时出错: {e}")
            return {}
    
    def _arrange_2x2_layout(self, shapes: List, left: float, top: float, width: float, height: float) -> Dict[str, Any]:
        """
        2x2布局排列
        """
        spacing = Inches(self.config.layout_margins['shape_spacing'])
        shape_width = width / 2 - spacing
        shape_height = height / 2 - spacing
        
        half_spacing = spacing / 2
        positions = [
            (left, top),  # 左上
            (left + width/2 + half_spacing, top),  # 右上
            (left, top + height/2 + half_spacing),  # 左下
            (left + width/2 + half_spacing, top + height/2 + half_spacing)  # 右下
        ]
        
        for i, shape in enumerate(shapes[:4]):
            if i < len(positions):
                x, y = positions[i]
                shape.left = int(x)
                shape.top = int(y)
                shape.width = int(shape_width)
                shape.height = int(shape_height)
                
                # 调整字体大小以适应新的形状大小
                self._adjust_text_size(shape, shape_width, shape_height)
        
        return {
            'layout_type': '2x2',
            'shape_count': len(shapes),
            'positions': positions[:len(shapes)]
        }
    
    def _arrange_2x3_layout(self, shapes: List, left: float, top: float, width: float, height: float) -> Dict[str, Any]:
        """
        2x3布局排列
        """
        spacing = Inches(self.config.layout_margins['shape_spacing'])
        shape_width = width / 3 - spacing
        shape_height = height / 2 - spacing
        
        positions = []
        for row in range(2):
            for col in range(3):
                x = left + col * (width/3 + spacing/3)
                y = top + row * (height/2 + spacing/2)
                positions.append((x, y))
        
        for i, shape in enumerate(shapes[:6]):
            if i < len(positions):
                x, y = positions[i]
                shape.left = int(x)
                shape.top = int(y)
                shape.width = int(shape_width)
                shape.height = int(shape_height)
                
                self._adjust_text_size(shape, shape_width, shape_height)
        
        return {
            'layout_type': '2x3',
            'shape_count': len(shapes),
            'positions': positions[:len(shapes)]
        }
    
    def _arrange_3x3_layout(self, shapes: List, left: float, top: float, width: float, height: float) -> Dict[str, Any]:
        """
        3x3布局排列
        """
        spacing = Inches(self.config.layout_margins['shape_spacing'])
        shape_width = width / 3 - spacing
        shape_height = height / 3 - spacing
        
        positions = []
        for row in range(3):
            for col in range(3):
                x = left + col * (width/3 + spacing/3)
                y = top + row * (height/3 + spacing/3)
                positions.append((x, y))
        
        for i, shape in enumerate(shapes[:9]):
            if i < len(positions):
                x, y = positions[i]
                shape.left = int(x)
                shape.top = int(y)
                shape.width = int(shape_width)
                shape.height = int(shape_height)
                
                self._adjust_text_size(shape, shape_width, shape_height)
        
        return {
            'layout_type': '3x3',
            'shape_count': len(shapes),
            'positions': positions[:len(shapes)]
        }
    
    def _adjust_text_size(self, shape, shape_width: float, shape_height: float):
        """
        根据形状大小调整文本大小
        """
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return
            
            # 根据形状大小决定字体大小
            area = (shape_width / Inches(1)) * (shape_height / Inches(1))
            
            thresholds = self.config.layout_thresholds
            font_sizes = self.config.font_sizes
            
            if area > thresholds['large_area']:
                font_size = Pt(font_sizes['large_area'])
            elif area > thresholds['medium_area']:
                font_size = Pt(font_sizes['medium_area'])
            else:
                font_size = Pt(font_sizes['small_area'])
            
            # 应用字体大小
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size
                
                # 设置段落对齐
                paragraph.alignment = PP_ALIGN.LEFT
            
            # 调整文本框边距
            text_frame = shape.text_frame
            margin = Inches(self.config.layout_margins['shape_margin'])
            text_frame.margin_left = margin
            text_frame.margin_right = margin
            text_frame.margin_top = margin
            text_frame.margin_bottom = margin
            
        except Exception as e:
            self.logger.error(f"调整文本大小时出错: {e}")
    
    def remove_empty_slides(self) -> List[int]:
        """
        删除空的幻灯片
        
        Returns:
            List[int]: 被删除的幻灯片索引列表
        """
        removed_slides = []
        
        # 从后往前遍历，避免索引变化问题
        for i in range(len(self.presentation.slides) - 1, -1, -1):
            slide = self.presentation.slides[i]
            
            # 检查幻灯片是否为空
            has_content = False
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    shape_text = getattr(shape, 'text', '')
                    if shape_text and shape_text.strip():
                        # 检查是否还有未填充的占位符
                        if not re.search(r'\{[^}]+\}', shape_text):
                            has_content = True
                            break
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_content = True
                    break
            
            if not has_content:
                # 删除空幻灯片
                xml_slides = self.presentation.slides._sldIdLst
                xml_slides.remove(xml_slides[i])
                removed_slides.append(i)
                self.logger.info(f"删除空幻灯片: {i+1}")
        
        return list(reversed(removed_slides))  # 返回原始顺序
    
    def optimize_slide_sequence(self) -> Dict[str, Any]:
        """
        优化幻灯片序列
        
        Returns:
            Dict: 优化结果
        """
        results = {
            'total_slides_before': len(self.presentation.slides),
            'removed_empty_slides': [],
            'final_slide_count': 0
        }
        
        # 删除空幻灯片
        results['removed_empty_slides'] = self.remove_empty_slides()
        results['final_slide_count'] = len(self.presentation.slides)
        
        self.logger.info(f"幻灯片优化完成: {results['total_slides_before']} -> {results['final_slide_count']}")
        
        return results