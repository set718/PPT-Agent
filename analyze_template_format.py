#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT模板格式分析工具
分析split_presentations_1.pptx的布局、颜色、字体等信息，并优化合并策略
"""

import os
import sys
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def analyze_template_format(template_path: str) -> Dict[str, Any]:
    """
    深度分析PPT模板的格式信息
    
    Args:
        template_path: 模板文件路径
        
    Returns:
        Dict: 格式分析结果
    """
    if not os.path.exists(template_path):
        return {"error": f"模板文件不存在: {template_path}"}
    
    try:
        ppt = Presentation(template_path)
        analysis = {
            "file_path": template_path,
            "slide_count": len(ppt.slides),
            "slide_size": {
                "width": ppt.slide_width,
                "height": ppt.slide_height
            },
            "slides": []
        }
        
        for slide_idx, slide in enumerate(ppt.slides):
            slide_info = {
                "slide_index": slide_idx,
                "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
                "background": analyze_background(slide),
                "shapes": [],
                "color_scheme": analyze_color_scheme(slide),
                "font_info": {}
            }
            
            # 分析所有形状
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = analyze_shape(shape, shape_idx)
                slide_info["shapes"].append(shape_info)
            
            # 提取字体信息统计
            slide_info["font_info"] = extract_font_statistics(slide_info["shapes"])
            
            analysis["slides"].append(slide_info)
        
        return analysis
        
    except Exception as e:
        return {"error": f"分析失败: {str(e)}"}

def analyze_background(slide) -> Dict[str, Any]:
    """分析幻灯片背景"""
    background_info = {
        "type": "unknown",
        "color": None,
        "fill_type": None
    }
    
    try:
        if hasattr(slide, 'background'):
            bg = slide.background
            if hasattr(bg, 'fill'):
                fill = bg.fill
                background_info["fill_type"] = str(fill.type) if hasattr(fill, 'type') else None
                
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    try:
                        rgb = fill.fore_color.rgb
                        background_info["color"] = f"RGB({rgb.red}, {rgb.green}, {rgb.blue})"
                    except:
                        background_info["color"] = "Cannot determine"
    except:
        pass
    
    return background_info

def analyze_color_scheme(slide) -> Dict[str, Any]:
    """分析颜色方案"""
    colors = {
        "text_colors": [],
        "fill_colors": [],
        "accent_colors": []
    }
    
    try:
        if hasattr(slide, 'color_scheme'):
            # 尝试提取颜色方案信息
            pass
    except:
        pass
    
    return colors

def analyze_shape(shape, shape_idx: int) -> Dict[str, Any]:
    """分析单个形状的详细信息"""
    shape_info = {
        "index": shape_idx,
        "type": str(shape.shape_type) if hasattr(shape, 'shape_type') else "Unknown",
        "position": {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height
        },
        "has_text": hasattr(shape, 'text'),
        "text_info": None,
        "fill_info": None
    }
    
    # 分析文本信息
    if hasattr(shape, 'text') and hasattr(shape, 'text_frame'):
        shape_info["text_info"] = analyze_text_format(shape)
    
    # 分析填充信息
    if hasattr(shape, 'fill'):
        shape_info["fill_info"] = analyze_fill_format(shape.fill)
    
    return shape_info

def analyze_text_format(shape) -> Dict[str, Any]:
    """分析文本格式"""
    text_info = {
        "text_content": shape.text[:100] + "..." if len(shape.text) > 100 else shape.text,
        "text_length": len(shape.text),
        "paragraphs": []
    }
    
    try:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                para_info = {
                    "index": para_idx,
                    "text": paragraph.text[:50] + "..." if len(paragraph.text) > 50 else paragraph.text,
                    "alignment": str(paragraph.alignment) if hasattr(paragraph, 'alignment') else None,
                    "level": paragraph.level if hasattr(paragraph, 'level') else None,
                    "font_info": None
                }
                
                # 分析字体
                if hasattr(paragraph, 'font') and paragraph.font:
                    para_info["font_info"] = {
                        "name": paragraph.font.name,
                        "size": paragraph.font.size.pt if paragraph.font.size else None,
                        "bold": paragraph.font.bold,
                        "italic": paragraph.font.italic,
                        "color": None
                    }
                    
                    # 尝试获取字体颜色
                    try:
                        if paragraph.font.color and paragraph.font.color.rgb:
                            rgb = paragraph.font.color.rgb
                            para_info["font_info"]["color"] = f"RGB({rgb.red}, {rgb.green}, {rgb.blue})"
                    except:
                        para_info["font_info"]["color"] = "Cannot determine"
                
                text_info["paragraphs"].append(para_info)
    except Exception as e:
        text_info["error"] = str(e)
    
    return text_info

def analyze_fill_format(fill) -> Dict[str, Any]:
    """分析填充格式"""
    fill_info = {
        "type": str(fill.type) if hasattr(fill, 'type') else None,
        "color": None,
        "transparency": None
    }
    
    try:
        if hasattr(fill, 'fore_color') and fill.fore_color:
            try:
                rgb = fill.fore_color.rgb
                fill_info["color"] = f"RGB({rgb.red}, {rgb.green}, {rgb.blue})"
            except:
                fill_info["color"] = "Cannot determine"
                
        if hasattr(fill, 'transparency'):
            fill_info["transparency"] = fill.transparency
    except:
        pass
    
    return fill_info

def extract_font_statistics(shapes: List[Dict]) -> Dict[str, Any]:
    """提取字体统计信息"""
    fonts = {}
    colors = {}
    sizes = {}
    
    for shape in shapes:
        if shape.get("text_info") and shape["text_info"].get("paragraphs"):
            for para in shape["text_info"]["paragraphs"]:
                font_info = para.get("font_info")
                if font_info:
                    # 统计字体名称
                    font_name = font_info.get("name")
                    if font_name:
                        fonts[font_name] = fonts.get(font_name, 0) + 1
                    
                    # 统计字体颜色
                    font_color = font_info.get("color")
                    if font_color:
                        colors[font_color] = colors.get(font_color, 0) + 1
                    
                    # 统计字体大小
                    font_size = font_info.get("size")
                    if font_size:
                        sizes[str(font_size)] = sizes.get(str(font_size), 0) + 1
    
    return {
        "most_common_fonts": sorted(fonts.items(), key=lambda x: x[1], reverse=True)[:5],
        "most_common_colors": sorted(colors.items(), key=lambda x: x[1], reverse=True)[:5],
        "most_common_sizes": sorted(sizes.items(), key=lambda x: x[1], reverse=True)[:5]
    }

def print_analysis_report(analysis: Dict[str, Any]):
    """打印分析报告"""
    if "error" in analysis:
        print(f"❌ {analysis['error']}")
        return
    
    print("🔍 PPT模板格式分析报告")
    print("=" * 60)
    print(f"文件: {analysis['file_path']}")
    print(f"幻灯片数量: {analysis['slide_count']}")
    print(f"幻灯片尺寸: {analysis['slide_size']['width']} x {analysis['slide_size']['height']}")
    
    for slide in analysis["slides"]:
        print(f"\n📄 第 {slide['slide_index'] + 1} 页 - {slide['layout_name']}")
        print("-" * 40)
        
        # 背景信息
        bg = slide["background"]
        if bg["color"]:
            print(f"背景: {bg['fill_type']} - {bg['color']}")
        
        # 字体统计
        font_stats = slide["font_info"]
        if font_stats["most_common_fonts"]:
            print("主要字体:")
            for font, count in font_stats["most_common_fonts"]:
                print(f"  - {font} (使用 {count} 次)")
        
        if font_stats["most_common_colors"]:
            print("主要颜色:")
            for color, count in font_stats["most_common_colors"]:
                print(f"  - {color} (使用 {count} 次)")
        
        if font_stats["most_common_sizes"]:
            print("主要字号:")
            for size, count in font_stats["most_common_sizes"]:
                print(f"  - {size}pt (使用 {count} 次)")
        
        # 形状信息
        print(f"形状数量: {len(slide['shapes'])}")
        for shape in slide["shapes"]:
            if shape["has_text"] and shape["text_info"]:
                text_info = shape["text_info"]
                if text_info["text_content"].strip():
                    print(f"  📝 文本: '{text_info['text_content'][:30]}...' ({text_info['text_length']} 字符)")

def main():
    """主函数"""
    # 分析split_presentations_1.pptx
    template_path = os.path.join("templates", "ppt_template", "split_presentations_1.pptx")
    
    print("🎯 分析基准模板: split_presentations_1.pptx")
    print("=" * 60)
    
    analysis = analyze_template_format(template_path)
    print_analysis_report(analysis)
    
    # 保存分析结果
    import json
    output_file = "template_format_analysis.json"
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(analysis, f, ensure_ascii=False, indent=2, default=str)
        print(f"\n📁 详细分析结果已保存到: {output_file}")
    except Exception as e:
        print(f"❌ 保存分析结果失败: {e}")
    
    return analysis

if __name__ == "__main__":
    main()
