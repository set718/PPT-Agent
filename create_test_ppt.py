#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
创建包含占位符的测试PPT文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def create_test_ppt():
    """创建包含占位符的测试PPT"""
    presentation = Presentation()
    
    # 第一页：包含文本框占位符
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    
    # 添加标题文本框
    title_shape = slide1.shapes.title
    title_shape.text = "{title} - 这是一个测试标题"
    
    # 添加内容文本框
    content_shape = slide1.placeholders[1]
    content_shape.text = """
这是测试内容：
• {bullet_1} - 第一个要点
• {bullet_2} - 第二个要点
• {content} - 主要内容
• {description} - 详细描述
"""
    
    # 第二页：包含表格占位符
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])  # 空白布局
    
    # 添加标题
    title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "{slide2_title} - 表格测试页"
    
    # 添加表格
    table = slide2.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(8), Inches(4))
    table = table.table
    
    # 填入表格占位符
    table.cell(0, 0).text = "时间"
    table.cell(0, 1).text = "活动"
    table.cell(0, 2).text = "负责人"
    
    table.cell(1, 0).text = "{bullet_1_time_1}"
    table.cell(1, 1).text = "{bullet_1_activity_1}"
    table.cell(1, 2).text = "{bullet_1_person_1}"
    
    table.cell(2, 0).text = "{bullet_2_time_1}"
    table.cell(2, 1).text = "{bullet_2_activity_1}"  
    table.cell(2, 2).text = "{bullet_2_person_1}"
    
    # 添加更多文本框占位符
    text_box = slide2.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = "备注: {notes} | 总结: {summary}"
    
    # 保存文件
    presentation.save("test_placeholder_template.pptx")
    print("已创建测试文件: test_placeholder_template.pptx")
    
    # 统计占位符
    total_placeholders = []
    for slide_idx, slide in enumerate(presentation.slides):
        slide_placeholders = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                import re
                placeholders = re.findall(r'\{([^}]+)\}', shape.text)
                slide_placeholders.extend(placeholders)
            elif hasattr(shape, 'table'):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        placeholders = re.findall(r'\{([^}]+)\}', cell.text)
                        slide_placeholders.extend(placeholders)
        
        if slide_placeholders:
            print(f"第{slide_idx+1}页包含占位符: {slide_placeholders}")
            total_placeholders.extend(slide_placeholders)
    
    print(f"总共包含{len(total_placeholders)}个占位符: {set(total_placeholders)}")

if __name__ == "__main__":
    create_test_ppt()