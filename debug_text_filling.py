#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文本填充调试脚本
用于诊断AI文本填充功能的问题
"""

import os
import json
from pptx import Presentation
from config import get_config
from utils import AIProcessor, PPTProcessor, PPTAnalyzer

def debug_text_filling():
    """调试文本填充功能"""
    print("🔍 文本填充功能调试")
    print("=" * 50)
    
    config = get_config()
    
    # 检查PPT模板
    ppt_path = config.default_ppt_template
    if not os.path.exists(ppt_path):
        print(f"❌ PPT模板文件不存在: {ppt_path}")
        return
    
    print(f"📄 使用PPT文件: {ppt_path}")
    
    # 步骤1: 分析PPT结构
    print("\n🔍 步骤1: 分析PPT结构")
    print("-" * 30)
    
    presentation = Presentation(ppt_path)
    ppt_structure = PPTAnalyzer.analyze_ppt_structure(presentation)
    
    print(f"总页数: {ppt_structure['total_slides']}")
    
    for i, slide_info in enumerate(ppt_structure['slides']):
        print(f"\n📑 第{i+1}页:")
        print(f"  标题: {slide_info['title'] or '无'}")
        print(f"  占位符数量: {len(slide_info['placeholders'])}")
        
        if slide_info['placeholders']:
            print("  占位符详情:")
            for placeholder_name, placeholder_info in slide_info['placeholders'].items():
                original_text = placeholder_info['original_text']
                print(f"    - {{{placeholder_name}}} -> 原文: '{original_text}'")
    
    # 步骤2: 测试AI分析（如果提供API密钥）
    print("\n🤖 步骤2: 测试AI分析")
    print("-" * 30)
    
    api_key = input("请输入OpenRouter API密钥 (留空跳过AI测试): ").strip()
    
    if api_key:
        try:
            ai_processor = AIProcessor(api_key)
            
            # 测试文本
            test_text = """
            人工智能的发展历程
            
            人工智能技术经历了多个重要阶段。从1950年代的符号主义开始，到1980年代的专家系统，再到现在的深度学习。
            
            当前的大语言模型如GPT、Claude等展现出前所未有的能力。这些技术正在革新各个行业。
            
            未来，人工智能将继续向更加智能化的方向发展。
            """
            
            print("📝 测试文本:")
            print(f"'{test_text.strip()}'")
            print()
            
            print("🔄 正在调用AI分析...")
            assignments = ai_processor.analyze_text_for_ppt(test_text, ppt_structure)
            
            print("📊 AI返回的分配结果:")
            print(json.dumps(assignments, ensure_ascii=False, indent=2))
            
            # 检查assignments结构
            assignments_list = assignments.get('assignments', [])
            print(f"\n📋 分配方案数量: {len(assignments_list)}")
            
            for i, assignment in enumerate(assignments_list):
                print(f"\n分配方案 {i+1}:")
                print(f"  动作: {assignment.get('action', '未知')}")
                print(f"  幻灯片: {assignment.get('slide_index', 0) + 1}")
                print(f"  占位符: {assignment.get('placeholder', '未知')}")
                print(f"  内容: '{assignment.get('content', '')}'")
                print(f"  原因: {assignment.get('reason', '无')}")
            
            # 步骤3: 测试实际填充
            print("\n🔧 步骤3: 测试实际填充")
            print("-" * 30)
            
            if assignments_list:
                ppt_processor = PPTProcessor(presentation)
                results = ppt_processor.apply_assignments(assignments)
                
                print("📝 填充结果:")
                for result in results:
                    print(f"  {result}")
                
                # 检查填充后的内容
                print("\n📋 填充后验证:")
                for i, slide in enumerate(presentation.slides):
                    print(f"\n第{i+1}页内容:")
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            print(f"  文本框: '{shape.text.strip()}'")
            else:
                print("⚠️ AI没有返回任何分配方案")
                
        except Exception as e:
            print(f"❌ AI分析失败: {e}")
    else:
        print("⏭️ 跳过AI测试")
    
    # 步骤4: 检查美化过程
    print("\n🎨 步骤4: 检查美化过程")
    print("-" * 30)
    
    ppt_processor = PPTProcessor(presentation)
    
    # 模拟一些已填充的占位符
    ppt_processor.filled_placeholders = {0: {'title', 'content'}}
    
    print("模拟已填充占位符:", ppt_processor.filled_placeholders)
    
    beautify_results = ppt_processor.beautify_presentation()
    print("\n美化结果摘要:")
    summary = beautify_results.get('summary', {})
    for key, value in summary.items():
        print(f"  {key}: {value}")

def check_ppt_template():
    """检查PPT模板的占位符格式"""
    print("\n🔍 检查PPT模板占位符格式")
    print("=" * 40)
    
    config = get_config()
    ppt_path = config.default_ppt_template
    
    if not os.path.exists(ppt_path):
        print(f"❌ PPT模板文件不存在: {ppt_path}")
        return
    
    presentation = Presentation(ppt_path)
    
    for i, slide in enumerate(presentation.slides):
        print(f"\n📑 第{i+1}页所有文本内容:")
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    print(f"  文本框{j+1}: '{text}'")
                    
                    # 检查是否包含占位符
                    import re
                    placeholders = re.findall(r'\{([^}]+)\}', text)
                    if placeholders:
                        print(f"    发现占位符: {placeholders}")
                    else:
                        print(f"    无占位符")

if __name__ == "__main__":
    print("PPT文本填充调试工具")
    print("=" * 50)
    
    mode = input("选择模式: \n1. 完整调试\n2. 只检查模板\n请输入 (1/2): ").strip()
    
    if mode == "1":
        debug_text_filling()
    elif mode == "2":
        check_ppt_template()
    else:
        print("无效选择")