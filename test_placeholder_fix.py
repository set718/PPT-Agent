#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
测试占位符清理修复效果
"""

import os
from pptx import Presentation
from utils import PPTProcessor
from user_app import UserPPTGenerator

def test_placeholder_cleanup_fix():
    """测试占位符清理修复效果"""
    print("=== 测试占位符清理修复效果 ===")
    
    # 使用之前创建的测试文件
    test_file = "test_placeholder_template.pptx"
    
    if not os.path.exists(test_file):
        print(f"测试文件不存在: {test_file}")
        print("请先运行 create_test_ppt.py 创建测试文件")
        return
    
    try:
        # 1. 初始化用户PPT生成器（模拟用户界面使用场景）
        print("\n1. 初始化PPT生成器...")
        # 使用占位符API密钥模拟
        mock_api_key = "sk-test-placeholder-api-key"
        generator = UserPPTGenerator(mock_api_key)
        
        # 2. 加载PPT文件
        print("2. 加载PPT文件...")
        success, message = generator.load_ppt_from_path(test_file)
        if not success:
            print(f"加载PPT失败: {message}")
            return
        
        print(f"成功加载PPT: {message}")
        
        # 3. 分析PPT结构（查看检测到的占位符）
        print("\n3. 分析PPT结构...")
        ppt_structure = generator.ppt_processor.ppt_structure
        total_placeholders = 0
        
        for slide_idx, slide_info in enumerate(ppt_structure['slides']):
            placeholders = slide_info.get('placeholders', {})
            if placeholders:
                print(f"第{slide_idx+1}页检测到{len(placeholders)}个占位符:")
                for placeholder_name, info in placeholders.items():
                    placeholder_type = info.get('type', 'text_box')
                    if placeholder_type == 'table_cell':
                        print(f"  - {{{placeholder_name}}} (表格{info['row_idx']+1},{info['col_idx']+1})")
                    else:
                        print(f"  - {{{placeholder_name}}} (文本框)")
                    total_placeholders += 1
        
        print(f"总计检测到{total_placeholders}个占位符")
        
        # 4. 模拟AI分配（创建测试分配）
        print("\n4. 模拟AI分配...")
        test_assignments = {
            "assignments": [
                {"slide_index": 0, "action": "replace_placeholder", "placeholder": "title", "content": "AI测试标题", "reason": "测试标题填充"},
                {"slide_index": 0, "action": "replace_placeholder", "placeholder": "bullet_1", "content": "测试要点1", "reason": "测试要点填充"},
                {"slide_index": 0, "action": "replace_placeholder", "placeholder": "content", "content": "测试主要内容", "reason": "测试内容填充"},
                # 故意跳过一些占位符，如 bullet_2, description
                
                {"slide_index": 1, "action": "replace_placeholder", "placeholder": "slide2_title", "content": "第二页测试标题", "reason": "测试标题填充"},
                {"slide_index": 1, "action": "replace_placeholder", "placeholder": "bullet_1_time_1", "content": "09:00", "reason": "测试表格时间填充"},
                {"slide_index": 1, "action": "replace_placeholder", "placeholder": "bullet_1_activity_1", "content": "开始活动", "reason": "测试表格活动填充"},
                {"slide_index": 1, "action": "replace_placeholder", "placeholder": "notes", "content": "重要备注信息", "reason": "测试备注填充"},
                # 故意跳过表格中的其他占位符和summary
            ]
        }
        
        # 5. 应用分配（模拟填充过程）
        print("5. 应用占位符分配...")
        success, results = generator.apply_text_assignments(test_assignments, "原始测试文本")
        
        if success:
            print("占位符填充完成")
            print(f"处理结果: {len(results)}条")
            
            # 显示处理结果
            for result in results:
                print(f"  - {result}")
        else:
            print(f"占位符填充失败: {results}")
            return
        
        # 6. 检查filled_placeholders记录
        print("\n6. 检查填充记录...")
        filled_placeholders = generator.ppt_processor.filled_placeholders
        total_filled = 0
        
        for slide_idx, filled_set in filled_placeholders.items():
            if filled_set:
                print(f"第{slide_idx+1}页已填充: {', '.join([f'{{{p}}}' for p in filled_set])}")
                total_filled += len(filled_set)
        
        print(f"总计记录已填充{total_filled}个占位符")
        
        # 7. 执行清理功能
        print("\n7. 执行占位符清理...")
        cleanup_results = generator.cleanup_unfilled_placeholders()
        
        if cleanup_results.get('success'):
            cleaned_count = cleanup_results.get('cleaned_placeholders', 0)
            cleaned_list = cleanup_results.get('cleaned_placeholder_list', [])
            
            print(f"清理完成，共清理{cleaned_count}个文本框/单元格中的未填充占位符")
            
            if cleaned_list:
                print("清理详情:")
                for item in cleaned_list:
                    print(f"  - {item}")
            else:
                print("没有需要清理的占位符")
        else:
            error_msg = cleanup_results.get('error', '未知错误')
            print(f"清理失败: {error_msg}")
        
        # 8. 验证清理效果
        print("\n8. 验证清理效果...")
        verification_results = verify_cleanup_results(generator.presentation, filled_placeholders)
        
        if verification_results['success']:
            print("✅ 清理效果验证通过")
            print(f"验证详情: {verification_results['message']}")
        else:
            print("❌ 清理效果验证失败")
            print(f"问题: {verification_results['issues']}")
        
        print("\n=== 测试完成 ===")
        
    except Exception as e:
        print(f"测试过程中出现异常: {e}")
        import traceback
        traceback.print_exc()

def verify_cleanup_results(presentation, filled_placeholders):
    """验证清理结果是否正确"""
    try:
        issues = []
        total_remaining = 0
        total_expected_remaining = 0
        
        for slide_idx, slide in enumerate(presentation.slides):
            filled_set = filled_placeholders.get(slide_idx, set())
            
            # 检查文本框
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    import re
                    remaining_placeholders = re.findall(r'\{([^}]+)\}', shape.text)
                    
                    for placeholder in remaining_placeholders:
                        total_remaining += 1
                        if placeholder in filled_set:
                            # 已填充的占位符不应该还存在
                            total_expected_remaining += 1
                        else:
                            # 未填充的占位符被清理是正常的
                            issues.append(f"第{slide_idx+1}页文本框仍有未填充占位符{{{placeholder}}}，但这是预期的")
                
                # 检查表格
                elif hasattr(shape, 'shape_type') and shape.shape_type == 19:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if cell_text:
                                remaining_placeholders = re.findall(r'\{([^}]+)\}', cell_text)
                                
                                for placeholder in remaining_placeholders:
                                    total_remaining += 1
                                    if placeholder in filled_set:
                                        # 已填充的占位符不应该还存在
                                        total_expected_remaining += 1
                                    else:
                                        # 未填充的占位符被清理是正常的
                                        issues.append(f"第{slide_idx+1}页表格{row_idx+1},{col_idx+1}仍有未填充占位符{{{placeholder}}}，但这是预期的")
        
        # 如果有已填充但仍然存在的占位符，说明清理有问题
        if total_expected_remaining > 0:
            return {
                'success': False,
                'issues': f"发现{total_expected_remaining}个已填充但未被替换的占位符，清理逻辑可能有问题"
            }
        else:
            return {
                'success': True,
                'message': f"清理正常，剩余{total_remaining}个占位符都是未填充的占位符（已被清理）"
            }
    
    except Exception as e:
        return {
            'success': False,
            'issues': f"验证过程出错: {e}"
        }

if __name__ == "__main__":
    test_placeholder_cleanup_fix()