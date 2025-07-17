#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT美化功能演示脚本
展示占位符清理和重新排版功能
"""

import os
import sys
from pptx import Presentation
from config import get_config
from utils import PPTProcessor, PPTAnalyzer
from ppt_beautifier import PPTBeautifier
from logger import get_logger

def main():
    """演示主函数"""
    print("=" * 60)
    print("        PPT美化功能演示")
    print("    展示占位符清理和重新排版功能")
    print("=" * 60)
    
    config = get_config()
    logger = get_logger()
    
    # 获取PPT文件路径
    ppt_path = config.default_ppt_template
    
    if not os.path.exists(ppt_path):
        print(f"[ERROR] PPT模板文件不存在: {ppt_path}")
        return
    
    try:
        # 加载PPT
        presentation = Presentation(ppt_path)
        print(f"[INFO] 加载PPT: {os.path.basename(ppt_path)}")
        
        # 分析PPT结构
        analyzer = PPTAnalyzer()
        ppt_structure = analyzer.analyze_ppt_structure(presentation)
        
        print(f"[INFO] PPT结构分析:")
        print(f"   总幻灯片数: {ppt_structure['total_slides']}")
        
        total_placeholders = 0
        for slide in ppt_structure['slides']:
            placeholders = slide.get('placeholders', {})
            total_placeholders += len(placeholders)
            
            if placeholders:
                print(f"   第{slide['slide_index']+1}页: {len(placeholders)} 个占位符")
                for placeholder_name in placeholders.keys():
                    print(f"      - {{{placeholder_name}}}")
        
        print(f"[INFO] 总占位符数: {total_placeholders}")
        
        # 创建PPT处理器
        processor = PPTProcessor(presentation)
        
        # 模拟填充部分占位符（仅填充前4个）
        print(f"\n[INFO] 模拟填充部分占位符...")
        filled_count = 0
        for slide_idx, slide in enumerate(ppt_structure['slides']):
            placeholders = slide.get('placeholders', {})
            if placeholders and filled_count < 4:
                # 只填充前几个占位符
                for placeholder_name, placeholder_info in list(placeholders.items())[:2]:
                    if filled_count < 4:
                        # 模拟填充
                        test_content = f"测试内容 {filled_count + 1}"
                        success = processor._replace_placeholder_in_slide(placeholder_info, test_content)
                        
                        if success:
                            # 记录已填充的占位符
                            if slide_idx not in processor.filled_placeholders:
                                processor.filled_placeholders[slide_idx] = set()
                            processor.filled_placeholders[slide_idx].add(placeholder_name)
                            
                            print(f"   [OK] 填充第{slide_idx+1}页的 {{{placeholder_name}}}: {test_content}")
                            filled_count += 1
                        
                        if filled_count >= 4:
                            break
        
        print(f"[INFO] 模拟填充完成，已填充 {filled_count} 个占位符")
        
        # 显示填充情况
        print(f"\n[INFO] 已填充的占位符:")
        for slide_idx, placeholders in processor.filled_placeholders.items():
            print(f"   第{slide_idx+1}页: {list(placeholders)}")
        
        # 执行美化
        print(f"\n[INFO] 开始美化PPT...")
        beautify_results = processor.beautify_presentation()
        
        # 显示美化结果
        summary = beautify_results['summary']
        print(f"\n[INFO] 美化结果:")
        print(f"   删除未填充占位符: {summary['removed_placeholders_count']} 个")
        print(f"   重新排版幻灯片: {summary['reorganized_slides_count']} 页")
        print(f"   删除空幻灯片: {summary['removed_empty_slides_count']} 页")
        print(f"   最终幻灯片数: {summary['final_slide_count']} 页")
        
        # 显示详细结果
        if summary['removed_placeholders_count'] > 0:
            print(f"\n[INFO] 清理详情:")
            for item in beautify_results['beautify_results']['removed_placeholders']:
                print(f"   第{item['slide_index']+1}页: 删除了 {item['removed_count']} 个未填充占位符")
                for placeholder in item['removed_placeholders']:
                    print(f"      - {{{placeholder}}}")
        
        if summary['reorganized_slides_count'] > 0:
            print(f"\n[INFO] 重排版详情:")
            for item in beautify_results['beautify_results']['reorganized_slides']:
                layout_change = item['layout_change']
                print(f"   第{item['slide_index']+1}页: 使用 {layout_change['layout_type']} 布局重新排版了 {layout_change['shape_count']} 个元素")
        
        # 保存演示结果
        output_path = os.path.join(config.output_dir, "beautify_demo.pptx")
        presentation.save(output_path)
        print(f"\n[INFO] 演示结果已保存: {output_path}")
        
        print(f"\n[INFO] 演示完成！")
        print("您可以打开生成的PPT文件查看美化效果。")
        print("对比原始模板文件，您会发现:")
        print("1. 未填充的占位符已被删除")
        print("2. 剩余的内容已重新排版为更美观的布局")
        print("3. 空的幻灯片已被移除")
        
    except Exception as e:
        logger.exception(f"演示过程中出现错误: {e}")
        print(f"[ERROR] 演示过程中出现错误: {e}")

if __name__ == "__main__":
    main()