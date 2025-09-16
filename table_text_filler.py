#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
è¡¨æ ¼æ–‡æœ¬å¡«å……å™¨æ¨¡å—
ä¸“é—¨ç”¨äºå¤„ç†PPTæ¨¡æ¿ä¸­è¡¨æ ¼å ä½ç¬¦çš„æ™ºèƒ½å¡«å……
"""

import re
import json
from typing import Dict, List, Any, Optional, Tuple
from pptx import Presentation
from pptx.util import Pt
from utils import AIProcessor, PPTProcessor, FileManager
from logger import log_user_action, log_file_operation, LogContext


class TableTextFiller:
    """è¡¨æ ¼æ–‡æœ¬æ™ºèƒ½å¡«å……å™¨"""
    
    def __init__(self, api_key: Optional[str] = None):
        """åˆå§‹åŒ–è¡¨æ ¼æ–‡æœ¬å¡«å……å™¨"""
        self.api_key = api_key
        self.ai_processor = AIProcessor(api_key)
        self.presentation = None
        self.ppt_processor = None
        self.ppt_structure = None
        self.table_info = {}  # å­˜å‚¨è¡¨æ ¼ä¿¡æ¯
        
    def load_ppt_template(self, ppt_path: str) -> Tuple[bool, str]:
        """åŠ è½½PPTæ¨¡æ¿"""
        with LogContext(f"è¡¨æ ¼å¡«å……å™¨åŠ è½½PPTæ¨¡æ¿"):
            try:
                # éªŒè¯æ–‡ä»¶
                is_valid, error_msg = FileManager.validate_ppt_file(ppt_path)
                if not is_valid:
                    return False, error_msg
                
                self.presentation = Presentation(ppt_path)
                self.ppt_processor = PPTProcessor(self.presentation)
                self.ppt_structure = self.ppt_processor.ppt_structure
                
                # åˆ†æè¡¨æ ¼ç»“æ„
                self._analyze_table_structure()
                
                log_file_operation("load_ppt_table_filler", ppt_path, "success")
                return True, "æ¨¡æ¿åŠ è½½æˆåŠŸ"
                
            except Exception as e:
                log_file_operation("load_ppt_table_filler", ppt_path, "error", str(e))
                return False, f"åŠ è½½å¤±è´¥: {str(e)}"
    
    def _analyze_table_structure(self):
        """åˆ†æPPTä¸­æ‰€æœ‰è¡¨æ ¼çš„ç»“æ„"""
        self.table_info = {}
        
        if not self.presentation:
            return
            
        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_tables = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                # æ£€æŸ¥æ˜¯å¦ä¸ºè¡¨æ ¼
                if hasattr(shape, 'shape_type') and shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
                    table = shape.table
                    table_data = {
                        'slide_idx': slide_idx,
                        'shape_idx': shape_idx,
                        'rows': len(table.rows),
                        'cols': len(table.columns),
                        'placeholders': []
                    }
                    
                    # åˆ†æè¡¨æ ¼ä¸­çš„å ä½ç¬¦
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if cell_text:
                                placeholders = re.findall(r'\{([^}]+)\}', cell_text)
                                for placeholder in placeholders:
                                    table_data['placeholders'].append({
                                        'placeholder': placeholder,
                                        'row': row_idx,
                                        'col': col_idx,
                                        'cell_text': cell_text,
                                        'position': f"è¡Œ{row_idx+1}åˆ—{col_idx+1}"
                                    })
                    
                    if table_data['placeholders']:  # åªè®°å½•åŒ…å«å ä½ç¬¦çš„è¡¨æ ¼
                        slide_tables.append(table_data)
            
            if slide_tables:
                self.table_info[slide_idx] = slide_tables
    
    def get_table_analysis(self) -> Dict[str, Any]:
        """è·å–è¡¨æ ¼åˆ†æç»“æœ"""
        analysis = {
            'total_tables': 0,
            'total_table_placeholders': 0,
            'slides_with_tables': [],
            'table_details': []
        }
        
        for slide_idx, tables in self.table_info.items():
            slide_info = {
                'slide_number': slide_idx + 1,
                'table_count': len(tables),
                'tables': []
            }
            
            for table_idx, table_data in enumerate(tables):
                table_detail = {
                    'table_index': table_idx + 1,
                    'size': f"{table_data['rows']}è¡Œ{table_data['cols']}åˆ—",
                    'placeholder_count': len(table_data['placeholders']),
                    'placeholders': [
                        f"{{{p['placeholder']}}}"
                        for p in table_data['placeholders']
                    ]
                }
                slide_info['tables'].append(table_detail)
                analysis['total_table_placeholders'] += len(table_data['placeholders'])
            
            analysis['slides_with_tables'].append(slide_info)
            analysis['total_tables'] += len(tables)
        
        return analysis
    
    def process_table_text_with_ai(self, user_text: str) -> Dict[str, Any]:
        """ä½¿ç”¨AIåˆ†ææ–‡æœ¬å¹¶ç”Ÿæˆè¡¨æ ¼å¡«å……æ–¹æ¡ˆ"""
        if not self.table_info:
            return {"assignments": []}
        
        # æ„å»ºä¸“é—¨é’ˆå¯¹è¡¨æ ¼å¡«å……çš„ç³»ç»Ÿæç¤º
        system_prompt = self._build_table_filling_prompt()
        
        # è·å–è¡¨æ ¼ç»“æ„ä¿¡æ¯
        table_structure_info = self._get_table_structure_for_ai()
        
        # ç»„åˆå®Œæ•´çš„åˆ†ææç¤º
        full_prompt = f"{system_prompt}\n\n{table_structure_info}\n\nç”¨æˆ·æ–‡æœ¬å†…å®¹ï¼š\n{user_text}"
        
        log_user_action("è¡¨æ ¼æ–‡æœ¬AIåˆ†æ", f"æ–‡æœ¬é•¿åº¦: {len(user_text)}å­—ç¬¦, è¡¨æ ¼æ•°é‡: {sum(len(tables) for tables in self.table_info.values())}")
        
        try:
            # è°ƒç”¨AIè¿›è¡Œåˆ†æï¼Œä½¿ç”¨ç±»ä¼¼analyze_text_for_pptçš„æ–¹å¼
            self.ai_processor._ensure_client()
            
            # æ£€æŸ¥APIç±»å‹
            model_info = self.ai_processor.config.get_model_info()
            
            if model_info.get('request_format') == 'dify_compatible':
                # ä½¿ç”¨Liai APIæ ¼å¼
                content = self.ai_processor._call_liai_api(system_prompt, user_text)
            else:
                # ä½¿ç”¨OpenAIæ ¼å¼
                actual_model = model_info.get('actual_model', self.ai_processor.config.ai_model)
                
                response = self.ai_processor.client.chat.completions.create(
                    model=actual_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"è¡¨æ ¼ç»“æ„ä¿¡æ¯ï¼š\n{table_structure_info}\n\nç”¨æˆ·æ–‡æœ¬ï¼š\n{user_text}"}
                    ],
                    temperature=0.3,
                    stream=True
                )
                
                # æ”¶é›†æµå¼å“åº”å†…å®¹
                content = ""
                for chunk in response:
                    if chunk.choices and chunk.choices[0].delta and chunk.choices[0].delta.content:
                        content += chunk.choices[0].delta.content
                
                content = content.strip() if content else ""
            
            # è§£æAIè¿”å›ç»“æœ
            return self._parse_table_assignments(content)
            
        except Exception as e:
            log_user_action("è¡¨æ ¼æ–‡æœ¬AIåˆ†æå¤±è´¥", str(e))
            return {"error": f"AIåˆ†æå¤±è´¥: {str(e)}"}
    
    def _build_table_filling_prompt(self) -> str:
        """æ„å»ºè¡¨æ ¼å¡«å……çš„AIæç¤º"""
        return """ä½ æ˜¯ä¸€ä¸ªä¸“ä¸šçš„è¡¨æ ¼å†…å®¹å¡«å……ä¸“å®¶ã€‚ä½ çš„ä»»åŠ¡æ˜¯åˆ†æç”¨æˆ·æä¾›çš„æ–‡æœ¬å†…å®¹ï¼Œå¹¶æ™ºèƒ½åœ°å°†è¿™äº›å†…å®¹åˆ†é…åˆ°PPTæ¨¡æ¿çš„è¡¨æ ¼å ä½ç¬¦ä¸­ã€‚

**æ ¸å¿ƒä»»åŠ¡ï¼š**
1. ç†è§£ç”¨æˆ·æ–‡æœ¬çš„ç»“æ„å’Œå†…å®¹
2. è¯†åˆ«æ–‡æœ¬ä¸­å¯ä»¥æå–çš„æ•°æ®ã€æ¦‚å¿µã€åˆ†ç±»ç­‰ä¿¡æ¯
3. å°†è¿™äº›ä¿¡æ¯åˆç†åœ°åˆ†é…åˆ°è¡¨æ ¼çš„å„ä¸ªå ä½ç¬¦ä½ç½®

**è¡¨æ ¼å¡«å……åŸåˆ™ï¼š**
1. **é€»è¾‘åŒ¹é…**ï¼šæ ¹æ®å ä½ç¬¦åç§°å’Œä½ç½®ï¼Œæ¨æ–­åº”è¯¥å¡«å…¥ä»€ä¹ˆç±»å‹çš„å†…å®¹
2. **æ•°æ®æå–**ï¼šä»ç”¨æˆ·æ–‡æœ¬ä¸­æå–ç»“æ„åŒ–ä¿¡æ¯ï¼ˆå¦‚ï¼šåç§°ã€æ•°å€¼ã€åˆ†ç±»ã€æè¿°ç­‰ï¼‰
3. **åˆç†åˆ†é…**ï¼šç¡®ä¿è¡¨æ ¼å†…å®¹çš„é€»è¾‘æ€§å’Œä¸€è‡´æ€§
4. **å®Œæ•´åˆ©ç”¨**ï¼šå°½å¯èƒ½åˆ©ç”¨ç”¨æˆ·æä¾›çš„æ–‡æœ¬ä¿¡æ¯
5. **æ ¼å¼é€‚é…**ï¼šç”Ÿæˆé€‚åˆè¡¨æ ¼å•å…ƒæ ¼çš„ç®€æ´å†…å®¹

**ç‰¹æ®Šå¤„ç†è§„åˆ™ï¼š**
- è¡¨å¤´å ä½ç¬¦ï¼ˆå¦‚ç¬¬ä¸€è¡Œï¼‰ï¼šå¡«å…¥åˆ†ç±»ã€æ ‡é¢˜ã€å­—æ®µåç­‰
- æ•°æ®è¡Œå ä½ç¬¦ï¼šå¡«å…¥å…·ä½“çš„æ•°æ®ã€æ¡ˆä¾‹ã€å†…å®¹ç­‰
- ç»Ÿè®¡å ä½ç¬¦ï¼šå¡«å…¥æ•°é‡ã€ç™¾åˆ†æ¯”ã€æ€»è®¡ç­‰ç»Ÿè®¡ä¿¡æ¯
- æè¿°å ä½ç¬¦ï¼šå¡«å…¥è¯¦ç»†è¯´æ˜ã€å¤‡æ³¨ç­‰

**è¾“å‡ºæ ¼å¼ï¼š**
è¯·è¿”å›JSONæ ¼å¼çš„åˆ†é…æ–¹æ¡ˆï¼š

```json
{
  "assignments": [
    {
      "placeholder": "å ä½ç¬¦åç§°",
      "content": "å¡«å……å†…å®¹",
      "slide_number": 1,
      "table_position": "è¡Œ1åˆ—2",
      "reasoning": "åˆ†é…ç†ç”±"
    }
  ]
}
```"""
    
    def _get_table_structure_for_ai(self) -> str:
        """è·å–è¡¨æ ¼ç»“æ„ä¿¡æ¯ï¼Œä¾›AIåˆ†æä½¿ç”¨"""
        structure_info = "**å½“å‰æ¨¡æ¿ä¸­çš„è¡¨æ ¼ç»“æ„ï¼š**\n\n"
        
        for slide_idx, tables in self.table_info.items():
            structure_info += f"**ç¬¬{slide_idx + 1}é¡µï¼š**\n"
            
            for table_idx, table_data in enumerate(tables):
                structure_info += f"  è¡¨æ ¼{table_idx + 1}ï¼ˆ{table_data['rows']}è¡Œ{table_data['cols']}åˆ—ï¼‰ï¼š\n"
                
                for placeholder_info in table_data['placeholders']:
                    placeholder = placeholder_info['placeholder']
                    position = placeholder_info['position']
                    structure_info += f"    - {{{placeholder}}} (ä½ç½®: {position})\n"
                
                structure_info += "\n"
        
        return structure_info
    
    def _parse_table_assignments(self, ai_response: str) -> Dict[str, Any]:
        """è§£æAIè¿”å›çš„è¡¨æ ¼åˆ†é…æ–¹æ¡ˆ"""
        try:
            # å°è¯•æå–JSONå†…å®¹
            json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', ai_response, re.DOTALL)
            if json_match:
                json_str = json_match.group(1)
            else:
                # å¦‚æœæ²¡æœ‰ä»£ç å—ï¼Œå°è¯•ç›´æ¥è§£æ
                json_str = ai_response.strip()
            
            result = json.loads(json_str)
            
            # éªŒè¯æ ¼å¼
            if 'assignments' in result and isinstance(result['assignments'], list):
                return result
            else:
                raise ValueError("AIè¿”å›æ ¼å¼ä¸æ­£ç¡®")
                
        except Exception as e:
            log_user_action("è¡¨æ ¼åˆ†é…æ–¹æ¡ˆè§£æå¤±è´¥", str(e))
            return {
                "error": f"è§£æAIå“åº”å¤±è´¥: {str(e)}",
                "assignments": []
            }
    
    def apply_table_assignments(self, assignments: Dict[str, Any]) -> Tuple[bool, List[str]]:
        """åº”ç”¨è¡¨æ ¼å¡«å……æ–¹æ¡ˆ"""
        if not self.presentation or not assignments.get('assignments'):
            return False, ["æ²¡æœ‰å¯åº”ç”¨çš„åˆ†é…æ–¹æ¡ˆ"]
        
        results = []
        filled_count = 0
        
        try:
            for assignment in assignments['assignments']:
                placeholder = assignment.get('placeholder', '')
                content = assignment.get('content', '')
                slide_number = assignment.get('slide_number', 1) - 1  # è½¬æ¢ä¸º0ç´¢å¼•
                
                if not placeholder or not content:
                    continue
                
                # åœ¨æŒ‡å®šå¹»ç¯ç‰‡çš„è¡¨æ ¼ä¸­æŸ¥æ‰¾å¹¶æ›¿æ¢å ä½ç¬¦
                success = self._fill_table_placeholder(slide_number, placeholder, content)
                if success:
                    filled_count += 1
                    results.append(f"âœ… {{{placeholder}}} -> {content[:30]}...")
                else:
                    results.append(f"âŒ {{{placeholder}}} å¡«å……å¤±è´¥")
            
            # å¡«å……å®Œæˆåè‡ªåŠ¨æ¸…ç†æœªä½¿ç”¨çš„å ä½ç¬¦
            if filled_count > 0:
                cleanup_result = self.cleanup_unfilled_table_placeholders()
                if cleanup_result.get('success'):
                    cleaned_count = cleanup_result.get('cleaned_count', 0)
                    if cleaned_count > 0:
                        results.append(f"ğŸ§¹ è‡ªåŠ¨æ¸…ç†äº† {cleaned_count} ä¸ªæœªä½¿ç”¨çš„å ä½ç¬¦")
                
                results.append(f"\næ€»å…±æˆåŠŸå¡«å…… {filled_count} ä¸ªè¡¨æ ¼å ä½ç¬¦")
                return True, results
            else:
                return False, ["æ²¡æœ‰æˆåŠŸå¡«å……ä»»ä½•å ä½ç¬¦"]
                
        except Exception as e:
            log_user_action("åº”ç”¨è¡¨æ ¼åˆ†é…æ–¹æ¡ˆå¤±è´¥", str(e))
            return False, [f"åº”ç”¨å¤±è´¥: {str(e)}"]
    
    def _fill_table_placeholder(self, slide_idx: int, placeholder: str, content: str) -> bool:
        """åœ¨æŒ‡å®šå¹»ç¯ç‰‡çš„è¡¨æ ¼ä¸­å¡«å……å ä½ç¬¦"""
        try:
            if slide_idx >= len(self.presentation.slides):
                return False
            
            slide = self.presentation.slides[slide_idx]
            placeholder_pattern = f"{{{placeholder}}}"
            
            # éå†è¯¥å¹»ç¯ç‰‡ä¸­çš„æ‰€æœ‰è¡¨æ ¼
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type') and shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    table = shape.table
                    
                    for row in table.rows:
                        for cell in row.cells:
                            if placeholder_pattern in cell.text:
                                # æ›¿æ¢å ä½ç¬¦ï¼Œä¿æŒå…¶ä»–å†…å®¹ä¸å˜
                                cell.text = cell.text.replace(placeholder_pattern, content)
                                return True
            
            return False
            
        except Exception as e:
            log_user_action(f"å¡«å……è¡¨æ ¼å ä½ç¬¦å¤±è´¥", f"{placeholder}: {str(e)}")
            return False
    
    def cleanup_unfilled_table_placeholders(self) -> Dict[str, Any]:
        """æ¸…ç†è¡¨æ ¼ä¸­æœªå¡«å……çš„å ä½ç¬¦"""
        if not self.presentation:
            return {"error": "PPTæœªåŠ è½½"}
        
        cleaned_count = 0
        cleaned_placeholders = []
        
        try:
            for slide_idx, slide in enumerate(self.presentation.slides):
                for shape in slide.shapes:
                    if hasattr(shape, 'shape_type') and shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                        table = shape.table
                        
                        for row_idx, row in enumerate(table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                original_text = cell.text
                                
                                # æŸ¥æ‰¾æœªå¡«å……çš„å ä½ç¬¦
                                placeholders = re.findall(r'\{([^}]+)\}', original_text)
                                if placeholders:
                                    cleaned_text = original_text
                                    for placeholder in placeholders:
                                        pattern = f"{{{placeholder}}}"
                                        cleaned_text = cleaned_text.replace(pattern, "")
                                        cleaned_placeholders.append(
                                            f"ç¬¬{slide_idx+1}é¡µè¡¨æ ¼({row_idx+1},{col_idx+1}): {{{placeholder}}}"
                                        )
                                    
                                    # æ¸…ç†å¤šä½™ç©ºç™½
                                    cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()
                                    cell.text = cleaned_text
                                    cleaned_count += 1
            
            return {
                "success": True,
                "cleaned_count": len(cleaned_placeholders),
                "cleaned_placeholders": cleaned_placeholders,
                "message": f"æ¸…ç†äº†{len(cleaned_placeholders)}ä¸ªè¡¨æ ¼å ä½ç¬¦"
            }
            
        except Exception as e:
            return {"error": f"æ¸…ç†å¤±è´¥: {str(e)}"}
    
    def get_ppt_bytes(self) -> bytes:
        """è·å–å¤„ç†åçš„PPTæ–‡ä»¶å­—èŠ‚æ•°æ®"""
        if not self.presentation:
            raise ValueError("PPTæœªåŠ è½½")
        
        import io
        ppt_bytes = io.BytesIO()
        self.presentation.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes.getvalue()


class TableTextProcessor:
    """è¡¨æ ¼æ–‡æœ¬å¤„ç†å™¨ - è´Ÿè´£æ–‡æœ¬è§£æå’Œç»“æ„åŒ–"""
    
    @staticmethod
    def parse_table_data_from_text(text: str) -> Dict[str, Any]:
        """ä»æ–‡æœ¬ä¸­è§£æå¯èƒ½çš„è¡¨æ ¼æ•°æ®"""
        # å°è¯•è¯†åˆ«æ–‡æœ¬ä¸­çš„è¡¨æ ¼ç»“æ„æ•°æ®
        # ä¾‹å¦‚ï¼šåˆ—è¡¨ã€åˆ†ç±»ã€é”®å€¼å¯¹ç­‰
        
        # æŒ‰è¡Œåˆ†å‰²
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        
        # å°è¯•è¯†åˆ«ä¸åŒçš„æ•°æ®æ¨¡å¼
        patterns = {
            'list_items': [],      # åˆ—è¡¨é¡¹
            'key_values': [],      # é”®å€¼å¯¹
            'categories': [],      # åˆ†ç±»ä¿¡æ¯
            'numbers': [],         # æ•°å€¼ä¿¡æ¯
            'names': [],           # åç§°ä¿¡æ¯
            'descriptions': []     # æè¿°ä¿¡æ¯
        }
        
        for line in lines:
            # æ£€æµ‹åˆ—è¡¨é¡¹ï¼ˆå¦‚ï¼š- xxx, 1. xxxï¼‰
            if re.match(r'^[-â€¢\d+\.]\s*', line):
                patterns['list_items'].append(re.sub(r'^[-â€¢\d+\.]\s*', '', line))
            
            # æ£€æµ‹é”®å€¼å¯¹ï¼ˆå¦‚ï¼šåç§°ï¼šå€¼ï¼‰
            elif ':' in line or 'ï¼š' in line:
                parts = re.split(r'[:ï¼š]', line, 1)
                if len(parts) == 2:
                    patterns['key_values'].append({
                        'key': parts[0].strip(),
                        'value': parts[1].strip()
                    })
            
            # æ£€æµ‹æ•°å€¼ä¿¡æ¯
            elif re.search(r'\d+[%ï¼…]?', line):
                patterns['numbers'].append(line)
            
            # å…¶ä»–å†…å®¹ä½œä¸ºæè¿°
            else:
                patterns['descriptions'].append(line)
        
        return patterns
    
    @staticmethod
    def suggest_table_structure(data_patterns: Dict[str, Any], table_size: Tuple[int, int]) -> List[Dict[str, Any]]:
        """æ ¹æ®æ•°æ®æ¨¡å¼å»ºè®®è¡¨æ ¼å¡«å……ç»“æ„"""
        rows, cols = table_size
        suggestions = []
        
        # å¦‚æœæœ‰é”®å€¼å¯¹æ•°æ®ï¼Œå»ºè®®ç”¨äºä¸¤åˆ—è¡¨æ ¼
        if data_patterns['key_values'] and cols >= 2:
            for i, kv in enumerate(data_patterns['key_values'][:rows]):
                suggestions.extend([
                    {'row': i, 'col': 0, 'content': kv['key'], 'type': 'key'},
                    {'row': i, 'col': 1, 'content': kv['value'], 'type': 'value'}
                ])
        
        # å¦‚æœæœ‰åˆ—è¡¨é¡¹ï¼Œå»ºè®®å¡«å……åˆ°å•åˆ—æˆ–å¤šåˆ—
        elif data_patterns['list_items']:
            items_per_col = (len(data_patterns['list_items']) + cols - 1) // cols
            for i, item in enumerate(data_patterns['list_items'][:rows * cols]):
                row = i // cols
                col = i % cols
                if row < rows:
                    suggestions.append({
                        'row': row, 'col': col, 'content': item, 'type': 'list_item'
                    })
        
        return suggestions