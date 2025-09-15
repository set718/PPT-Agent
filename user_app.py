#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文本转PPT填充器 - 用户版Web界面
使用OpenAI GPT-5将文本填入现有PPT文件
集成AI智能分页与Dify-模板桥接功能
"""

import streamlit as st
import os
import sys

# 强制设置UTF-8编码
import locale
try:
    locale.setlocale(locale.LC_ALL, 'zh_CN.UTF-8')
except:
    try:
        locale.setlocale(locale.LC_ALL, 'C.UTF-8')
    except:
        pass

# 设置环境变量
os.environ['PYTHONIOENCODING'] = 'utf-8'
if hasattr(sys, 'setdefaultencoding'):
    sys.setdefaultencoding('utf-8')
import subprocess
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import re
from typing import Dict, List, Any, Optional
from config import get_config
from utils import AIProcessor, PPTProcessor, FileManager, PPTAnalyzer
from logger import get_logger, log_user_action, log_file_operation, LogContext
from table_text_filler import TableTextFiller, TableTextProcessor

# 依赖检查和安装函数
def check_dependencies_light():
    """轻量级依赖检查（不安装）"""
    try:
        import streamlit
        import pptx
        return True
    except ImportError:
        return False

def check_system_requirements():
    """检查系统要求"""
    print("🔍 检查系统要求...")
    
    # 检查Python版本
    if sys.version_info < (3, 8):
        print("❌ Python版本过低，需要Python 3.8或更高版本")
        return False
    
    print("✅ Python版本检查通过")
    
    # 检查必要的目录和文件
    required_files = [
        'config.py',
        'utils.py',
        'logger.py',
        'ai_page_splitter.py',
        'dify_template_bridge.py',
        'dify_api_client.py'
    ]
    
    missing_files = []
    for file in required_files:
        if not os.path.exists(file):
            missing_files.append(file)
    
    if missing_files:
        print(f"❌ 缺少必要的文件: {', '.join(missing_files)}")
        return False
    
    print("✅ 必要文件检查通过")
    
    # 检查模板目录
    templates_dir = os.path.join("templates", "ppt_template")
    if not os.path.exists(templates_dir):
        print("❌ 模板目录不存在: templates/ppt_template/")
        return False
    
    template_files = [f for f in os.listdir(templates_dir) if f.startswith("split_presentations_") and f.endswith(".pptx")]
    if len(template_files) == 0:
        print("❌ 模板目录中没有找到可用的PPT模板文件")
        return False
    
    print(f"✅ 模板库检查通过，发现 {len(template_files)} 个模板文件")
    
    return True

def check_dify_api_keys():
    """检查Dify API密钥配置，返回(是否有效, 有效密钥数量, 错误消息)"""
    import os
    
    dify_keys = [os.getenv(f"DIFY_API_KEY_{i}") for i in range(1, 6)]
    valid_dify_keys = [key for key in dify_keys if key]
    
    if len(valid_dify_keys) == 0:
        return False, 0, "⚠️ **Dify API密钥未配置**\n\n请配置环境变量 `DIFY_API_KEY_1` 到 `DIFY_API_KEY_5`。\n\n**配置方法：**\n1. 复制 `.env.example` 为 `.env`\n2. 填入实际的API密钥\n3. 重启应用\n\n详细说明请查看 `ENVIRONMENT_SETUP.md`"
    elif len(valid_dify_keys) < 5:
        return True, len(valid_dify_keys), f"⚠️ 当前配置了 {len(valid_dify_keys)}/5 个Dify API密钥，建议配置全部5个以获得最佳性能"
    else:
        return True, len(valid_dify_keys), None

def initialize_system():
    """轻量级系统初始化"""
    # 只做基础检查，不执行耗时操作
    if not check_dependencies_light():
        return False
    
    # 检查基础文件存在
    required_files = ['config.py', 'utils.py', 'logger.py']
    for file in required_files:
        if not os.path.exists(file):
            return False
    
    return True

def show_results_section(pages, page_results):
    """显示处理结果部分"""
    
    # PPT下载区域
    st.markdown("### 📥 下载完整PPT")
    pages_count = len(pages) if pages else len(page_results)
    
    # 初始化session state
    if 'ppt_merge_result' not in st.session_state:
        st.session_state.ppt_merge_result = None
    
    # 检查PPT整合结果
    if st.session_state.ppt_merge_result:
        merge_result = st.session_state.ppt_merge_result
        
        # 检查是否为分批处理结果
        is_batch_result = merge_result.get("batch_files") is not None
        
        if is_batch_result:
            # 分批处理结果显示
            st.success("🎉 PPT分批整合成功！")
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("📄 总页数", merge_result["total_pages"])
            with col2:
                st.metric("✅ 成功页面", merge_result["processed_pages"])
            with col3:
                st.metric("📦 成功批次", merge_result.get("successful_batches", 0))
            with col4:
                st.metric("🔄 总批次数", merge_result.get("batch_count", 0))
            
            # 显示分批处理说明
            st.info(f"📋 由于页面数超过10页，系统自动分批处理（每批最多10页），生成了 {len(merge_result['batch_files'])} 个PPT文件")
            
            # 显示所有批次文件的下载按钮
            st.markdown("### 📥 下载分批文件")
            from datetime import datetime
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            for batch_info in merge_result["batch_files"]:
                batch_index = batch_info["batch_index"]
                batch_name = batch_info["batch_name"]
                file_size_mb = batch_info["file_size_mb"]
                pages_in_batch = batch_info["pages_in_batch"]
                presentation_bytes = batch_info["presentation_bytes"]
                actual_start_page = batch_info.get("actual_start_page", (batch_index - 1) * 10 + 1)
                actual_end_page = batch_info.get("actual_end_page", min(batch_index * 10, merge_result["total_pages"]))
                
                with st.container():
                    col1, col2 = st.columns([3, 1])
                    with col1:
                        st.write(f"**第{batch_index}批次: 第{actual_start_page}-{actual_end_page}页** ({pages_in_batch}页, {file_size_mb:.2f}MB)")
                    with col2:
                        filename = f"AI智能生成PPT_第{actual_start_page}-{actual_end_page}页_{timestamp}.pptx"
                        st.download_button(
                            label=f"📥 下载第{actual_start_page}-{actual_end_page}页",
                            data=presentation_bytes,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key=f"download_batch_{batch_index}"
                        )
            
            st.markdown('<div class="success-box">🎉 <strong>PPT分批生成完成！</strong><br><br><strong>1. 分批说明：</strong>由于页面超过10页，系统自动分批处理以确保最佳质量，请分别下载各批次文件。<br><strong>2. 备注查看提示：</strong>您提供的原始文本已完整放置在每一页PPT的"备注"栏中，方便您核对和修改内容。<br><strong>3. 文本缩略说明：</strong>出于美观，部分填充处会限制填充字数，以...代替。因此您的原始文本会被截断，您可以根据备注里保留的原始文本自行调整。<br><strong>4. 水印处理提示：</strong>下载的PPT文件包含水印。请前往【PPT去水印工具】功能页面上传文件进行处理。</div>', unsafe_allow_html=True)
        else:
            # 单文件结果显示（原有逻辑）
            st.success("🎉 PPT整合成功！")
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("📄 总页数", merge_result["total_pages"])
            with col2:
                st.metric("✅ 成功页面", merge_result["processed_pages"])
            with col3:
                st.metric("⚠️ 跳过页面", merge_result["skipped_pages"])
            with col4:
                ppt_size_mb = len(merge_result["presentation_bytes"]) / (1024 * 1024)
                st.metric("📦 文件大小", f"{ppt_size_mb:.2f}MB")
            
            # 提供下载
            if merge_result["presentation_bytes"]:
                from datetime import datetime
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"AI智能生成PPT_{timestamp}.pptx"
                
                col1, col2, col3 = st.columns([1, 2, 1])
                with col2:
                    st.download_button(
                        label="📥 下载完整PPT文件",
                        data=merge_result["presentation_bytes"],
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                        key="download_merged_ppt"
                    )
                
                st.markdown('<div class="success-box">🎉 <strong>PPT自动生成完成！</strong><br><br><strong>1. 备注查看提示：</strong>您提供的原始文本已完整放置在每一页PPT的"备注"栏中，方便您核对和修改内容。<br><strong>2. 文本缩略说明：</strong>出于美观，部分填充处会限制填充字数，以...代替。因此您的原始文本会被截断，您可以根据备注里保留的原始文本自行调整。<br><strong>3. 水印处理提示：</strong>下载的PPT文件包含水印。请前往【PPT去水印工具】功能页面上传文件进行处理。</div>', unsafe_allow_html=True)
        
        # 显示错误信息（如果有）
        if merge_result.get("errors"):
            with st.expander("⚠️ 查看处理警告", expanded=False):
                for error in merge_result["errors"]:
                    st.warning(f"• {error}")
        
    
    else:
        # PPT整合正在进行或失败
        st.info("🔄 PPT正在自动整合中，请稍候...")
        st.markdown('<div class="info-box">📋 <strong>处理状态：</strong><br>• ✅ AI智能分页：成功将长文本分割为 {pages_count} 页<br>• ✅ 封面页处理：第1页自动使用 title_slides.pptx 固定模板<br>• ✅ Dify模板桥接：其他页面通过API获取最适合的模板<br>• 🔄 PPT整合：系统正在自动整合模板页面...<br>• ⏳ 请稍候：整合完成后将自动显示下载按钮</div>'.format(pages_count=pages_count), unsafe_allow_html=True)
    
    # 添加重新开始按钮
    st.markdown("---")
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        if st.button("🔄 重新开始", help="清除当前结果，重新输入内容", key="restart_process"):
            # 清除所有相关的session state
            if 'current_page_results' in st.session_state:
                del st.session_state.current_page_results
            if 'current_pages' in st.session_state:
                del st.session_state.current_pages
            if 'ppt_merge_result' in st.session_state:
                del st.session_state.ppt_merge_result
            if 'ppt_generation_completed' in st.session_state:
                del st.session_state.ppt_generation_completed
            st.rerun()
    
    # 调试信息
    with st.expander("🔍 查看完整处理数据（调试信息）", expanded=False):
        st.json({
            'pages': pages,
            'page_results': page_results
        })

# 获取配置 - 移除阻塞性初始化
config = get_config()
logger = get_logger()

# 云环境检测
def is_cloud_environment():
    """检测是否在云环境中运行"""
    return (os.getenv('STREAMLIT_CLOUD') or 
            '/home/adminuser/venv' in str(sys.executable) or
            '/mount/src/' in os.getcwd())

# 延迟初始化函数
@st.cache_resource
def lazy_initialize():
    """延迟初始化系统资源"""
    if is_cloud_environment():
        # 云环境只做基础检查
        return True
    else:
        # 本地环境执行完整初始化
        return initialize_system()

# 页面配置
st.set_page_config(
    page_title="AI PPT助手",
    page_icon="🎨",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# 自定义CSS样式
st.markdown("""
<style>
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        color: #2E86AB;
        text-align: center;
        margin-bottom: 0.5rem;
    }
    .sub-header {
        font-size: 1.3rem;
        color: #666;
        text-align: center;
        margin-bottom: 2rem;
    }
    .success-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #d4edda;
        border: 2px solid #c3e6cb;
        color: #155724;
        margin: 1rem 0;
    }
    .info-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #e8f4fd;
        border: 2px solid #bee5eb;
        color: #0c5460;
        margin: 1rem 0;
    }
    .warning-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #fff3cd;
        border: 2px solid #ffeaa7;
        color: #856404;
        margin: 1rem 0;
    }
    .error-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #f8d7da;
        border: 2px solid #f5c6cb;
        color: #721c24;
        margin: 1rem 0;
    }
    .feature-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #f8f9fa;
        border: 2px solid #e9ecef;
        margin: 1rem 0;
    }
    .steps-container {
        background-color: #f8f9fa;
        padding: 2rem;
        border-radius: 1rem;
        margin: 2rem 0;
    }
</style>
""", unsafe_allow_html=True)

class UserPPTGenerator:
    def __init__(self, api_key):
        """初始化生成器"""
        self.api_key = api_key
        # 直接传递给AIProcessor，让它处理内置密钥标识符
        self.ai_processor = AIProcessor(api_key)
        self.presentation = None
        self.ppt_processor = None
        self.ppt_structure = None
        logger.info(f"用户界面初始化PPT生成器")
    
    def load_ppt_from_path(self, ppt_path):
        """从文件路径加载PPT"""
        with LogContext(f"用户界面加载PPT文件"):
            try:
                # 验证文件
                is_valid, error_msg = FileManager.validate_ppt_file(ppt_path)
                if not is_valid:
                    return False, error_msg
                
                self.presentation = Presentation(ppt_path)
                self.ppt_processor = PPTProcessor(self.presentation)
                self.ppt_structure = self.ppt_processor.ppt_structure
                
                log_file_operation("load_ppt_user", ppt_path, "success")
                return True, "成功"
            except Exception as e:
                log_file_operation("load_ppt_user", ppt_path, "error", str(e))
                return False, str(e)
    
    def process_text_with_openai(self, user_text):
        """使用OpenAI API分析如何将用户文本填入PPT模板的占位符"""
        if not self.ppt_structure:
            return {"assignments": []}
        
        log_user_action("用户界面AI文本分析", f"文本长度: {len(user_text)}字符")
        return self.ai_processor.analyze_text_for_ppt(user_text, self.ppt_structure)
    
    def process_text_with_openai_enhanced(self, user_text):
        """使用增强的数字提取逻辑分析文本并填入PPT模板占位符"""
        if not self.ppt_structure:
            return {"assignments": []}
        
        log_user_action("用户界面增强AI文本分析", f"文本长度: {len(user_text)}字符")
        
        # 预处理：提取文本中的数字信息
        extracted_data = self._extract_numbers_and_data(user_text)
        
        # 使用专门的数字感知AI提示
        return self._analyze_text_with_number_extraction(user_text, extracted_data)
    
    def _extract_numbers_and_data(self, text: str):
        """从文本中提取数字和结构化数据"""
        import re
        
        extracted = {
            'numbers': [],          # 纯数字
            'percentages': [],      # 百分比
            'currencies': [],       # 货币/价格
            'dates': [],           # 日期
            'measurements': [],     # 尺寸/度量
            'ratios': [],          # 比例/分数
            'key_value_pairs': []   # 键值对
        }
        
        lines = text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 提取百分比
            percentages = re.findall(r'(\d+(?:\.\d+)?[%％])', line)
            extracted['percentages'].extend(percentages)
            
            # 提取货币/价格
            currencies = re.findall(r'(\d+(?:\.\d+)?(?:元|美元|USD|\$|￥))', line)
            extracted['currencies'].extend(currencies)
            
            # 提取日期
            dates = re.findall(r'(\d{4}年\d{1,2}月|\d{1,2}/\d{1,2}/\d{4}|\d{4}-\d{1,2}-\d{1,2})', line)
            extracted['dates'].extend(dates)
            
            # 提取尺寸/度量
            measurements = re.findall(r'(\d+(?:\.\d+)?(?:英寸|寸|cm|mm|米|MB|GB|TB))', line)
            extracted['measurements'].extend(measurements)
            
            # 提取纯数字（不包括已经匹配的特殊格式）
            pure_numbers = re.findall(r'\b(\d+(?:\.\d+)?)\b', line)
            # 过滤掉已经在其他类别中的数字
            for num in pure_numbers:
                if not any(num in item for item_list in [extracted['percentages'], extracted['currencies'], 
                          extracted['dates'], extracted['measurements']] for item in item_list):
                    extracted['numbers'].append(num)
            
            # 提取键值对
            if ':' in line or '：' in line:
                parts = re.split(r'[:：]', line, 1)
                if len(parts) == 2:
                    key = parts[0].strip()
                    value = parts[1].strip()
                    extracted['key_value_pairs'].append({'key': key, 'value': value})
        
        return extracted
    
    def _analyze_text_with_number_extraction(self, user_text: str, extracted_data):
        """使用数字提取信息进行AI分析"""
        # 构建专门的数字感知系统提示
        system_prompt = self._build_number_aware_prompt(extracted_data)
        
        # 创建PPT结构描述
        ppt_description = self.ai_processor._create_ppt_description(self.ppt_structure)
        
        # 组合完整提示
        full_prompt = f"{system_prompt}\n\n{ppt_description}\n\n用户原始文本：\n{user_text}"
        
        try:
            # 确保AI客户端已初始化
            self.ai_processor._ensure_client()
            
            # 调用AI进行分析
            model_info = self.ai_processor.config.get_model_info()
            
            if model_info.get('request_format') == 'dify_compatible':
                content = self.ai_processor._call_liai_api(system_prompt, f"{ppt_description}\n\n{user_text}")
            else:
                actual_model = model_info.get('actual_model', self.ai_processor.config.ai_model)
                
                response = self.ai_processor.client.chat.completions.create(
                    model=actual_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": f"{ppt_description}\n\n用户文本：\n{user_text}"}
                    ],
                    temperature=0.3,
                    max_tokens=4000,
                    stream=True
                )
                
                content = ""
                for chunk in response:
                    if chunk.choices and chunk.choices[0].delta and chunk.choices[0].delta.content:
                        content += chunk.choices[0].delta.content
                
                content = content.strip() if content else ""
            
            # 解析AI返回结果
            return self.ai_processor._extract_json_from_response(content, user_text)
            
        except Exception as e:
            log_user_action("数字感知AI分析失败", str(e))
            return {"error": f"AI分析失败: {str(e)}"}
    
    def _build_number_aware_prompt(self, extracted_data):
        """构建数字感知的AI提示"""
        prompt = """你是一个专业的PPT内容分配专家，特别擅长处理数字和数据信息。你的任务是将用户文本智能分配到PPT模板的占位符中，并特别注意数字的精确提取和分配。

**核心原则：**
1. **数字优先**：识别并单独提取所有数字信息，避免将包含数字的整段文本填入通用占位符
2. **精确匹配**：根据占位符名称推断应该填入的具体数据类型
3. **数据分离**：将数字、文本描述、标题等分别处理
4. **逻辑分配**：确保每个占位符得到最合适的内容

**数字处理规则：**
- 如果占位符名称包含"价格"、"金额"、"费用"等，优先填入货币数字
- 如果占位符名称包含"百分比"、"比例"、"率"等，优先填入百分比数字
- 如果占位符名称包含"数量"、"个数"、"量"等，优先填入纯数字
- 如果占位符名称包含"尺寸"、"大小"、"长度"等，优先填入度量数字
- 如果占位符名称包含"日期"、"时间"等，优先填入日期信息

**内容分配策略：**
- {标题}、{名称}、{title} -> 填入主要标题或名称文本（不包含数字）
- {内容}、{描述}、{content} -> 填入描述性文本（移除已单独提取的数字）
- {价格}、{金额}、{cost} -> 填入提取的货币数字
- {百分比}、{比例}、{percent} -> 填入百分比数字
- {数量}、{quantity}、{count} -> 填入纯数字
- {日期}、{时间}、{date} -> 填入日期信息"""
        
        # 添加提取到的数字信息
        if extracted_data:
            prompt += "\n\n**已提取的数字信息：**\n"
            
            if extracted_data['numbers']:
                prompt += f"- 纯数字: {', '.join(extracted_data['numbers'])}\n"
            if extracted_data['percentages']:
                prompt += f"- 百分比: {', '.join(extracted_data['percentages'])}\n"
            if extracted_data['currencies']:
                prompt += f"- 货币/价格: {', '.join(extracted_data['currencies'])}\n"
            if extracted_data['dates']:
                prompt += f"- 日期: {', '.join(extracted_data['dates'])}\n"
            if extracted_data['measurements']:
                prompt += f"- 度量/尺寸: {', '.join(extracted_data['measurements'])}\n"
            if extracted_data['key_value_pairs']:
                prompt += f"- 键值对: {len(extracted_data['key_value_pairs'])}组\n"
        
        prompt += """

**重要要求：**
1. 数字信息必须单独提取，不要将"价格999元的手机"整体填入{content}，而是将"999元"填入{价格}，"手机"相关描述填入{content}
2. 优先根据占位符名称的语义匹配对应的数据类型
3. 如果同一类型有多个数字，优先使用最相关或最重要的
4. 描述性占位符只填入文本内容，不包含已单独提取的数字

请严格按照以下JSON格式返回分配方案：

```json
{
  "assignments": [
    {
      "action": "replace_placeholder",
      "slide_index": 0,
      "placeholder": "占位符名称",
      "content": "填充内容",
      "reason": "分配理由"
    }
  ]
}
```

**字段说明：**
- action: 固定为 "replace_placeholder"
- slide_index: 幻灯片索引（从0开始）
- placeholder: 占位符名称（不包含大括号）
- content: 要填充的内容
- reason: 选择此内容的理由"""
        
        return prompt
    
    def apply_text_assignments(self, assignments, user_text: str = ""):
        """根据分配方案替换PPT模板中的占位符，并将原始文本添加到备注"""
        if not self.presentation or not self.ppt_processor:
            return False, ["PPT文件未正确加载"]
        
        log_user_action("用户界面应用文本分配", f"分配数量: {len(assignments.get('assignments', []))}")
        # 传递用户原始文本，用于添加到幻灯片备注
        results = self.ppt_processor.apply_assignments(assignments, user_text)
        
        # 文本填充完成，不立即美化
        return True, results
    
    def cleanup_unfilled_placeholders(self):
        """清理未填充的占位符"""
        if not self.ppt_processor:
            return {"error": "PPT处理器未初始化"}
        
        try:
            log_user_action("用户界面清理占位符", f"已填充: {len(self.ppt_processor.filled_placeholders)}")
            
            # 智能清理占位符，只清理未填充的
            cleanup_count = 0
            cleaned_placeholders = []
            
            for slide_idx, slide in enumerate(self.presentation.slides):
                # 获取该页已填充的占位符
                filled_placeholders_in_slide = self.ppt_processor.filled_placeholders.get(slide_idx, set())
                
                for shape in slide.shapes:
                    # 处理普通文本框
                    if hasattr(shape, 'text') and shape.text:
                        original_text = shape.text
                        
                        # 找出文本中的所有占位符 - 识别所有{}格式的占位符
                        import re
                        placeholder_matches = re.findall(r'\{([^}]+)\}', original_text)
                        
                        if placeholder_matches:
                            # 检查哪些占位符未被填充
                            unfilled_placeholders = [
                                p for p in placeholder_matches 
                                if p not in filled_placeholders_in_slide
                            ]
                            
                            # 只移除未填充的占位符
                            if unfilled_placeholders:
                                cleaned_text = original_text
                                for unfilled_placeholder in unfilled_placeholders:
                                    pattern = f"{{{unfilled_placeholder}}}"
                                    cleaned_text = cleaned_text.replace(pattern, "")
                                    cleaned_placeholders.append(f"第{slide_idx+1}页(文本框): {{{unfilled_placeholder}}}")
                                
                                # 清理多余的空白
                                cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()
                                
                                if cleaned_text != original_text:
                                    shape.text = cleaned_text
                                    cleanup_count += 1
                    
                    # 处理表格中的占位符
                    elif hasattr(shape, 'shape_type') and shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                original_cell_text = cell.text.strip()
                                if original_cell_text:
                                    # 找出表格单元格中的占位符
                                    placeholder_matches = re.findall(r'\{([^}]+)\}', original_cell_text)
                                    
                                    if placeholder_matches:
                                        # 检查哪些占位符未被填充
                                        unfilled_placeholders = [
                                            p for p in placeholder_matches 
                                            if p not in filled_placeholders_in_slide
                                        ]
                                        
                                        # 只移除未填充的占位符
                                        if unfilled_placeholders:
                                            cleaned_cell_text = original_cell_text
                                            for unfilled_placeholder in unfilled_placeholders:
                                                pattern = f"{{{unfilled_placeholder}}}"
                                                cleaned_cell_text = cleaned_cell_text.replace(pattern, "")
                                                cleaned_placeholders.append(f"第{slide_idx+1}页(表格{row_idx+1},{col_idx+1}): {{{unfilled_placeholder}}}")
                                            
                                            # 清理多余的空白
                                            cleaned_cell_text = re.sub(r'\s+', ' ', cleaned_cell_text).strip()
                                            
                                            if cleaned_cell_text != original_cell_text:
                                                cell.text = cleaned_cell_text
                                                cleanup_count += 1
            
            # 使用实际清理的占位符数量，而不是修改的文本框数量
            actual_cleaned_count = len(cleaned_placeholders)
            
            return {
                "success": True,
                "cleaned_placeholders": actual_cleaned_count,
                "cleaned_placeholder_list": cleaned_placeholders,
                "message": f"清理了{actual_cleaned_count}个占位符，涉及{cleanup_count}个文本框和表格单元格"
            }
            
        except Exception as e:
            log_user_action("用户界面清理占位符失败", str(e))
            return {"error": f"清理占位符失败: {e}"}
    
    def apply_basic_beautification(self):
        """应用基础美化"""
        if not self.ppt_processor:
            return {"error": "PPT处理器未初始化"}
        
        try:
            log_user_action("用户界面基础美化")
            # 进行基础的美化处理
            beautify_results = self.ppt_processor.beautify_presentation()
            
            return beautify_results
            
        except Exception as e:
            log_user_action("用户界面基础美化失败", str(e))
            return {"error": f"基础美化失败: {e}"}
    
    
    def get_ppt_bytes(self):
        """获取修改后的PPT字节数据"""
        if not self.presentation:
            raise ValueError("PPT文件未正确加载")
        
        log_user_action("用户界面获取PPT字节数据")
        return FileManager.save_ppt_to_bytes(self.presentation)

def display_processing_summary(optimization_results, cleanup_results):
    """显示处理结果摘要"""
    if not optimization_results or "error" in optimization_results:
        if "error" in optimization_results:
            st.warning(f"⚠️ 处理过程中出现问题: {optimization_results['error']}")
        return
    
    summary = optimization_results.get('summary', {})
    
    # 基础处理信息
    st.markdown("### 📊 处理结果摘要")
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        final_slide_count = summary.get('final_slide_count', 1)  # 默认至少1页
        st.metric("📑 最终页数", final_slide_count)
    
    with col2:
        # 显示手动清理的占位符数量
        cleanup_count = cleanup_results.get('cleaned_placeholders', 0) if cleanup_results else 0
        st.metric("🧹 清理占位符", cleanup_count)
    
    with col3:
        removed_empty_slides = summary.get('removed_empty_slides_count', 0)
        st.metric("🗑️ 删除空页", removed_empty_slides)
    
    with col4:
        reorganized_slides = summary.get('reorganized_slides_count', 0)
        st.metric("🔄 重新排版", reorganized_slides)
    

def main():
    import os
    # 延迟初始化系统
    if not lazy_initialize():
        st.error("❌ 系统初始化失败，请刷新页面重试")
        return
    
    # 页面标题
    st.markdown('<div class="main-header">🎨 AI PPT助手</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">智能将您的文本内容转换为精美的PPT演示文稿</div>', unsafe_allow_html=True)
    
    # 加载环境变量
    try:
        from dotenv import load_dotenv
        load_dotenv()
    except ImportError:
        pass
    
    # 模型选择区域
    st.markdown("### 🤖 选择AI模型")
    
    # 获取配置（在函数内重新获取）
    from config import get_config
    config = get_config()
    available_models = getattr(config, 'available_models', {})
    
    # 如果没有可用模型，使用默认配置
    if not available_models:
        available_models = {
            "deepseek-v3": {
                "name": "DeepSeek V3（非保密场景请选择此模型）",
                "description": "火山引擎提供的DeepSeek V3模型，支持中英文对话，性能优异",
                "cost": "",
                "base_url": "https://ark.cn-beijing.volces.com/api/v3",
                "api_provider": "Volces",
                "api_key_env": "ARK_API_KEY",
                "actual_model": "deepseek-v3-250324",
                "request_format": "streaming_compatible"
            }
        }
    
    model_options = {}
    for model_key, model_info in available_models.items():
        # 只对有成本信息的模型显示成本
        if model_info['cost']:
            display_name = f"{model_info['name']} ({model_info['cost']}成本)"
        else:
            display_name = model_info['name']
        model_options[display_name] = model_key
    
    model_col1, model_col2 = st.columns([2, 1])
    with model_col1:
        selected_display = st.selectbox(
            "选择适合您需求的AI模型：",
            options=list(model_options.keys()),
            index=0,
            help="不同模型有不同的功能特点"
        )
        
        selected_model = model_options[selected_display]
        model_info = available_models[selected_model]
        
        # 动态更新配置
        if selected_model != config.ai_model:
            config.set_model(selected_model)
    
    with model_col2:
        st.markdown("**模型对比**")
        if selected_model == "liai-chat":
            st.info("🏢 调用公司融合云AgentOps私有化模型\n🔒 数据安全保障\n🌐 需要连接公司网络使用")
        else:  # DeepSeek V3
            st.success("🚀 火山引擎DeepSeek V3模型\n⚡ 性能优异\n🌐 支持中英文对话")
    

    
    st.markdown("---")
    
    # API密钥输入区域
    st.markdown("### 🔑 开始使用")
    
    # 根据选择的模型动态显示API密钥输入信息
    current_model_info = config.get_model_info()
    api_provider = current_model_info.get('api_provider', 'OpenAI')
    api_key_url = current_model_info.get('api_key_url', 'https://platform.openai.com/api-keys')
    
    col1, col2 = st.columns([2, 1])
    with col1:
        if api_provider == "Liai":
            # Liai自动填充API密钥（从环境变量读取，无需显示任何提示）
            import random
            import os
            
            # 强制重新加载环境变量以确保读取到最新的.env文件
            try:
                from dotenv import load_dotenv
                import os
                
                # 尝试多个可能的路径
                script_dir = os.path.dirname(os.path.abspath(__file__))
                current_work_dir = os.getcwd()
                
                possible_paths = [
                    os.path.join(script_dir, '.env'),
                    os.path.join(current_work_dir, '.env'),
                    '.env'
                ]
                
                
                found_env = False
                for env_path in possible_paths:
                    
                    if os.path.exists(env_path):
                        try:
                            with open(env_path, 'r', encoding='utf-8') as f:
                                content = f.read()
                                liai_lines = [line for line in content.split('\n') if 'LIAI_API_KEY' in line and not line.strip().startswith('#')]
                                
                                if len(liai_lines) > 0:
                                    
                                    load_dotenv(dotenv_path=env_path, override=True, encoding='utf-8')
                                    found_env = True
                                    break
                                    
                        except Exception as e:
                            pass
                
            except ImportError as e:
                pass
            except Exception as e:
                pass
            
            liai_api_keys = []
            for i in range(1, 6):  # 读取LIAI_API_KEY_1到LIAI_API_KEY_5
                key_name = f"LIAI_API_KEY_{i}"
                key = os.getenv(key_name)
                if key:
                    liai_api_keys.append(key)
            
            if not liai_api_keys:
                st.error("❌ 未找到Liai API密钥配置，请检查环境变量")
                return
            
            # 随机选择一个API密钥
            auto_api_key = random.choice(liai_api_keys)
            api_key = auto_api_key  # 直接使用自动选择的密钥
        elif api_provider == "Volces":
            # 火山引擎从环境变量读取（无需显示任何提示）
            import os
            ark_keys = [os.getenv(f"ARK_API_KEY_{i}") for i in range(1, 6)]
            valid_keys = [key for key in ark_keys if key]
            
            if not valid_keys:
                st.error("❌ 未找到火山引擎API密钥配置，请检查环境变量ARK_API_KEY_1到ARK_API_KEY_5")
                return
            
            # 使用第一个可用密钥（实际轮询由ai_page_splitter处理）
            api_key = valid_keys[0]
        else:  # 其他平台需要用户输入
            placeholder_text = "sk-xxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxx"
            help_text = f"通过{api_provider}平台访问AI模型，API密钥不会被保存"
        
        # 只有需要用户输入密钥的情况才显示密钥输入框
        needs_user_input = api_provider not in ["Liai", "Volces"]
        
        if needs_user_input:
            api_key = st.text_input(
                f"请输入您的{api_provider} API密钥",
                type="password",
                placeholder=placeholder_text,
                help=help_text
            )
    with col2:
        # 只有需要用户输入密钥的情况才显示获取密钥链接
        if needs_user_input:
            st.markdown("**获取API密钥**")
            st.markdown(f"[🔗 {api_provider}平台]({api_key_url})")
        
        # API密钥测试按钮（只有需要用户输入时才显示）
        if needs_user_input and api_key and api_key.strip():
            if st.button("🔍 测试API密钥", help="快速验证密钥是否有效"):
                with st.spinner("正在验证API密钥..."):
                    try:
                        # 创建一个临时的AIProcessor来测试
                        test_processor = AIProcessor(api_key.strip())
                        test_processor._ensure_client()
                        st.success("✅ API密钥验证通过！")
                    except ValueError as e:
                        st.error(f"❌ API密钥验证失败: {str(e)}")
                    except Exception as e:
                        error_msg = str(e)
                        if hasattr(e, 'status_code'):
                            status_code = e.status_code
                            if status_code == 401:
                                st.error("❌ API认证失败 (401): API密钥无效")
                            elif status_code == 402:
                                st.error("❌ API付费限制 (402): 账户余额不足")
                            elif status_code == 429:
                                st.error("❌ API请求频率限制 (429): 请求过于频繁")
                            else:
                                st.error(f"❌ API错误 ({status_code}): 这是API服务的问题")
                        elif "authentication" in error_msg.lower() or "unauthorized" in error_msg.lower():
                            st.error("❌ API密钥认证失败，请检查密钥是否正确")
                        elif "network" in error_msg.lower() or "connection" in error_msg.lower():
                            st.error("❌ 网络连接异常，请检查网络连接")
                        else:
                            st.error("❌ API调用异常，这不是应用程序的问题")
                        st.error(f"详细错误: {error_msg}")
    
    # 检查API密钥
    if not api_key or not api_key.strip():
        # 显示功能介绍
        st.markdown("---")
        
        # 使用步骤
        st.markdown('<div class="steps-container">', unsafe_allow_html=True)
        st.markdown("### 📝 四步轻松制作PPT")
        
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.markdown("""
            **第一步：选择模型** 🤖
            - DeepSeek V3：火山引擎先进模型，非保密场景推荐
            - Liai Chat：保密信息专用模型，安全可靠
            """)
        
        with col2:
            st.markdown("""
            **第二步：准备API密钥** 🔑
            - 根据选择的模型注册相应平台账号
            - OpenAI/Liai平台获取API密钥
            - 在上方输入密钥
            """)
        
        with col3:
            st.markdown("""
            **第三步：输入内容** ✏️
            - 粘贴您的文本内容
            - 可以是任何主题
            - 无需特殊格式
            """)
        
        with col4:
            st.markdown("""
            **第四步：生成下载** 🚀
            - 点击开始处理
            - 等待AI智能分析
            - 下载精美PPT
            """)
        st.markdown('</div>', unsafe_allow_html=True)
        
        # 功能特色
        st.markdown("### ✨ 产品特色")
        
        col1, col2 = st.columns(2)
        with col1:
            st.markdown('<div class="feature-box">', unsafe_allow_html=True)
            st.markdown("""
            **🤖 AI智能分析**
            - 自动理解文本结构
            - 智能匹配PPT模板
            - 保持内容完整性
            """)
            st.markdown('</div>', unsafe_allow_html=True)
            
            st.markdown('<div class="feature-box">', unsafe_allow_html=True)
            st.markdown("""
            **🎨 专业美化**
            - 自动优化布局
            - 清理多余元素
            - 统一设计风格
            """)
            st.markdown('</div>', unsafe_allow_html=True)
        
        with col2:
            st.markdown('<div class="feature-box">', unsafe_allow_html=True)
            st.markdown("""
            **⚡ 快速高效**
            - 一键生成PPT
            - 无需手动排版
            - 节省大量时间
            """)
            st.markdown('</div>', unsafe_allow_html=True)
            
            st.markdown('<div class="feature-box">', unsafe_allow_html=True)
            st.markdown("""
            **📱 简单易用**
            - 界面简洁明了
            - 操作步骤清晰
            - 适合所有用户
            """)
            st.markdown('</div>', unsafe_allow_html=True)
        
        # 适用场景
        st.markdown("### 🎯 适用场景")
        
        scenario_col1, scenario_col2, scenario_col3, scenario_col4 = st.columns(4)
        with scenario_col1:
            st.markdown("**📚 学术报告**\n研究成果展示")
        with scenario_col2:
            st.markdown("**💼 商业提案**\n项目方案介绍")
        with scenario_col3:
            st.markdown("**🎓 教学课件**\n课程内容整理")
        with scenario_col4:
            st.markdown("**📊 工作汇报**\n数据结果展示")
        
        return
    
    # 验证API密钥格式（根据选择的API提供商）
    # 只对需要用户输入的API提供商进行格式验证
    if needs_user_input and api_key:
        if not api_key.startswith('sk-'):
            st.markdown('<div class="warning-box">⚠️ API密钥格式可能不正确，通常以"sk-"开头</div>', unsafe_allow_html=True)
            return
    # elif api_provider == "Liai":
    #     # Liai API密钥格式检查已移除，直接通过格式验证
    
    # 跳过系统默认模板检查，直接使用Dify API和模板库
    # 注释掉原有的模板检查，改为检查模板库是否可用
    templates_dir = os.path.join(os.getcwd(), "templates", "ppt_template")
    if not os.path.exists(templates_dir):
        st.markdown('<div class="error-box">❌ 模板库文件夹不存在，请检查templates/ppt_template目录</div>', unsafe_allow_html=True)
        return
    
    # 检查模板库中是否有可用的模板文件
    template_files = [f for f in os.listdir(templates_dir) if f.startswith("split_presentations_") and f.endswith(".pptx")]
    if len(template_files) == 0:
        st.markdown('<div class="error-box">❌ 模板库中没有找到可用的PPT模板文件</div>', unsafe_allow_html=True)
        return
    
    # 模板库检查通过，不显示成功提示
    
    # 初始化AI处理器（不依赖默认模板）
    try:
        with st.spinner("正在验证API密钥..."):
            # 直接初始化AI处理器用于Dify API调用
            ai_processor = AIProcessor(api_key)
            # 测试API密钥有效性
            ai_processor._ensure_client()
            
    except ValueError as e:
        if "API密钥" in str(e):
            st.markdown('<div class="error-box">❌ API密钥验证失败，请检查密钥是否正确</div>', unsafe_allow_html=True)
        else:
            st.markdown(f'<div class="error-box">❌ 初始化失败: {str(e)}</div>', unsafe_allow_html=True)
        return
    except Exception as e:
        error_msg = str(e)
        if "authentication" in error_msg.lower() or "unauthorized" in error_msg.lower():
            st.markdown('<div class="error-box">❌ API密钥认证失败，请检查密钥是否正确或是否有足够余额</div>', unsafe_allow_html=True)
        elif "network" in error_msg.lower() or "connection" in error_msg.lower():
            st.markdown('<div class="error-box">❌ 网络连接异常，请检查网络连接后重试</div>', unsafe_allow_html=True)
        else:
            st.markdown('<div class="error-box">❌ 系统初始化异常，请稍后重试</div>', unsafe_allow_html=True)
        st.error(f"详细错误: {error_msg}")
        return
    
    # AI助手初始化成功，不显示成功提示
    
    # 角色选择
    st.markdown("---")
    st.markdown("#### 👤 请选择您的角色")
    user_role = st.selectbox(
        "选择角色以获得相应的功能界面",
        options=["用户", "开发者"],
        help="用户：仅显示核心功能；开发者：显示全部功能包括测试工具",
        key="user_role_selectbox"
    )
    
    # 功能选择选项卡
    st.markdown("---")
    if user_role == "用户":
        # 用户角色：只显示核心功能
        tab1, tab_watermark = st.tabs(["🎨 智能PPT生成", "🧽 PPT去水印工具"])
    else:
        # 开发者角色：显示全部功能
        tab1, tab_watermark, tab3, tab_table, tab_format, tab_ai_test = st.tabs(["🎨 智能PPT生成", "🧽 PPT去水印工具", "🧪 自定义模板测试", "📊 表格文本填充", "🔍 PPT格式读取展示", "🤖 AI分页测试"])
    
    with tab1:
        # 智能PPT生成功能 - AI分页 + 模板匹配
        st.markdown("### 🚀 智能PPT生成")
        
        # 检查是否有保存的处理结果
        if 'current_page_results' in st.session_state and 'current_pages' in st.session_state:
            # 显示保存的结果
            pages = st.session_state.current_pages
            page_results = st.session_state.current_page_results
            
            st.markdown('<div class="success-box">🎉 智能PPT生成完成！</div>', unsafe_allow_html=True)
            
            # 跳转到结果显示部分
            show_results_section(pages, page_results)
        
        # 初始化变量
        user_text = ""
        process_button = False
        
        # 只在没有生成结果时显示输入界面
        if 'current_page_results' not in st.session_state or 'current_pages' not in st.session_state:
            st.markdown('''<div class="info-box">📋 <strong>使用前请知悉</strong><br>
            <strong>1. 文本生成说明：</strong>本产品专注于PPT样式的智能生成。您提供的原始文本将被直接使用，AI不会对其进行修改或扩充。如果您需要AI辅助撰写或优化文本，欢迎在后续的问卷中向我们反馈该需求。<br>
            <strong>2. 图表支持说明：</strong>请注意，当前版本暂不支持自动图表数据填充。如果您的文本中包含图表描述，生成后需要您手动补充相关数据。<br>
            <strong>3. 样式调整说明：</strong>出于兼容性考虑，生成PPT的字体、大小和颜色等样式可能有瑕疵，您可以在下载后手动进行美化调整。<br>
            <strong>4. 问题与支持：</strong>使用过程中如有任何问题或建议，请随时联系 @贾轶涵 获取帮助。
            </div>''', unsafe_allow_html=True)
        
            # 文本输入
            st.markdown("#### 📝 输入您的内容")
            
            user_text = st.text_area(
                "请输入您想要制作成PPT的文本内容：",
                height=250,
                placeholder="""例如：

人工智能的发展历程与未来趋势

人工智能技术的发展经历了多个重要阶段。从1950年代的符号主义开始，强调逻辑推理和知识表示，到1980年代的专家系统兴起，再到近年来深度学习的突破性进展。

技术发展阶段：
- 符号主义时代：基于规则和逻辑推理
- 连接主义时代：神经网络和机器学习
- 深度学习时代：大数据驱动的智能系统
- 大模型时代：通用人工智能的探索

当前，大语言模型如GPT、Claude等展现出了前所未有的能力，能够进行复杂的文本理解、生成和推理。这些技术正在革新各个行业，从教育、医疗到金融、娱乐，都能看到AI的身影。

未来发展趋势：
人工智能将继续向更加智能化、人性化的方向发展，实现更好的人机协作，为人类社会带来更多便利和创新可能性。同时需要关注AI安全和伦理问题。""",
                help="AI将分析文本结构进行智能分页，每页内容调用AI模型获取对应模板"
            )

            # 分页选项 - 简化布局
            st.markdown("#### ⚙️ 分页选项")
            
            col1, col2 = st.columns([1, 1])
            with col1:
                # 使用与AI分页测试一致的下拉选择框
                page_options = ["AI自动判断"] + [str(i) for i in range(4, 26)]
                selected_option = st.selectbox(
                    "目标页面数量",
                    options=page_options,
                    index=0,
                    help="选择AI自动判断或手动设置页面数量（最少4页：封面+目录+内容+结尾）",
                    key="smart_ppt_page_count"
                )
                target_pages = 0 if selected_option == "AI自动判断" else int(selected_option)
            
            with col2:
                # 页面数量限制提醒 - 移至右侧
                st.info("📋 页面限制：最少4页（封面+目录+内容+结尾），最多25页")
            
            # 页数建议 - 使用更简洁的布局
            st.markdown("""
            <div style="background-color: #f0f2f6; padding: 0.75rem; border-radius: 0.5rem; margin: 0.5rem 0;">
            <small>💡 <strong>页数建议：</strong>
            5分钟演示：4-6页 • 10分钟演示：6-8页 • 15分钟演示：8-12页 • 30分钟演示：15-20页 • 学术报告：20-25页</small>
            </div>
            """, unsafe_allow_html=True)
            
            # 生成按钮 - 居中显示
            st.markdown("---")
            col1, col2, col3 = st.columns([1, 2, 1])
            with col2:
                process_button = st.button(
                    "🚀 开始生成PPT",
                    type="primary",
                    use_container_width=True,
                    disabled=not user_text.strip(),
                    help="AI智能分页 → 模板匹配 → 自动整合PPT → 可直接下载"
                )
        
        
        # 处理逻辑 - AI分页 + 智能模板匹配
        if process_button and user_text.strip():
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                # 步骤1：AI智能分页
                status_text.text("🤖 AI正在分析文本结构并进行智能分页...")
                progress_bar.progress(20)
                
                from ai_page_splitter import AIPageSplitter
                page_splitter = AIPageSplitter(api_key)
                # 验证页面数设置：手动设置时最少4页（封面+目录+内容+结尾）
                if target_pages > 0 and target_pages < 4:
                    st.error("❌ 页面数量不能少于4页（封面页+目录页+内容页+结尾页）")
                    return
                target_page_count = int(target_pages) if target_pages > 0 else None
                split_result = page_splitter.split_text_to_pages(user_text.strip(), target_page_count)
                
                if not split_result.get('success'):
                    st.error(f"❌ AI分页失败: {split_result.get('error', '未知错误')}")
                    return
                
                pages = split_result.get('pages', [])
                if not pages:
                    st.error("❌ 分页结果为空，请检查输入文本")
                    return
                
                
                # 显示简化的进度提示
                st.markdown('<div class="info-box">⏳ <strong>正在生成PPT...</strong><br><br><strong>1. 耐心等待：</strong>生成需要一些时间，请耐心等候。您可以最小化此页面，处理其他工作。<br><strong>2. 多文件生成提示：</strong>由于功能限制，如果您的PPT超过10页，系统会自动以10页为单位拆分成多个文件供您下载。</div>', unsafe_allow_html=True)
                
                # 步骤2：为每页内容调用AI模型获取模板
                status_text.text("🔗 正在为每页内容调用AI模型获取对应模板...")
                progress_bar.progress(40)
                
                # 检查Dify API密钥配置
                dify_valid, dify_count, dify_message = check_dify_api_keys()
                if not dify_valid:
                    st.error(dify_message)
                    return
                elif dify_message:  # 有警告消息
                    st.warning(dify_message)
                
                from dify_template_bridge import sync_test_dify_template_bridge
                from dify_api_client import BatchProcessor, DifyAPIConfig
                
                # 检查是否启用分批处理（超过5页时自动启用）
                if len(pages) > 5:
                    # 检查当前API提供商
                    current_model_info = st.session_state.get('selected_model_info', {})
                    api_provider = current_model_info.get('api_provider', 'OpenAI')
                    
                    if api_provider == "Liai":
                        # 使用Liai分批处理
                        try:
                            # 准备Liai批处理的页面数据
                            liai_pages_data = []
                            for i, page in enumerate(pages):
                                page_content = page.get('original_text_segment', '')
                                if not page_content:
                                    title = page.get('title', '')
                                    key_points = page.get('key_points', [])
                                    page_content = f"{title}\n\n" + "\n".join(key_points)
                                
                                page_type = page.get('page_type', 'content')
                                page_number = page.get('page_number', i + 1)
                                
                                # 跳过特殊页面，只处理需要AI分析的内容页
                                if page_type not in ['title', 'ending'] and page_number != 1 and len(pages) > 1:
                                    liai_pages_data.append({
                                        'page_number': page_number,
                                        'content': page_content,
                                        'ppt_structure': {'slides': [{'placeholders': {}}]},  # 简化的结构
                                        'page_data': page
                                    })
                            
                            # 如果有需要处理的页面，使用Liai批处理
                            if liai_pages_data:
                                st.info(f"🔄 开始Liai分批处理{len(liai_pages_data)}个内容页面，每批5个...")
                                
                                # 创建AI处理器并进行批处理
                                ai_processor = AIProcessor(api_key.strip())
                                batch_results = ai_processor.batch_analyze_pages_for_liai(liai_pages_data, 5)
                                
                                # 处理批处理结果
                                page_results = []
                                for result in batch_results:
                                    if result.get('success'):
                                        page_results.append({
                                            'page_number': result['page_number'],
                                            'content': result['content'],
                                            'template_number': 'liai_analyzed',
                                            'template_path': None,
                                            'template_filename': 'Liai智能分析',
                                            'dify_response': str(result.get('analysis_result', ''))[:200] + '...',
                                            'processing_time': result.get('processing_time', 0),
                                            'is_title_page': False
                                        })
                                    else:
                                        pass
                                
                        except Exception as e:
                            st.error(f"Liai分批处理出错: {str(e)}")
                            page_results = []
                    else:
                        # 使用Dify分批处理
                        if api_provider != "Liai":
                            try:
                                dify_config = DifyAPIConfig()
                                dify_config.batch_size = 5  # 每批5个
                                
                                page_results = []
                                batch_index = 0
                                
                                # 准备需要调用API的页面（排除title和ending页）
                                api_pages = []
                                for i, page in enumerate(pages):
                                    page_content = page.get('original_text_segment', '')
                                    if not page_content:
                                        title = page.get('title', '')
                                        key_points = page.get('key_points', [])
                                        page_content = f"{title}\n\n" + "\n".join(key_points)
                                    
                                    page_type = page.get('page_type', 'content')
                                    page_number = page.get('page_number', i + 1)
                                
                                    # 特殊页面处理
                                    if page_type == 'title' or page_number == 1:
                                        title_template_path = os.path.join("templates", "title_slides.pptx")
                                        page_results.append({
                                            'page_number': page_number,
                                            'content': page_content,
                                            'template_number': 'title',
                                            'template_path': title_template_path,
                                            'template_filename': "title_slides.pptx",
                                            'dify_response': '封面页使用固定标题模板',
                                            'processing_time': 0,
                                            'is_title_page': True
                                        })
                                    elif page_type == 'table_of_contents':
                                        toc_template_path = page.get('template_path', os.path.join("templates", "table_of_contents_slides.pptx"))
                                        page_results.append({
                                            'page_number': page_number,
                                            'content': page_content,
                                            'template_number': 'table_of_contents',
                                            'template_path': toc_template_path,
                                            'template_filename': "table_of_contents_slides.pptx",
                                            'dify_response': '目录页使用提取的内容页标题动态生成',
                                            'processing_time': 0,
                                            'is_toc_page': True
                                        })
                                    elif page_type == 'ending' or page.get('skip_dify_api', False):
                                        ending_template_path = page.get('template_path', os.path.join("templates", "ending_slides.pptx"))
                                        page_results.append({
                                            'page_number': page_number,
                                            'content': page_content,
                                            'template_number': 'ending',
                                            'template_path': ending_template_path,
                                            'template_filename': "ending_slides.pptx",
                                            'dify_response': '结尾页使用固定感谢模板',
                                            'processing_time': 0,
                                            'is_ending_page': True
                                        })
                                    elif page_content:
                                        # 需要调用API的页面
                                        api_pages.append({
                                            'page_index': i,
                                            'page_data': page,
                                            'page_content': page_content,
                                            'page_number': page_number
                                        })
                                
                                # 分批处理API调用
                                if api_pages:
                                    
                                    # 创建进度跟踪
                                    batch_progress = st.progress(0)
                                    batch_status = st.empty()
                                    
                                    total_batches = (len(api_pages) + 4) // 5  # 向上取整
                                    
                                    for batch_start in range(0, len(api_pages), 5):
                                        batch_end = min(batch_start + 5, len(api_pages))
                                        batch_pages = api_pages[batch_start:batch_end]
                                        batch_index += 1
                                        
                                        
                                        # 处理当前批次
                                        for page_info in batch_pages:
                                            # 合并title和content作为完整输入
                                            page_title = page_info['page_data'].get('title', '')
                                            page_content = page_info['page_content']
                                            full_content = f"标题: {page_title}\n\n{page_content}" if page_title else page_content
                                            
                                            # 获取当前模型配置
                                            from config import get_config
                                            current_config = get_config()
                                            model_config = current_config.get_model_info()
                                            bridge_result = sync_test_dify_template_bridge(full_content, model_config=model_config)
                                            
                                            # 如果成功且有title，强制添加title占位符填充
                                            if bridge_result.get('success') and page_title:
                                                step_3_result = bridge_result.get('step_3_template_fill', {})
                                                if step_3_result.get('success'):
                                                    assignments = step_3_result.get('assignments', {}).get('assignments', [])
                                                    # 直接添加title占位符填充（内容页都有title占位符）
                                                    assignments.append({
                                                        'action': 'replace_placeholder',
                                                        'slide_index': 0,
                                                        'placeholder': 'title',
                                                        'content': page_title,
                                                        'reason': '自动填充页面标题'
                                                    })
                                            
                                            if bridge_result.get('success'):
                                                dify_result = bridge_result["step_1_dify_api"]
                                                template_result = bridge_result["step_2_template_lookup"]
                                                page_results.append({
                                                    'page_number': page_info['page_number'],
                                                    'content': page_info['page_content'],
                                                    'template_number': dify_result.get('template_number'),
                                                    'template_path': template_result.get('file_path'),
                                                    'template_filename': template_result.get('filename'),
                                                    'dify_response': dify_result.get('response_text', ''),
                                                    'processing_time': bridge_result.get('processing_time', 0),
                                                    'is_title_page': False
                                                })
                                            else:
                                                page_results.append({
                                                    'page_number': page_info['page_number'],
                                                    'content': page_info['page_content'],
                                                    'template_number': None,
                                                    'template_path': None,
                                                    'template_filename': None,
                                                    'dify_response': f"错误: {bridge_result.get('error')}",
                                                    'processing_time': bridge_result.get('processing_time', 0),
                                                    'is_title_page': False,
                                                    'error': True
                                                })
                                        
                                        # 更新进度
                                        progress = batch_index / total_batches
                                        batch_progress.progress(progress)
                                        
                                        # 批次间延迟
                                        if batch_index < total_batches:
                                            batch_status.text(f"⏳ 批次间等待{dify_config.batch_delay}秒...")
                                            import time
                                            time.sleep(dify_config.batch_delay)
                                    
                                    # 清理进度显示
                                    batch_progress.empty()
                                    batch_status.empty()
                                    
                            except Exception as e:
                                st.error(f"❌ 分批处理异常: {str(e)}")
                                st.info("🔄 降级到逐页处理模式...")
                                # 降级到原来的逐页处理
                                page_results = []
                                for i, page in enumerate(pages):
                                    # 原来的逐页处理逻辑
                                    page_content = page.get('original_text_segment', '')
                                    if not page_content:
                                        title = page.get('title', '')
                                        key_points = page.get('key_points', [])
                                        page_content = f"{title}\n\n" + "\n".join(key_points)
                                    
                                    page_type = page.get('page_type', 'content')
                                    page_number = page.get('page_number', i + 1)
                                    
                                    if page_type == 'title' or page_number == 1:
                                        title_template_path = os.path.join("templates", "title_slides.pptx")
                                        page_results.append({
                                            'page_number': page_number,
                                            'content': page_content,
                                            'template_number': 'title',
                                            'template_path': title_template_path,
                                            'template_filename': "title_slides.pptx",
                                            'dify_response': '封面页使用固定标题模板',
                                            'processing_time': 0,
                                            'is_title_page': True
                                        })
                                    elif page_type == 'ending' or page.get('skip_dify_api', False):
                                        ending_template_path = page.get('template_path', os.path.join("templates", "ending_slides.pptx"))
                                        page_results.append({
                                            'page_number': page_number,
                                            'content': page_content,
                                            'template_number': 'ending',
                                            'template_path': ending_template_path,
                                            'template_filename': "ending_slides.pptx",
                                            'dify_response': '结尾页使用固定感谢模板',
                                            'processing_time': 0,
                                            'is_ending_page': True
                                        })
                                    elif page_content:
                                        # 获取当前模型配置
                                        from config import get_config
                                        current_config = get_config()
                                        model_config = current_config.get_model_info()
                                        bridge_result = sync_test_dify_template_bridge(page_content, model_config=model_config)
                                        if bridge_result.get('success'):
                                            dify_result = bridge_result["step_1_dify_api"]
                                            template_result = bridge_result["step_2_template_lookup"]
                                            page_results.append({
                                                'page_number': page_number,
                                                'content': page_content,
                                                'template_number': dify_result.get('template_number'),
                                                'template_path': template_result.get('file_path'),
                                                'template_filename': template_result.get('filename'),
                                                'dify_response': dify_result.get('response_text', ''),
                                                'processing_time': bridge_result.get('processing_time', 0),
                                                'is_title_page': False
                                            })
                                        else:
                                            st.error("🚫 无法继续处理，请检查Dify API配置或稍后重试")
                                            return
                else:
                    # 页面数少于等于5页，使用原来的逐页处理
                    page_results = []
                    
                    for i, page in enumerate(pages):
                        # 获取页面内容，优先使用original_text_segment，如果没有则使用title和key_points组合
                        page_content = page.get('original_text_segment', '')
                        if not page_content:
                            # 如果没有original_text_segment，则组合title和key_points
                            title = page.get('title', '')
                            key_points = page.get('key_points', [])
                            page_content = f"{title}\n\n" + "\n".join(key_points)
                        
                        page_type = page.get('page_type', 'content')
                        page_number = page.get('page_number', i + 1)
                        
                        # 封面页直接使用 title_slides.pptx，不调用Dify API
                        if page_type == 'title' or page_number == 1:
                            title_template_path = os.path.join("templates", "title_slides.pptx")
                            page_results.append({
                                'page_number': page_number,
                                'content': page_content,
                                'template_number': 'title',
                                'template_path': title_template_path,
                                'template_filename': "title_slides.pptx",
                                'dify_response': '封面页使用固定标题模板',
                                'processing_time': 0,
                                'is_title_page': True
                            })
                        
                        # 目录页直接使用 table_of_contents_slides.pptx，不调用Dify API
                        elif page_type == 'table_of_contents':
                            toc_template_path = page.get('template_path', os.path.join("templates", "table_of_contents_slides.pptx"))
                            page_results.append({
                                'page_number': page_number,
                                'content': page_content,
                                'template_number': 'table_of_contents',
                                'template_path': toc_template_path,
                                'template_filename': "table_of_contents_slides.pptx",
                                'dify_response': '目录页使用提取的内容页标题动态生成',
                                'processing_time': 0,
                                'is_toc_page': True
                            })
                        
                        # 结尾页直接使用 ending_slides.pptx，不调用Dify API
                        elif page_type == 'ending' or page.get('skip_dify_api', False):
                            ending_template_path = page.get('template_path', os.path.join("templates", "ending_slides.pptx"))
                            page_results.append({
                                'page_number': page_number,
                                'content': page_content,
                                'template_number': 'ending',
                                'template_path': ending_template_path,
                                'template_filename': "ending_slides.pptx",
                                'dify_response': '结尾页使用固定感谢模板',
                                'processing_time': 0,
                                'is_ending_page': True
                            })
                        
                        elif page_content:
                            # 其他页面调用API（支持Dify和Liai）
                            # 获取当前模型配置
                            from config import get_config
                            current_config = get_config()
                            model_config = current_config.get_model_info()
                            bridge_result = sync_test_dify_template_bridge(page_content, model_config=model_config)
                            if bridge_result.get('success'):
                                dify_result = bridge_result["step_1_dify_api"]
                                template_result = bridge_result["step_2_template_lookup"]
                                page_results.append({
                                    'page_number': page_number,
                                    'content': page_content,
                                    'template_number': dify_result.get('template_number'),
                                    'template_path': template_result.get('file_path'),
                                    'template_filename': template_result.get('filename'),
                                    'dify_response': dify_result.get('response_text', ''),
                                    'processing_time': bridge_result.get('processing_time', 0),
                                    'is_title_page': False
                                })
                            else:
                                # 记录失败但继续处理其他页面
                                page_results.append({
                                    'page_number': page_number,
                                    'content': page_content,
                                    'template_number': None,
                                    'template_path': None,
                                    'template_filename': None,
                                    'dify_response': f'处理失败: {bridge_result.get("error", "未知错误")}',
                                    'processing_time': bridge_result.get('processing_time', 0),
                                    'is_title_page': False,
                                    'error': True
                                })
                
                # 步骤3：文本填充（新增）
                status_text.text("📝 正在对每个模板进行智能文本填充...")
                progress_bar.progress(70)
                
                filled_page_results = []
                from pptx import Presentation
                
                # 导入PPT处理器（AIProcessor已在文件顶部导入）
                from utils import PPTProcessor
                
                for i, page_result in enumerate(page_results):
                    try:
                        template_path = page_result.get('template_path')
                        page_content = page_result.get('content', '')
                        page_number = page_result.get('page_number', i+1)
                        
                        if template_path and os.path.exists(template_path):
                            # 加载模板
                            template_prs = Presentation(template_path)
                            
                            # 检查是否为结尾页（只有结尾页完全跳过文本填充）
                            if (page_result.get('is_ending_page') or page_result.get('page_type') == 'ending'):
                                # 结尾页直接使用模板，不进行文本填充
                                fill_results = []
                                print(f"🔍 跳过结尾页文本填充: 第{page_number}页")
                            else:
                                # 创建PPT处理器并进行文本填充
                                print(f"🔍 开始文本填充: 第{page_number}页 - {page_result.get('page_type', 'content')}")
                                print(f"📄 页面内容长度: {len(page_content)}字")
                                
                                processor = PPTProcessor(template_prs)
                                
                                # 使用完整的文本填充流程（会自动使用当前选择的AI模型）
                                # 1. 创建AI处理器来分析文本并生成分配方案
                                ai_processor = AIProcessor()
                                
                                # 2. 分析PPT结构
                                try:
                                    print(f"🔍 分析PPT结构...")
                                    print(f"📁 模板路径: {template_path}")
                                    print(f"📑 模板slides数量: {len(template_prs.slides)}")
                                    
                                    ppt_structure = PPTAnalyzer.analyze_ppt_structure(template_prs)
                                    # 从slides中收集所有占位符
                                    all_placeholders = {}
                                    for slide in ppt_structure.get('slides', []):
                                        all_placeholders.update(slide.get('placeholders', {}))
                                    
                                    print(f"📊 检测到占位符数量: {len(all_placeholders)}")
                                    if all_placeholders:
                                        print(f"🔍 占位符列表: {list(all_placeholders.keys())}")
                                    else:
                                        print(f"⚠️ 未检测到任何占位符，模板可能没有{{placeholder}}格式的内容")
                                    
                                    # 3. 生成文本分配方案
                                    print(f"🤖 调用AI生成分配方案...")
                                    assignments = ai_processor.analyze_text_for_ppt(page_content, ppt_structure)
                                    print(f"📋 生成分配方案数量: {len(assignments.get('assignments', []))}")
                                    
                                    # 4. 应用分配方案
                                    print(f"✏️ 应用分配方案...")
                                    fill_results = processor.apply_assignments(assignments, page_content)
                                    print(f"✅ 文本填充完成，结果数量: {len(fill_results)}")
                                except Exception as fill_error:
                                    print(f"❌ 文本填充过程异常: {fill_error}")
                                    st.error(f"文本填充过程异常: {fill_error}")
                                    fill_results = []
                            
                            # 更新结果信息（为合并器保存临时文件）
                            filled_result = page_result.copy()
                            filled_result['fill_results'] = fill_results
                            
                            # 为所有页面保存临时文件用于合并（确保合并器能正确处理）
                            import tempfile
                            temp_dir = tempfile.gettempdir()
                            filled_temp_path = os.path.join(temp_dir, f"filled_temp_{page_number}_{os.path.basename(template_path)}")
                            template_prs.save(filled_temp_path)
                            filled_result['template_path'] = filled_temp_path  # 使用处理后的临时文件路径
                            
                            filled_page_results.append(filled_result)
                            
                        else:
                            # 没有模板的页面直接传递
                            filled_page_results.append(page_result)
                            
                    except Exception as e:
                        # 失败时使用原始模板
                        filled_page_results.append(page_result)
                
                # 步骤4：清理未填充的占位符
                status_text.text("🧹 正在清理未填充的占位符...")
                progress_bar.progress(75)
                
                # 对每个填充后的页面调用现有的清理功能
                for filled_result in filled_page_results:
                    if filled_result.get('template_path') and os.path.exists(filled_result['template_path']):
                        try:
                            # 加载已填充的模板并创建临时生成器实例用于清理
                            filled_prs = Presentation(filled_result['template_path'])
                            
                            # 创建临时的UserPPTGenerator实例来调用清理方法
                            temp_generator = UserPPTGenerator(api_key)
                            temp_generator.presentation = filled_prs
                            temp_generator.ppt_processor = PPTProcessor(filled_prs)
                            
                            # 调用清理功能
                            cleanup_results = temp_generator.cleanup_unfilled_placeholders()
                            
                            # 保存清理后的结果
                            filled_prs.save(filled_result['template_path'])
                            
                        except Exception as cleanup_error:
                            print(f"⚠️ 第{filled_result.get('page_number')}页占位符清理失败: {cleanup_error}")
                
                # 步骤5：整合PPT页面
                status_text.text("🔗 正在整合填充后的PPT页面...")
                progress_bar.progress(80)
                
                # 保存填充后的页面结果到session state
                st.session_state.current_page_results = filled_page_results
                st.session_state.current_pages = pages
                
                # 自动执行PPT整合
                try:
                    # 使用增强版合并器，自动选择最佳方法
                    from ppt_merger import merge_dify_templates_to_ppt_enhanced
                    status_text.text("🔗 正在整合PPT页面(增强格式保留)...")
                    progress_bar.progress(90)
                    merge_result = merge_dify_templates_to_ppt_enhanced(filled_page_results)
                    
                    # 整合PPT结果处理保持不变
                    
                    if merge_result["success"]:
                        # 保存整合结果
                        st.session_state.ppt_merge_result = merge_result
                        
                        # 完成处理流程
                        progress_bar.progress(100)
                        status_text.text("✅ PPT整合完成，可以下载！")
                        
                        # 清除进度显示
                        progress_bar.empty()
                        status_text.empty()
                        
                        # 显示完成提示，触发页面更新
                        st.success("🎉 PPT生成完成！请查看下载按钮。")
                        
                        # PPT整合完成，结果已保存到session_state，刷新页面显示结果
                        st.rerun()
                    else:
                        progress_bar.empty()
                        status_text.empty()
                        st.error(f"❌ PPT整合失败: {merge_result.get('error', '未知错误')}")
                        
                        if merge_result.get("errors"):
                            with st.expander("🔍 查看详细错误信息", expanded=False):
                                for error in merge_result["errors"]:
                                    st.error(f"• {error}")
                        pass  # 显示错误但不退出，继续显示其他功能
                
                except ImportError:
                    progress_bar.empty()
                    status_text.empty()
                    st.error("❌ PPT整合模块未找到，请检查 ppt_merger.py 文件")
                    pass  # 显示错误但不退出，继续显示其他功能
                except Exception as e:
                    progress_bar.empty()
                    status_text.empty()
                    st.error(f"❌ PPT整合过程中出现异常: {str(e)}")
                    pass  # 显示错误但不退出，继续显示其他功能
                
            except ImportError as e:
                st.error(f"❌ 模块导入失败: {str(e)}")
            except Exception as e:
                progress_bar.empty()
                status_text.empty()
                st.error(f"❌ 智能PPT生成过程中出现异常: {str(e)}")
                logger.error("智能PPT生成异常: %s", str(e))
    
    # 开发者专用功能：自定义模板测试
    if user_role == "开发者":
        with tab3:
            # 自定义模板测试功能
            st.markdown("### 🧪 自定义模板测试")
            st.markdown('<div class="info-box">🎯 <strong>功能说明</strong><br>此功能独立于智能分页和Dify API，专门用于测试您自己的PPT模板。您可以上传自定义模板，输入文本内容，系统将智能填充到您的模板中。</div>', unsafe_allow_html=True)
            
            # 模板上传区域
            st.markdown("#### 📁 上传您的PPT模板")
            
            uploaded_files = st.file_uploader(
                "选择您的PPT模板文件（可选择多个单页PPT）",
                type=['pptx'],
                help="请上传.pptx格式的PPT模板文件，支持同时上传多个单页PPT文件",
                accept_multiple_files=True,
                key="custom_template_uploader"
            )
            
            if uploaded_files:
                st.success(f"✅ 已上传 {len(uploaded_files)} 个PPT模板文件")
                
                # 处理并存储所有文件的信息
                processed_files = []
                import tempfile
                import re
                from pptx import Presentation
                
                # 分析每个上传的文件
                for file_idx, uploaded_file in enumerate(uploaded_files):
                    try:
                        # 创建临时文件
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                            tmp_file.write(uploaded_file.getvalue())
                            temp_ppt_path = tmp_file.name
                        
                        # 验证PPT文件
                        is_valid, error_msg = FileManager.validate_ppt_file(temp_ppt_path)
                        
                        if is_valid:
                            # 分析模板结构
                            temp_presentation = Presentation(temp_ppt_path)
                            slide_count = len(temp_presentation.slides)
                        
                            # 分析占位符
                            total_placeholders = 0
                            placeholder_info = []
                        
                            for i, slide in enumerate(temp_presentation.slides):
                                slide_placeholders = []
                                table_placeholders = []
                            
                                for shape in slide.shapes:
                                    # 处理普通文本框中的占位符
                                    if hasattr(shape, 'text') and shape.text:
                                        placeholders = re.findall(r'\{([^}]+)\}', shape.text)
                                        if placeholders:
                                            slide_placeholders.extend(placeholders)
                                            total_placeholders += len(placeholders)
                                
                                    # 处理表格中的占位符
                                    elif hasattr(shape, 'shape_type') and shape.shape_type == 19:
                                        table = shape.table
                                        for row_idx, row in enumerate(table.rows):
                                            for col_idx, cell in enumerate(row.cells):
                                                cell_text = cell.text.strip()
                                                if cell_text:
                                                    placeholders = re.findall(r'\{([^}]+)\}', cell_text)
                                                    if placeholders:
                                                        for placeholder in placeholders:
                                                            table_placeholders.append(f"{placeholder}(表格{row_idx+1},{col_idx+1})")
                                                            total_placeholders += 1
                            
                                # 合并占位符
                                all_slide_placeholders = slide_placeholders + table_placeholders
                                if all_slide_placeholders:
                                    placeholder_info.append({
                                        'slide_num': i + 1,
                                        'placeholders': slide_placeholders,
                                        'table_placeholders': table_placeholders,
                                        'total_count': len(all_slide_placeholders)
                                    })
                        
                            # 存储文件信息
                            processed_files.append({
                                'index': file_idx,
                                'name': uploaded_file.name,
                                'size': f"{uploaded_file.size / 1024:.1f} KB",
                                'temp_path': temp_ppt_path,
                                'slide_count': slide_count,
                                'placeholder_count': total_placeholders,
                                'placeholder_info': placeholder_info,
                                'is_valid': True
                            })
                        else:
                            processed_files.append({
                                'index': file_idx,
                                'name': uploaded_file.name,
                                'size': f"{uploaded_file.size / 1024:.1f} KB",
                                'temp_path': None,
                                'error': error_msg,
                                'is_valid': False
                            })
                            
                    except Exception as e:
                        processed_files.append({
                            'index': file_idx,
                            'name': uploaded_file.name,
                            'size': f"{uploaded_file.size / 1024:.1f} KB",
                            'temp_path': None,
                            'error': str(e),
                            'is_valid': False
                        })
                
                # 显示文件分析结果
                st.markdown("#### 📋 文件分析结果")
                
                valid_files = [f for f in processed_files if f['is_valid']]
                invalid_files = [f for f in processed_files if not f['is_valid']]
                
                if valid_files:
                    with st.expander(f"✅ 有效文件 ({len(valid_files)})", expanded=True):
                        for file_info in valid_files:
                            col1, col2, col3, col4 = st.columns([2, 1, 1, 1])
                            with col1:
                                st.text(f"📄 {file_info['name']}")
                            with col2:
                                st.text(f"📑 {file_info['slide_count']}页")
                            with col3:
                                st.text(f"🎯 {file_info['placeholder_count']}占位符")
                            with col4:
                                st.text(f"📊 {file_info['size']}")
                
                if invalid_files:
                    with st.expander(f"❌ 无效文件 ({len(invalid_files)})", expanded=False):
                        for file_info in invalid_files:
                            st.error(f"📄 {file_info['name']}: {file_info.get('error', '未知错误')}")
                
                # 为每个有效文件提供独立的测试界面
                if valid_files:
                    st.markdown("---")
                    st.markdown("#### 📝 为每个模板输入测试内容")
                    
                    
                    # 为每个有效文件创建独立的测试界面
                    for file_info in valid_files:
                        st.markdown(f"##### 📄 {file_info['name']}")
                        
                        col1, col2, col3 = st.columns([2, 1, 1])
                        with col1:
                            st.info(f"🎯 发现 {file_info['placeholder_count']} 个占位符")
                        with col2:
                            st.text(f"📑 {file_info['slide_count']} 页")
                        with col3:
                            st.text(f"📊 {file_info['size']}")
                        
                        # 显示占位符详情
                        if file_info['placeholder_info']:
                            with st.expander(f"🔍 查看 {file_info['name']} 的占位符", expanded=False):
                                for info in file_info['placeholder_info'][:3]:  # 只显示前3页
                                    slide_num = info['slide_num']
                                    text_placeholders = info['placeholders']
                                    table_placeholders = info['table_placeholders']
                                    
                                    st.write(f"**第{slide_num}页（{info['total_count']}个占位符）：**")
                                    if text_placeholders:
                                        st.write(f"  📝 文本框：{', '.join([f'{{{p}}}' for p in text_placeholders])}")
                                    if table_placeholders:
                                        st.write(f"  📊 表格：{', '.join([f'{{{p}}}' for p in table_placeholders])}")
                                
                                if len(file_info['placeholder_info']) > 3:
                                    remaining = len(file_info['placeholder_info']) - 3
                                    st.write(f"... 还有 {remaining} 页包含占位符")
                        
                        # 文本输入
                        test_text = st.text_area(
                            f"为 {file_info['name']} 输入测试内容：",
                            height=150,
                            placeholder=f"""例如：

针对 {file_info['name']} 的测试内容

AI将分析您的文本结构，并智能地将内容分配到该模板的 {file_info['placeholder_count']} 个占位符位置。

主要特点：
- 保持原有模板设计风格
- 智能内容分配
- 支持文本框和表格占位符

这个测试将展示AI如何理解您的内容并填充到该模板的对应位置。""",
                            help=f"AI将分析您的文本并智能分配到该模板的所有占位符中",
                            key=f"test_text_{file_info['index']}"
                        )
                        
                        # 测试按钮
                        
                        st.markdown("---")  # 分隔线，用于分隔不同文件的测试区域
                
                    # 统一测试按钮和处理逻辑
                    st.markdown("#### 🚀 开始批量测试")
                    
                    # 检查是否所有文件都有文本输入
                    all_texts_filled = True
                    text_inputs = {}
                    visual_options = {}
                    
                    for file_info in valid_files:
                        # 获取每个文件的文本输入
                        text_key = f"test_text_{file_info['index']}"
                        visual_key = f"visual_opt_{file_info['index']}"
                        
                        # 检查session_state中是否有这些值
                        if text_key in st.session_state:
                            text_inputs[file_info['index']] = st.session_state[text_key]
                            if not st.session_state[text_key].strip():
                                all_texts_filled = False
                        else:
                            all_texts_filled = False
                        
                        if visual_key in st.session_state:
                            visual_options[file_info['index']] = st.session_state[visual_key]
                        else:
                            visual_options[file_info['index']] = False
                    
                    # 统一测试按钮
                    batch_test_button = st.button(
                        f"🧪 批量测试所有模板 ({len(valid_files)}个)",
                        type="primary",
                        disabled=not all_texts_filled,
                        help="同时测试所有模板并合并为一个PPT文件",
                        use_container_width=True,
                        key="batch_test_btn"
                    )
                    
                    # 处理批量测试
                    if batch_test_button and all_texts_filled:
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        
                        try:
                            # 存储处理后的模板文件路径，用于Spire合并
                            processed_template_paths = []
                            processed_files = []
                            total_files = len(valid_files)
                            
                            for idx, file_info in enumerate(valid_files):
                                # 更新进度
                                progress = int((idx / total_files) * 70) + 10
                                progress_bar.progress(progress)
                                status_text.text(f"🔧 正在处理 {file_info['name']} ({idx + 1}/{total_files})")
                                
                                # 获取对应的文本输入
                                test_text = text_inputs.get(file_info['index'], '')
                                
                                if not test_text.strip():
                                    continue
                                
                                try:
                                    # 创建模板生成器
                                    custom_generator = UserPPTGenerator(api_key)
                                    success, message = custom_generator.load_ppt_from_path(file_info['temp_path'])
                                    
                                    if not success:
                                        st.error(f"❌ 模板 {file_info['name']} 加载失败: {message}")
                                        continue
                                    
                                    # AI分析和填充
                                    assignments = custom_generator.process_text_with_openai(test_text)
                                    success, results = custom_generator.apply_text_assignments(assignments, test_text)
                                    
                                    if not success:
                                        st.error(f"❌ {file_info['name']} 内容填充失败")
                                        continue
                                    
                                    # 清理占位符
                                    cleanup_results = custom_generator.cleanup_unfilled_placeholders()
                                    
                                        # 应用基础美化
                                    optimization_results = custom_generator.apply_basic_beautification()
                                    
                                    # 保存处理后的PPT到临时文件
                                    import tempfile
                                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                                    temp_dir = tempfile.gettempdir()
                                    processed_filename = f"processed_{idx}_{timestamp}.pptx"
                                    processed_path = os.path.join(temp_dir, processed_filename)
                                    
                                    # 保存处理后的PPT
                                    custom_generator.presentation.save(processed_path)
                                    processed_template_paths.append({
                                        'template_path': processed_path,
                                        'page_number': idx + 1
                                    })
                                    
                                    # 记录处理成功的文件
                                    processed_files.append({
                                        'name': file_info['name'],
                                        'success': True,
                                        'cleanup_count': cleanup_results.get('cleaned_placeholders', 0) if cleanup_results else 0,
                                        'processed_path': processed_path
                                    })
                                    
                                except Exception as e:
                                    st.error(f"❌ 处理 {file_info['name']} 时出现错误: {str(e)}")
                                    processed_files.append({
                                        'name': file_info['name'],
                                        'success': False,
                                        'error': str(e)
                                    })
                        
                            # 使用Spire合并所有处理后的PPT文件
                            if processed_template_paths:
                                status_text.text("📦 正在使用Spire合并文件（保持格式）...")
                                progress_bar.progress(85)
                                
                                try:
                                    from ppt_merger_spire import merge_dify_templates_to_ppt_spire
                                    merge_result = merge_dify_templates_to_ppt_spire(processed_template_paths)
                                    
                                    if merge_result.get('success') and merge_result.get('presentation_bytes'):
                                        merged_ppt_bytes = merge_result['presentation_bytes']
                                    else:
                                        # Spire合并失败，回退到简单合并
                                        st.warning("⚠️ Spire合并失败，回退到基本合并模式")
                                        from pptx import Presentation
                                        if processed_template_paths:
                                            first_ppt = Presentation(processed_template_paths[0]['template_path'])
                                            import io
                                            temp_bytes = io.BytesIO()
                                            first_ppt.save(temp_bytes)
                                            merged_ppt_bytes = temp_bytes.getvalue()
                                        else:
                                            raise Exception("没有可合并的文件")
                                
                                except ImportError:
                                    st.warning("⚠️ Spire.Presentation未安装，使用基本合并模式")
                                    # 回退到基本合并
                                    from pptx import Presentation
                                    if processed_template_paths:
                                        first_ppt = Presentation(processed_template_paths[0]['template_path'])
                                        import io
                                        temp_bytes = io.BytesIO()
                                        first_ppt.save(temp_bytes)
                                        merged_ppt_bytes = temp_bytes.getvalue()
                                    else:
                                        raise Exception("没有可合并的文件")
                                        
                                except Exception as e:
                                    st.error(f"❌ 合并失败: {str(e)}")
                                    raise e
                                    
                                # 清理临时处理文件
                                for file_info in processed_files:
                                    if file_info.get('success') and file_info.get('processed_path'):
                                        try:
                                            os.unlink(file_info['processed_path'])
                                        except:
                                            pass
                            else:
                                raise Exception("没有成功处理的文件可供合并")
                        
                            # 清除进度显示
                            progress_bar.empty()
                            status_text.empty()
                            
                            # 显示成功信息
                            st.markdown('<div class="success-box">🎉 批量模板测试完成！</div>', unsafe_allow_html=True)
                            
                            # 显示处理摘要
                            st.markdown("### 📊 处理结果")
                            
                            successful_files = [f for f in processed_files if f['success']]
                            failed_files = [f for f in processed_files if not f['success']]
                            
                            col1, col2, col3 = st.columns(3)
                            
                            with col1:
                                st.metric("✅ 成功处理", len(successful_files))
                            
                            with col2:
                                st.metric("❌ 处理失败", len(failed_files))
                            
                            with col3:
                                total_cleanup = sum(f.get('cleanup_count', 0) for f in successful_files)
                                st.metric("🧹 清理占位符", total_cleanup)
                            
                            # 显示详细结果
                            if successful_files:
                                with st.expander("✅ 成功处理的文件", expanded=True):
                                    for file_info in successful_files:
                                        st.success(f"📄 {file_info['name']} - 清理了{file_info['cleanup_count']}个占位符")
                            
                            if failed_files:
                                with st.expander("❌ 处理失败的文件", expanded=False):
                                    for file_info in failed_files:
                                        st.error(f"📄 {file_info['name']}: {file_info.get('error', '未知错误')}")
                            
                            # 下载合并后的文件
                            if successful_files:
                                st.markdown("### 💾 下载合并结果")
                                
                                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                                filename = f"批量测试结果_{timestamp}.pptx"
                                
                                col1, col2, col3 = st.columns([1, 2, 1])
                                with col2:
                                    st.download_button(
                                        label=f"📥 下载合并后的PPT文件",
                                        data=merged_ppt_bytes,
                                        file_name=filename,
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                        use_container_width=True,
                                        key="download_merged_result"
                                    )
                                
                                st.info(f"📁 **文件名：** {filename}")
                                st.info(f"📑 **包含：** {len(successful_files)} 个模板的测试结果")
                            
                            # 清理所有临时文件
                            for file_info in valid_files:
                                try:
                                    import os
                                    os.unlink(file_info['temp_path'])
                                except:
                                    pass
                                
                        except Exception as e:
                            progress_bar.empty()
                            status_text.empty()
                            st.error(f"❌ 批量处理过程中出现错误: {str(e)}")
                            
                            # 清理临时文件
                            for file_info in valid_files:
                                try:
                                    import os
                                    os.unlink(file_info['temp_path'])
                                except:
                                    pass
                
            else:
                # 未上传文件时的说明
                st.markdown("#### 🎯 使用说明")
                
                col1, col2 = st.columns(2)
                with col1:
                    st.markdown("""
                    **📋 模板要求：**
                    - 文件格式：.pptx
                    - 文件大小：<50MB
                    - 包含占位符：{标题}、{内容}等
                    - 支持多个单页PPT文件
                    """)
                
                with col2:
                    st.markdown("""
                    **🔄 处理流程：**
                    1. 同时上传多个PPT模板文件
                    2. 系统验证和分析每个模板
                    3. 为每个模板输入测试内容
                    4. AI智能分配内容到占位符
                    5. 分别下载每个测试结果
                """)
            
                st.markdown("#### ✨ 功能特色")
            
                col1, col2, col3 = st.columns(3)
                with col1:
                    st.markdown("""
                **🎨 保持设计风格**
                - 完全保留您的模板样式
                - 不改变颜色、字体、布局
                - 只填充内容到指定位置
                """)
            
                with col2:
                    st.markdown("""
                    **🤖 智能内容分配**
                    - AI理解文本结构和含义
                    - 自动匹配最合适的占位符
                    - 支持多种内容类型处理
                    """)
            
                with col3:
                    st.markdown("""
                    **📊 多文件批量测试**
                    - 支持同时上传多个模板
                    - 每个模板独立文本输入
                    - 单独生成测试结果
                    """)
            
                st.markdown('<div class="warning-box">💡 <strong>提示：</strong> 现在支持同时上传多个单页PPT模板！您可以为每个模板输入不同的测试内容，系统会分别处理并生成独立的测试结果。请确保每个模板都包含形如 {标题}、{内容}、{要点} 等占位符。</div>', unsafe_allow_html=True)
    
    # 开发者专用功能：表格文本填充
    if user_role == "开发者":
        with tab_table:
            # 数字智能提取填充功能
            st.markdown("### 📊 数字智能提取填充")
            
            st.markdown('<div class="info-box">🎯 <strong>功能说明</strong><br>专门用于处理包含数字信息的文本填充。AI会特别关注并提取所有数字（价格、百分比、尺寸、日期等），将数字单独填充到对应的占位符中，而不是将包含数字的整段文本都填入{content}等通用占位符。</div>', unsafe_allow_html=True)
            
            # 模板上传区域
            st.markdown("#### 📁 上传您的PPT模板")
            
            table_uploaded_files = st.file_uploader(
                "选择您的PPT模板文件（可选择多个单页PPT）",
                type=['pptx'],
                help="请上传.pptx格式的PPT模板文件，支持同时上传多个单页PPT文件",
                accept_multiple_files=True,
                key="table_template_uploader"
            )
            
            if table_uploaded_files:
                st.success(f"✅ 已上传 {len(table_uploaded_files)} 个PPT模板文件")
                
                # 处理并存储所有文件的信息
                processed_files = []
                import tempfile
                import re
                from pptx import Presentation
                
                # 为每个文件创建临时文件并验证
                for idx, uploaded_file in enumerate(table_uploaded_files):
                    try:
                        # 创建临时文件
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                            tmp_file.write(uploaded_file.getvalue())
                            temp_path = tmp_file.name
                        
                        # 验证PPT文件
                        is_valid, error_msg = FileManager.validate_ppt_file(temp_path)
                        
                        if is_valid:
                            # 分析模板结构
                            presentation = Presentation(temp_path)
                            slide_count = len(presentation.slides)
                            
                            # 分析占位符 - 支持文本框和表格中的占位符
                            total_placeholders = 0
                            placeholder_info = []
                            
                            for i, slide in enumerate(presentation.slides):
                                slide_placeholders = []
                                table_placeholders = []
                                
                                for shape in slide.shapes:
                                    # 处理普通文本框中的占位符
                                    if hasattr(shape, 'text') and shape.text:
                                        placeholders = re.findall(r'\{([^}]+)\}', shape.text)
                                        if placeholders:
                                            slide_placeholders.extend(placeholders)
                                            total_placeholders += len(placeholders)
                                    
                                    # 处理表格中的占位符
                                    elif hasattr(shape, 'shape_type') and shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
                                        table = shape.table
                                        for row_idx, row in enumerate(table.rows):
                                            for col_idx, cell in enumerate(row.cells):
                                                cell_text = cell.text.strip()
                                                if cell_text:
                                                    placeholders = re.findall(r'\{([^}]+)\}', cell_text)
                                                    if placeholders:
                                                        for placeholder in placeholders:
                                                            table_placeholders.append(f"{placeholder}(表格{row_idx+1},{col_idx+1})")
                                                            total_placeholders += 1
                                
                                # 合并文本框和表格占位符
                                all_slide_placeholders = slide_placeholders + table_placeholders
                                if all_slide_placeholders:
                                    placeholder_info.append({
                                        'slide_num': i + 1,
                                        'placeholders': slide_placeholders,
                                        'table_placeholders': table_placeholders,
                                        'total_count': len(all_slide_placeholders)
                                    })
                            
                            processed_files.append({
                                'index': idx,
                                'filename': uploaded_file.name,
                                'temp_path': temp_path,
                                'slide_count': slide_count,
                                'placeholder_count': total_placeholders,
                                'placeholder_info': placeholder_info,
                                'is_valid': True,
                                'error': None
                            })
                        else:
                            processed_files.append({
                                'index': idx,
                                'filename': uploaded_file.name,
                                'temp_path': None,
                                'is_valid': False,
                                'error': error_msg
                            })
                            
                    except Exception as e:
                        processed_files.append({
                            'index': idx,
                            'filename': uploaded_file.name,
                            'temp_path': None,
                            'is_valid': False,
                            'error': str(e)
                        })
                
                # 显示文件验证结果
                valid_files = [f for f in processed_files if f['is_valid']]
                invalid_files = [f for f in processed_files if not f['is_valid']]
                
                col1, col2 = st.columns([1, 2])
                with col1:
                    st.metric("✅ 有效文件", len(valid_files))
                    st.metric("❌ 无效文件", len(invalid_files))
                    
                with col2:
                    if invalid_files:
                        with st.expander("❌ 文件验证失败", expanded=True):
                            for file_info in invalid_files:
                                st.error(f"**{file_info['filename']}**: {file_info['error']}")
                
                if valid_files:
                    with st.expander("✅ 文件结构分析", expanded=False):
                        for file_info in valid_files:
                            st.write(f"**{file_info['filename']}**")
                            st.write(f"  📑 幻灯片: {file_info['slide_count']} 页")
                            st.write(f"  🎯 占位符: {file_info['placeholder_count']} 个")
                            
                            # 显示占位符详情
                            if file_info['placeholder_info']:
                                for info in file_info['placeholder_info'][:3]:  # 显示前3页
                                    slide_num = info['slide_num']
                                    text_placeholders = info['placeholders']
                                    table_placeholders = info['table_placeholders']
                                    
                                    st.write(f"    第{slide_num}页（{info['total_count']}个占位符）：")
                                    
                                    if text_placeholders:
                                        st.write(f"      📝 文本框：{', '.join([f'{{{p}}}' for p in text_placeholders])}")
                                    
                                    if table_placeholders:
                                        st.write(f"      📊 表格：{', '.join([f'{{{p}}}' for p in table_placeholders])}")
                                
                                if len(file_info['placeholder_info']) > 3:
                                    remaining = len(file_info['placeholder_info']) - 3
                                    st.write(f"    ... 还有 {remaining} 页")
                
                # 如果有有效文件，显示文本输入区域
                if valid_files:
                    st.markdown("---")
                    st.markdown("#### 📝 为每个文件输入文本")
                
                    # 为每个有效文件创建文本输入框
                    text_inputs = {}
                    for file_info in valid_files:
                        st.markdown(f"**{file_info['filename']}** (📑 {file_info['slide_count']} 页, 🎯 {file_info['placeholder_count']} 个占位符)")
                        
                        text_key = f"table_text_{file_info['index']}"
                        text_inputs[file_info['index']] = st.text_area(
                            f"为 {file_info['filename']} 输入要填充的文本内容：",
                            height=120,
                            placeholder="""例如（产品信息）：
            iPhone 15 Pro
            价格：999美元
            屏幕尺寸：6.1英寸
            处理器：A17 Pro芯片

            AI将自动提取数字信息并分别填入对应的占位符""",
                            help="AI将智能提取数字信息并分别填充，文本描述和数字数据会分开处理",
                            key=text_key
                        )
                        st.markdown("---")
                
                    # 处理选项（保留结构以便未来扩展）
                    col1, col2 = st.columns(2)
                    with col1:
                        pass  # 留空
                    
                    with col2:
                        pass  # 留空
                
                    # 处理按钮
                    st.markdown("#### 🚀 统一批量处理")
                    
                    # 检查是否所有文件都有文本输入
                    has_all_text = all(text_inputs.get(file_info['index'], '').strip() for file_info in valid_files)
                    
                    table_batch_button = st.button(
                        "📊 批量智能数字填充并合并",
                        type="primary",
                        use_container_width=True,
                        disabled=not has_all_text,
                        help="对所有文件批量处理，AI将提取数字信息并分别填充，然后合并为一个PPT",
                        key="table_batch_btn"
                    )
                
                    # 批量处理逻辑
                    if table_batch_button and has_all_text:
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        
                        try:
                            status_text.text("🚀 开始批量处理...")
                            progress_bar.progress(10)
                            
                            # 存储处理结果
                            processed_results = []
                            
                            total_files = len(valid_files)
                            for idx, file_info in enumerate(valid_files):
                                file_progress = 20 + (idx * 50) // total_files
                                status_text.text(f"📝 正在处理 {file_info['filename']} ({idx+1}/{total_files})...")
                                progress_bar.progress(file_progress)
                                
                                try:
                                    # 创建生成器
                                    generator = UserPPTGenerator(api_key)
                                    success, message = generator.load_ppt_from_path(file_info['temp_path'])
                                    
                                    if not success:
                                        st.warning(f"⚠️ {file_info['filename']} 加载失败: {message}")
                                        continue
                                    
                                    # 获取对应的文本输入
                                    text_content = text_inputs.get(file_info['index'], '').strip()
                                    if not text_content:
                                        continue
                                    
                                    # AI分析
                                    assignments = generator.process_text_with_openai_enhanced(text_content)
                                    
                                    if assignments.get('error'):
                                        st.warning(f"⚠️ {file_info['filename']} AI分析失败: {assignments['error']}")
                                        continue
                                    
                                    # 填充内容
                                    success, results = generator.apply_text_assignments(assignments, text_content)
                                    
                                    if not success:
                                        st.warning(f"⚠️ {file_info['filename']} 内容填充失败")
                                        continue
                                    
                                    # 清理占位符
                                    cleanup_results = generator.cleanup_unfilled_placeholders()
                                    
                                    # 应用基础美化
                                    optimization_results = generator.apply_basic_beautification()
                                    
                                    # 保存处理结果到临时文件（使用与自定义模板测试相同的方法）
                                    import tempfile
                                    import os
                                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                                    temp_dir = tempfile.gettempdir()
                                    processed_filename = f"table_processed_{idx}_{timestamp}.pptx"
                                    temp_result_path = os.path.join(temp_dir, processed_filename)
                                    
                                    # 使用presentation.save()方法保存文件
                                    generator.presentation.save(temp_result_path)
                                    
                                    processed_results.append({
                                        'page_number': idx + 1,
                                        'template_number': idx + 1,
                                        'template_path': temp_result_path,
                                        'template_filename': file_info['filename'],
                                        'original_filename': file_info['filename'],
                                        'success': True,
                                        'cleanup_count': cleanup_results.get('cleaned_placeholders', 0) if cleanup_results else 0
                                    })
                                    
                                except Exception as e:
                                    st.warning(f"⚠️ {file_info['filename']} 处理失败: {str(e)}")
                                    processed_results.append({
                                        'page_number': idx + 1,
                                        'template_path': file_info['temp_path'],
                                        'template_filename': file_info['filename'],
                                        'success': False,
                                        'error': str(e)
                                    })
                            
                            # 合并结果
                            successful_results = [r for r in processed_results if r.get('success', False)]
                            
                            if not successful_results:
                                st.error("❌ 没有成功处理的文件，无法合并")
                                return
                            
                            status_text.text(f"🔍 正在合并 {len(successful_results)} 个处理结果...")
                            progress_bar.progress(80)
                            
                            # 使用Spire合并器进行格式保持合并
                            from ppt_merger_spire import merge_dify_templates_to_ppt_spire
                            
                            merge_result = merge_dify_templates_to_ppt_spire(successful_results)
                            
                            # 完成处理
                            status_text.text("📦 正在准备下载...")
                            progress_bar.progress(100)
                            
                            # 清除进度显示
                            progress_bar.empty()
                            status_text.empty()
                            
                            if merge_result.get('success'):
                                st.markdown('<div class="success-box">🎉 批量表格数字填充完成！</div>', unsafe_allow_html=True)
                                
                                # 显示处理结果
                                st.markdown("### 📊 批量处理结果")
                                
                                col1, col2, col3, col4 = st.columns(4)
                                
                                with col1:
                                    st.metric("📑 最终页数", merge_result.get('total_pages', 0))
                                
                                with col2:
                                    successful_count = len(successful_results)
                                    st.metric("✅ 成功文件", successful_count)
                                
                                with col3:
                                    failed_count = len(processed_results) - successful_count
                                    st.metric("❌ 失败文件", failed_count)
                                
                                with col4:
                                    total_cleanup = sum(r.get('cleanup_count', 0) for r in successful_results)
                                    st.metric("🧹 清理占位符", total_cleanup)
                                
                                # 显示详细结果
                                if len(processed_results) > successful_count:
                                    failed_results = [r for r in processed_results if not r.get('success', False)]
                                    with st.expander("⚠️ 处理失败的文件", expanded=False):
                                        for result in failed_results:
                                            st.error(f"**{result['template_filename']}**: {result.get('error', '未知错误')}")
                                
                                # 下载文件
                                st.markdown("### 💾 下载合并结果")
                                
                                try:
                                    if merge_result.get('presentation_bytes'):
                                        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                                        filename = f"表格数字填充批量处理结果_{timestamp}.pptx"
                                        
                                        col1, col2, col3 = st.columns([1, 2, 1])
                                        with col2:
                                            st.download_button(
                                                label="📥 下载合并后的PPT",
                                                data=merge_result['presentation_bytes'],
                                                file_name=filename,
                                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                                use_container_width=True,
                                                key="download_table_batch_result"
                                            )
                                    else:
                                        st.error("❌ 未能获取合并后的PPT数据")
                                        
                                except Exception as e:
                                    st.error(f"❌ 生成下载文件失败: {str(e)}")
                            
                            else:
                                st.error(f"❌ 批量处理过程中出现错误: {merge_result.get('error', '未知错误')}")
                            
                            # 清理临时文件
                            for result in processed_results:
                                if result.get('success') and 'temp_result_path' in locals():
                                    try:
                                        import os
                                        if os.path.exists(result['template_path']):
                                            os.unlink(result['template_path'])
                                    except:
                                        pass
                            
                        except Exception as e:
                            progress_bar.empty()
                            status_text.empty()
                            st.error(f"❌ 批量处理过程中出现异常: {str(e)}")
                            logger.error("批量表格填充异常: %s", str(e))
            
            else:
                st.markdown("### 📖 功能特点")
                st.markdown("""
                **🔢 数字智能处理**
                - 自动提取文本中的所有数字信息
                - 将数字和文本分开填充到对应占位符
                - 支持价格、百分比、尺寸、日期等多种数据类型
                - 避免将包含数字的整段文本填入通用占位符
                
                **🎯 精确匹配**
                - 根据占位符名称智能匹配数据类型
                - {价格} 填入货币数字，{描述} 填入文本描述
                - 数据和内容完全分离，提高填充精度
                """)
                
                st.markdown('<div class="warning-box">💡 <strong>提示：</strong> 推荐使用具体的占位符名称，如 {产品名称}、{价格}、{百分比}、{尺寸}、{数量}、{日期} 等。AI将根据占位符名称智能提取对应的数字或文本信息。避免使用{content}这样的通用占位符来包含数字数据。</div>', unsafe_allow_html=True)
    
    # 开发者专用功能：PPT格式读取展示
    if user_role == "开发者":
        with tab_format:
            # PPT格式读取展示功能
            st.markdown("### 🔍 PPT格式读取展示")
            st.markdown("**上传一个PPT文件，查看我们的格式读取功能能识别到什么信息**")
            
            col1, col2 = st.columns([1, 2])
            
            with col1:
                st.markdown("#### 📤 上传文件")
                uploaded_file = st.file_uploader(
                    "选择PPT文件",
                    type=['pptx'],
                    help="支持.pptx格式的PowerPoint文件"
                )
                
                if uploaded_file is not None:
                    st.success(f"✅ 已上传：{uploaded_file.name}")
                    
                    # 分析按钮
                    if st.button("🔍 开始分析格式", type="primary"):
                        with st.spinner("正在分析PPT格式..."):
                            try:
                                # 保存上传的文件到临时位置
                                import tempfile
                                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                                    temp_file.write(uploaded_file.getbuffer())
                                    temp_path = temp_file.name
                                
                                # 使用现有的PPT分析功能
                                from pptx import Presentation as PptxPresentation
                                presentation = PptxPresentation(temp_path)
                                ppt_structure = PPTAnalyzer.analyze_ppt_structure(presentation)
                                
                                # 将结果存储到session state，包括临时文件路径用于后续格式提取
                                st.session_state.format_analysis_result = {
                                    'filename': uploaded_file.name,
                                    'structure': ppt_structure,
                                    'temp_path': temp_path,
                                    'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                                }
                                
                                # 注意：临时文件暂时保留，用于后续格式提取
                                # 文件会在清除结果或会话结束时清理
                                    
                                st.success("🎉 分析完成！")
                                st.rerun()
                                
                            except Exception as e:
                                st.error(f"❌ 分析失败：{e}")
                                # 清理临时文件
                                try:
                                    os.remove(temp_path)
                                except:
                                    pass
        
            with col2:
                st.markdown("#### 📊 分析结果")
            
                if 'format_analysis_result' in st.session_state:
                    result = st.session_state.format_analysis_result
                    
                    st.markdown(f"**文件名：** {result['filename']}")
                    st.markdown(f"**分析时间：** {result['timestamp']}")
                    
                    structure = result['structure']
                    total_slides = structure.get('total_slides', 0)
                    total_placeholders = structure.get('total_placeholders', 0)
                    
                    # 基本统计
                    st.markdown("---")
                    st.markdown("### 📈 基本统计")
                
                    metric_cols = st.columns(3)
                    with metric_cols[0]:
                        st.metric("幻灯片数量", total_slides)
                    with metric_cols[1]:
                        st.metric("占位符总数", total_placeholders)
                    with metric_cols[2]:
                        all_placeholders = []
                        for slide in structure.get('slides', []):
                            all_placeholders.extend(slide.get('placeholders', {}).keys())
                        unique_placeholders = len(set(all_placeholders))
                        st.metric("不同占位符", unique_placeholders)
                
                    # 详细信息展开
                    st.markdown("---")
                    st.markdown("### 🔍 详细分析")
                
                    with st.expander("📋 占位符详情"):
                        for i, slide in enumerate(structure.get('slides', [])):
                            placeholders = slide.get('placeholders', {})
                            if placeholders:
                                st.markdown(f"**第 {i+1} 页：**")
                            
                                for placeholder_name, placeholder_info in placeholders.items():
                                    st.markdown(f"- **{{{placeholder_name}}}**")
                                    
                                    # 显示类型信息
                                    ph_type = placeholder_info.get('type', 'unknown')
                                    st.markdown(f"  - 类型：{ph_type}")
                                    
                                    # 显示原始文本
                                    original_text = placeholder_info.get('original_text', '')
                                    if original_text:
                                        st.markdown(f"  - 原始文本：`{original_text[:100]}{'...' if len(original_text) > 100 else ''}`")
                                    
                                    # 实时提取字体格式信息
                                    try:
                                        # 重新加载presentation来提取格式
                                        from pptx import Presentation as PptxPresentation
                                        temp_path = result.get('temp_path')
                                        if not temp_path or not os.path.exists(temp_path):
                                            st.markdown(f"  - 🎨 **格式：** ❌ 临时文件不存在")
                                            continue
                                            
                                        temp_presentation = PptxPresentation(temp_path)
                                        slide_obj = temp_presentation.slides[i]
                                        
                                        # 创建临时的PPTProcessor来提取格式
                                        from utils import PPTProcessor
                                        temp_processor = PPTProcessor(temp_presentation)
                                        
                                        # 获取容器对象
                                        container = placeholder_info.get('shape')
                                        if placeholder_info.get('type') == 'table_cell':
                                            container = placeholder_info.get('cell')
                                        
                                        # 如果没有容器信息，尝试重新查找
                                        if not container:
                                            # 在slide中查找包含这个占位符的shape
                                            placeholder_pattern = f"{{{placeholder_name}}}"
                                            for shape in slide_obj.shapes:
                                                if hasattr(shape, 'text') and placeholder_pattern in shape.text:
                                                    container = shape
                                                    break
                                                elif hasattr(shape, 'table'):
                                                    # 检查表格
                                                    for row in shape.table.rows:
                                                        for cell in row.cells:
                                                            if placeholder_pattern in cell.text:
                                                                container = cell
                                                                break
                                                        if container:
                                                            break
                                                    if container:
                                                        break
                                        
                                        # 如果找到容器，提取格式信息
                                        if container:
                                            format_info = temp_processor._extract_placeholder_format(container, placeholder_name)
                                            
                                            # 格式化显示
                                            font_details = []
                                            
                                            # 字体名称
                                            font_name = format_info.get('font_name')
                                            if font_name:
                                                font_details.append(f"字体: {font_name}")
                                            else:
                                                font_details.append("字体: None")
                                            
                                            # 字体大小
                                            font_size = format_info.get('font_size')
                                            if font_size:
                                                font_details.append(f"大小: {font_size}pt")
                                            else:
                                                font_details.append("大小: None")
                                            
                                            # 字体颜色
                                            font_color = format_info.get('font_color')
                                            if font_color:
                                                font_details.append(f"颜色: {font_color}")
                                            else:
                                                font_details.append("颜色: None")
                                            
                                            # 粗体和斜体
                                            style_details = []
                                            if format_info.get('font_bold'):
                                                style_details.append("粗体")
                                            if format_info.get('font_italic'):
                                                style_details.append("斜体")
                                            
                                            if style_details:
                                                font_details.append(f"样式: {', '.join(style_details)}")
                                            else:
                                                font_details.append("样式: 普通")
                                            
                                            st.markdown(f"  - 🎨 **格式：** {' | '.join(font_details)}")
                                            
                                            # 如果有问题的格式，用颜色标出
                                            problems = []
                                            if not font_name:
                                                problems.append("字体名称")
                                            if not font_size:
                                                problems.append("字体大小")
                                            if not font_color:
                                                problems.append("字体颜色")
                                        
                                            if problems:
                                                st.markdown(f"    ⚠️ *无法读取: {', '.join(problems)}*")
                                        else:
                                            st.markdown(f"  - 🎨 **格式：** ❌ 无法定位占位符容器")
                                            
                                    except Exception as format_error:
                                        st.markdown(f"  - 🎨 **格式：** ❌ 提取失败 ({str(format_error)[:50]})")
                                    
                                    st.markdown("")
                            else:
                                st.markdown(f"**第 {i+1} 页：** 无占位符")
                
                    with st.expander("🗂️ 原始结构数据"):
                        st.json(structure, expanded=False)
                    
                    # 清除结果按钮
                    if st.button("🗑️ 清除结果"):
                        # 清理临时文件
                        temp_path = result.get('temp_path')
                        if temp_path and os.path.exists(temp_path):
                            try:
                                os.remove(temp_path)
                            except:
                                pass
                        del st.session_state.format_analysis_result
                        st.rerun()
                    
                else:
                    st.markdown("👆 请先上传PPT文件并点击分析按钮")
                
                # 功能说明
                st.markdown("---")
                st.markdown("#### 💡 功能说明")
                st.markdown("""
                **这个工具会显示：**
                
                1. **📊 基本统计**：幻灯片数量、占位符总数等
                2. **🔍 占位符详情**：每个占位符的类型、位置、原始文本
                3. **🎨 格式信息**：字体名称、大小、颜色等（如果可读取）
                4. **📂 原始数据**：完整的结构分析结果
                
                **支持的格式：**
                - 文本框中的 `{占位符}`
                - 表格单元格中的 `{占位符}`
                - 多种字体格式识别
                
                **用途：**
                - 调试模板兼容性
                - 验证占位符识别准确性
                - 了解格式读取能力的边界
                """)

    with tab_watermark:
        # PPT去水印工具功能
        st.markdown("### 🧽 PPT去水印工具")
        
        # 去水印小提示框
        st.markdown('''<div class="info-box">💡 <strong>去水印小提示</strong><br>
        <strong>1. 专用范围：</strong>本功能专为处理本平台生成的Spire.Presentation水印设计<br>
        <strong>2. 安全无忧：</strong>系统会自动为您生成一个全新无水印的文件，您的原始文件不会被修改<br>
        <strong>3. 常见提示：</strong>下载后，若文件提示需要修复，这是正常现象，请别担心，简单点击"修复"即可正常使用<br>
        <strong>4. 最后检查：</strong>处理完成后，建议您检查一下内容是否完整
        </div>''', unsafe_allow_html=True)
        
        st.markdown("**上传含有Spire.Presentation水印的PPT文件，自动去除水印后提供下载**")
        
        col1, col2 = st.columns([1, 1])
        
        with col1:
            st.markdown("#### 📤 上传含水印的PPT")
            watermark_uploaded_file = st.file_uploader(
                "选择含水印的PPT文件",
                type=['pptx'],
                help="支持去除Spire.Presentation生成的红色文字和白色框水印",
                key="watermark_uploader"
            )
            
            if watermark_uploaded_file is not None:
                st.success(f"✅ 已上传：{watermark_uploaded_file.name}")
                
                # 显示文件信息
                file_size = len(watermark_uploaded_file.getvalue()) / 1024 / 1024
                st.info(f"📊 文件大小：{file_size:.2f} MB")
                
                # 去水印按钮
                if st.button("🧽 开始去除水印", type="primary", key="remove_watermark_btn"):
                    with st.spinner("正在去除水印，请稍候..."):
                        try:
                            # 保存上传的文件到临时位置
                            import tempfile
                            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_input:
                                temp_input.write(watermark_uploaded_file.getbuffer())
                                temp_input_path = temp_input.name
                            
                            # 创建输出文件路径
                            with tempfile.NamedTemporaryFile(suffix="_clean.pptx", delete=False) as temp_output:
                                temp_output_path = temp_output.name
                            
                            # 导入去水印模块
                            from watermark_remover import remove_spire_watermark
                            
                            # 执行去水印操作
                            result_path = remove_spire_watermark(temp_input_path, temp_output_path)
                            
                            # 读取处理后的文件
                            with open(result_path, 'rb') as f:
                                clean_file_data = f.read()
                            
                            # 计算清理后的文件大小
                            clean_file_size = len(clean_file_data) / 1024 / 1024
                            
                            # 将结果存储到session state
                            original_filename = watermark_uploaded_file.name
                            clean_filename = original_filename.replace('.pptx', '_无水印.pptx')
                            
                            st.session_state.watermark_removal_result = {
                                'original_filename': original_filename,
                                'clean_filename': clean_filename,
                                'clean_file_data': clean_file_data,
                                'original_size': file_size,
                                'clean_size': clean_file_size,
                                'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                            }
                            
                            # 清理临时文件
                            try:
                                os.remove(temp_input_path)
                                os.remove(temp_output_path)
                            except:
                                pass
                            
                            st.success("🎉 水印去除完成！")
                            st.rerun()
                            
                        except Exception as e:
                            st.error(f"❌ 去水印失败：{str(e)}")
                            # 清理临时文件
                            try:
                                if 'temp_input_path' in locals():
                                    os.remove(temp_input_path)
                                if 'temp_output_path' in locals():
                                    os.remove(temp_output_path)
                            except:
                                pass
        
        with col2:
            st.markdown("#### 📥 下载清理结果")
            
            if 'watermark_removal_result' in st.session_state:
                result = st.session_state.watermark_removal_result
                
                st.markdown(f"**原始文件：** {result['original_filename']}")
                st.markdown(f"**清理时间：** {result['timestamp']}")
                
                # 文件大小对比
                st.markdown("---")
                st.markdown("### 📊 处理结果")
                
                size_cols = st.columns(2)
                with size_cols[0]:
                    st.metric("原始大小", f"{result['original_size']:.2f} MB")
                with size_cols[1]:
                    st.metric("清理后大小", f"{result['clean_size']:.2f} MB")
                
                # 大小变化
                size_diff = result['clean_size'] - result['original_size']
                if size_diff < 0:
                    st.success(f"✅ 文件缩小了 {abs(size_diff):.2f} MB")
                elif size_diff > 0:
                    st.info(f"ℹ️ 文件增大了 {size_diff:.2f} MB")
                else:
                    st.info("ℹ️ 文件大小无变化")
                
                # 下载按钮
                st.markdown("---")
                st.download_button(
                    label="📥 下载无水印PPT",
                    data=result['clean_file_data'],
                    file_name=result['clean_filename'],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    type="primary"
                )
                
                # 清除结果按钮
                if st.button("🗑️ 清除结果", key="clear_watermark_result"):
                    del st.session_state.watermark_removal_result
                    st.rerun()
                    
            else:
                st.markdown("👆 请先上传PPT文件并点击去水印按钮")
                

    # 开发者专用功能：AI分页测试
    if user_role == "开发者":
        with tab_ai_test:
            # AI分页测试功能
            st.markdown("### 🤖 AI分页测试")
            st.markdown('<div class="info-box">🎯 <strong>功能说明</strong><br>此功能专门用于测试AI分页算法。系统默认使用两次调用策略：第一次AI专注内容逻辑分析，第二次AI优化页数控制（指定页数时精确调整，未指定时减少过度分页）。您可以观察完整的AI分页过程和结果对比。</div>', unsafe_allow_html=True)
            
            # 文本输入区域
            st.markdown("#### 📝 输入测试文本")
            test_text = st.text_area(
                "请输入您要测试的文本内容：",
                height=200,
                placeholder="请输入要测试AI分页的文本内容...",
                key="ai_paging_test_text"
            )
            
            # 页面数量选择
            st.markdown("#### ⚙️ 分页选项")
            col1, col2, col3 = st.columns([1, 1, 1])
            
            with col1:
                # 创建有效的页面数选项（与主功能保持一致）
                page_options = ["AI自动判断"] + [str(i) for i in range(4, 26)]
                selected_option = st.selectbox(
                    "目标页面数量",
                    options=page_options,
                    index=0,
                    help="选择AI自动判断或手动设置页面数量",
                    key="ai_test_page_count"
                )
                test_target_pages = 0 if selected_option == "AI自动判断" else int(selected_option)
            
            with col2:
                # 直接使用两次调用策略，不需要用户选择
                st.info("🔄 默认使用两次调用策略：第一次注重逻辑性，第二次优化页数")
            
            with col3:
                st.info("📋 测试说明：此功能将显示AI分页的完整过程和结果")
            
            # 测试按钮
            if st.button("🚀 开始AI分页测试", type="primary", key="start_ai_paging_test"):
                st.write("🔧 按钮被点击，开始执行测试...")
                if test_text.strip():
                    st.write(f"🔍 调试信息：文本长度 {len(test_text)} 字符，目标页数 {test_target_pages}，两次调用策略")
                    with st.spinner("正在调用AI分页算法..."):
                        try:
                            # 获取API配置
                            from config import get_config
                            config = get_config()
                            model_info = config.get_model_info()
                            
                            # 创建AI分页处理器（使用测试版本，不传递API密钥，让它从环境变量获取）
                            st.write("🔧 正在初始化AI分页处理器...")
                            from ai_page_splitter_test import AIPageSplitterTest
                            page_splitter = AIPageSplitterTest()
                            st.write("✅ AI分页处理器初始化成功")
                            
                            # 调用AI分页
                            target_page_count = int(test_target_pages) if test_target_pages > 0 else None
                            result = page_splitter.split_text_to_pages(test_text.strip(), target_page_count)
                            
                            # 显示结果
                            if result.get('success'):
                                # 检查是否使用了备用方案
                                if result.get('is_fallback'):
                                    st.warning("⚠️ AI分页失败，已使用备用分页方案")
                                else:
                                    st.success("✅ AI分页测试完成！")
                                
                                # 显示两次调用策略的信息
                                if result.get('is_two_pass_result'):
                                    first_pages = result.get('first_pass_pages')
                                    final_pages = result.get('final_pass_pages')
                                    if test_target_pages > 0:
                                        st.info(f"🔄 使用了两次调用策略（精确调整）：第一次生成 {first_pages} 页 → 第二次调整为 {final_pages} 页")
                                    else:
                                        if final_pages < first_pages:
                                            st.info(f"🔄 使用了两次调用策略（页数优化）：第一次生成 {first_pages} 页 → 第二次优化为 {final_pages} 页，减少了 {first_pages - final_pages} 页")
                                        else:
                                            st.info(f"🔄 使用了两次调用策略（页数优化）：第一次生成 {first_pages} 页 → 第二次保持 {final_pages} 页（已是合理分页）")
                                
                                # 显示分析结果
                                st.markdown("#### 📊 分页分析结果")
                                analysis = result.get('analysis', {})
                                
                                col1, col2, col3 = st.columns(3)
                                with col1:
                                    st.metric("总页面数", analysis.get('total_pages', 'N/A'))
                                with col2:
                                    st.metric("内容类型", analysis.get('content_type', 'N/A'))
                                with col3:
                                    st.metric("分割策略", analysis.get('split_strategy', 'N/A'))
                                
                                if analysis.get('reasoning'):
                                    st.markdown(f"**分页原因：** {analysis.get('reasoning')}")
                                
                                # 显示每页详情
                                st.markdown("#### 📄 页面详情")
                                pages = result.get('pages', [])
                                
                                for i, page in enumerate(pages):
                                    page_type = page.get('page_type', 'content')
                                    page_number = page.get('page_number', i + 1)
                                    
                                    # 根据页面类型设置标题
                                    if page_type == 'title':
                                        title = f"第{page_number}页 - 📋 封面页"
                                    elif page_type == 'table_of_contents':
                                        title = f"第{page_number}页 - 📑 目录页"
                                    elif page_type == 'ending':
                                        title = f"第{page_number}页 - 🔚 结尾页"
                                    else:
                                        title = f"第{page_number}页 - 📄 内容页"
                                    
                                    with st.expander(title, expanded=i < 2):
                                        # 显示页面基本信息
                                        st.markdown(f"**页面标题：** {page.get('title', '无')}")
                                        st.markdown(f"**页面类型：** {page_type}")
                                        if page.get('date'):
                                            st.markdown(f"**日期：** {page.get('date')}")
                                        
                                        # 显示原文内容
                                        original_text = page.get('original_text_segment', '')
                                        if original_text:
                                            st.markdown("**原文内容：**")
                                            st.text_area(
                                                "原文内容",
                                                value=original_text,
                                                height=min(max(len(original_text) // 10, 60), 200),
                                                disabled=True,
                                                key=f"page_content_{page_number}_{i}",
                                                label_visibility="collapsed"
                                            )
                                        else:
                                            st.markdown("**原文内容：** 无（使用固定模板）")
                                
                                # 显示完整的AI返回结果（JSON格式）
                                with st.expander("🔍 查看完整AI返回结果（JSON）", expanded=False):
                                    st.json(result)
                            else:
                                st.error("❌ AI分页测试失败，请检查输入内容和配置")
                                if 'error' in result:
                                    st.error(f"错误详情：{result['error']}")
                        
                        except Exception as e:
                            st.error(f"❌ AI分页测试过程中出现异常: {str(e)}")
                            # 显示详细错误信息帮助调试
                            with st.expander("🔧 错误详情", expanded=True):
                                st.code(str(e))
                                import traceback
                                st.code(traceback.format_exc())
                else:
                    st.warning("⚠️ 请输入要测试的文本内容")
    
    # 页脚信息 - 显示在所有功能页面下方
    st.markdown("---")
    st.markdown(
        '<div style="text-align: center; color: #666; padding: 2rem;">'
        '💡 由AI驱动 | 🎨 专业PPT自动生成'
        '</div>', 
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()