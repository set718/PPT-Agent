#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文本转PPT工具
使用OpenAI GPT-4V API处理文本并生成PowerPoint演示文稿
"""

import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
import json
import re
from config import get_config
from utils import AIProcessor, PPTProcessor, FileManager
from logger import get_logger, log_user_action, log_system_info, LogContext

class TextToPPTGenerator:
    def __init__(self, api_key=None, ppt_path=None):
        """
        初始化文本转PPT生成器
        
        Args:
            api_key (str): OpenAI API密钥
            ppt_path (str): 现有PPT文件路径
        """
        self.config = get_config()
        self.logger = get_logger()
        
        # 初始化API处理器
        self.api_key = api_key
        if not self.api_key:
            raise ValueError("请提供OpenAI API密钥")
        
        self.ai_processor = AIProcessor(self.api_key)
        
        # 验证并加载PPT文件
        self.ppt_path = ppt_path
        is_valid, error_msg = FileManager.validate_ppt_file(ppt_path)
        if not is_valid:
            raise ValueError(f"PPT文件验证失败: {error_msg}")
        
        # 加载现有PPT
        self.presentation = Presentation(self.ppt_path)
        self.ppt_processor = PPTProcessor(self.presentation)
        self.ppt_structure = self.ppt_processor.ppt_structure
        
        self.logger.info(f"初始化文本转PPT生成器，加载文件: {ppt_path}")
    
    
    def process_text_with_deepseek(self, user_text):
        """
        使用OpenAI API分析如何将用户文本填入现有PPT的合适位置
        
        Args:
            user_text (str): 用户输入的文本
            
        Returns:
            dict: 文本分配方案
        """
        log_user_action("AI文本分析", f"文本长度: {len(user_text)}字符")
        
        # 获取增强的结构信息
        enhanced_info = self.ppt_processor.get_enhanced_structure_info()
        
        # 使用增强信息进行分析
        return self.ai_processor.analyze_text_for_ppt(user_text, self.ppt_structure, enhanced_info)
    
    def apply_text_assignments(self, assignments):
        """
        根据分配方案修改现有PPT
        
        Args:
            assignments (dict): 文本分配方案
            
        Returns:
            str: 修改后的PPT文件路径
        """
        log_user_action("应用文本分配", f"分配数量: {len(assignments.get('assignments', []))}")
        
        # 应用分配
        results = self.ppt_processor.apply_assignments(assignments)
        
        # 打印结果
        for result in results:
            print(result)
        
        # 美化演示文稿
        print("\n正在美化PPT布局...")
        beautify_results = self.ppt_processor.beautify_presentation()
        
        # 打印美化结果
        summary = beautify_results['summary']
        print(f"[INFO] 美化完成:")
        print(f"   删除未填充占位符: {summary['removed_placeholders_count']} 个")
        print(f"   重新排版幻灯片: {summary['reorganized_slides_count']} 页")
        print(f"   删除空幻灯片: {summary['removed_empty_slides_count']} 页")
        print(f"   最终幻灯片数: {summary['final_slide_count']} 页")
        
        # 保存修改后的PPT
        return FileManager.save_ppt_to_file(self.presentation)
    
    
    
    def generate_ppt_from_text(self, user_text):
        """
        将用户文本填入现有PPT的完整流程
        
        Args:
            user_text (str): 用户输入的文本
            
        Returns:
            str: 修改后的PPT文件路径
        """
        with LogContext(f"生成PPT文本填充"):
            print("正在使用OpenAI API分析文本结构...")
            assignments = self.process_text_with_deepseek(user_text)
            
            print("正在将您的原始文本填入现有PPT...")
            filepath = self.apply_text_assignments(assignments)
            
            return filepath

def main():
    """主函数"""
    config = get_config()
    logger = get_logger()
    
    print("=" * 50)
    print("        文本转PPT填充器")
    print("    使用OpenAI GPT-4V将文本填入现有PPT")
    print("=" * 50)
    
    # 获取PPT文件路径
    ppt_path = config.default_ppt_template
    
    # 验证PPT文件
    is_valid, error_msg = FileManager.validate_ppt_file(ppt_path)
    if not is_valid:
        print(f"\n[WARNING] PPT文件验证失败:")
        print(f"错误: {error_msg}")
        print(f"路径: {ppt_path}")
        
        print("\n请确认PPT模板文件存在且格式正确。")
        print("您可以：")
        print("1. 将PPT模板文件放置在指定路径")
        print("2. 修改config.py中的default_ppt_template路径")
        print("3. 创建config.json文件进行配置")
        sys.exit(1)
    else:
        print(f"\n[OK] 已找到PPT文件: {os.path.basename(ppt_path)}")
        log_system_info(f"PPT文件验证成功: {ppt_path}")
    
    # 获取用户输入的API密钥
    print("\n" + "="*50)
    print("请输入您的OpenAI API密钥")
    print("获取地址：https://platform.deepseek.com/api_keys")
    print("="*50)
    
    api_key = input("请输入API密钥（sk-开头）: ").strip()
    
    if not api_key:
        print("\n[ERROR] 未输入API密钥，程序退出")
        sys.exit(1)
    
    if not api_key.startswith('sk-'):
        print("\n[WARNING] API密钥格式可能不正确，请确认是否以'sk-'开头")
        confirm = input("是否继续？(y/n): ").strip().lower()
        if confirm not in ['y', 'yes', '是']:
            print("程序退出")
            sys.exit(1)
    
    try:
        # 初始化生成器
        generator = TextToPPTGenerator(api_key, ppt_path)
        
        # 显示现有PPT信息
        ppt_info = generator.ppt_structure
        print(f"\n[INFO] PPT信息:")
        print(f"   总共 {ppt_info['total_slides']} 张幻灯片")
        
        # 显示占位符信息
        total_placeholders = 0
        for slide in ppt_info['slides'][:3]:  # 只显示前3张
            title = slide['title'] if slide['title'] else "（无标题）"
            placeholders = slide.get('placeholders', {})
            total_placeholders += len(placeholders)
            
            if placeholders:
                placeholder_list = ', '.join([f"{{{name}}}" for name in placeholders.keys()])
                print(f"   第{slide['slide_index']+1}页: {title} - 占位符: {placeholder_list}")
            else:
                print(f"   第{slide['slide_index']+1}页: {title} - 无占位符")
        
        if ppt_info['total_slides'] > 3:
            print(f"   ... 还有 {ppt_info['total_slides']-3} 张幻灯片")
        
        print(f"\n[INFO] 总共找到 {total_placeholders} 个占位符")
        
        print("\n请输入您想要填入PPT的文本内容：")
        print("(输入'quit'或'exit'退出)")
        print("-" * 50)
        
        while True:
            try:
                user_input = input("\n请输入文本: ").strip()
                
                if user_input.lower() in ['quit', 'exit', '退出']:
                    print("\n感谢使用！再见！")
                    break
                
                if not user_input:
                    print("请输入有效的文本内容。")
                    continue
                
                # 填入PPT
                filepath = generator.generate_ppt_from_text(user_input)
                
                print(f"\n[OK] PPT更新成功！")
                print(f"[INFO] 文件路径: {filepath}")
                print(f"[INFO] 您可以在 {os.path.abspath(filepath)} 找到更新的PPT文件")
                
                print("\n是否继续添加更多文本？(y/n)")
                continue_choice = input().strip().lower()
                if continue_choice in ['n', 'no', '否', '不']:
                    print("\n感谢使用！再见！")
                    break
                    
            except KeyboardInterrupt:
                print("\n\n程序被用户中断。再见！")
                break
                
            except Exception as e:
                logger.exception(f"生成过程中出现错误: {e}")
                print(f"\n[ERROR] 生成过程中出现错误: {e}")
                print("请重试或检查您的输入。")
    
    except Exception as e:
        logger.exception(f"初始化失败: {e}")
        print(f"\n[ERROR] 初始化失败: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main() 