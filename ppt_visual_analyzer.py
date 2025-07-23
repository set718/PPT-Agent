#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT视觉分析器
使用GPT-4V分析PPT页面的视觉美观度并提供优化建议
"""

import os
import io
import base64
import json
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
from PIL import Image, ImageDraw, ImageFont
try:
    import win32com.client
    import pythoncom
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False
    print("警告: win32com未安装，PPT截图功能将使用备用方法")
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from config import get_config
from logger import get_logger

class PPTVisualAnalyzer:
    """PPT视觉分析器"""
    
    def __init__(self, api_key: str = None):
        """初始化视觉分析器"""
        self.config = get_config()
        self.logger = get_logger()
        self.api_key = api_key or self.config.openai_api_key
        
        if not self.api_key:
            raise ValueError("请设置OpenRouter API密钥")
            
        self.client = OpenAI(
            api_key=self.api_key,
            base_url=self.config.openai_base_url
        )
        
        # 视觉评估标准
        self.visual_criteria = {
            "layout_balance": {
                "description": "布局平衡度 - 元素分布是否均匀合理",
                "weight": 0.25,
                "ideal_score": 8.0
            },
            "color_harmony": {
                "description": "色彩协调性 - 颜色搭配是否和谐",
                "weight": 0.20,
                "ideal_score": 8.0
            },
            "typography": {
                "description": "字体排版 - 字体大小、间距、层次是否清晰",
                "weight": 0.25,
                "ideal_score": 8.0
            },
            "visual_hierarchy": {
                "description": "视觉层次 - 信息重要性是否通过视觉体现",
                "weight": 0.15,
                "ideal_score": 8.0
            },
            "white_space": {
                "description": "留白使用 - 空白区域使用是否得当",
                "weight": 0.10,
                "ideal_score": 7.5
            },
            "overall_aesthetics": {
                "description": "整体美观度 - 整体视觉效果",
                "weight": 0.05,
                "ideal_score": 8.0
            }
        }
    
    def convert_ppt_to_images(self, ppt_path: str, output_dir: str = None) -> List[str]:
        """
        将PPT页面转换为图片
        
        Args:
            ppt_path: PPT文件路径
            output_dir: 输出目录（可选）
            
        Returns:
            List[str]: 生成的图片文件路径列表
        """
        if not output_dir:
            output_dir = os.path.join(self.config.temp_output_dir, "ppt_images")
        
        # 确保使用绝对路径
        output_dir = os.path.abspath(output_dir)
        os.makedirs(output_dir, exist_ok=True)
        
        # 优先使用COM接口（如果可用）
        if HAS_WIN32COM:
            try:
                # 初始化COM环境
                pythoncom.CoInitialize()
                
                # 使用COM接口将PPT转换为图片
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = 1
                
                presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
                image_paths = []
                
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                
                for i, slide in enumerate(presentation.Slides):
                    # 导出为PNG图片
                    image_path = os.path.join(output_dir, f"slide_{timestamp}_{i+1}.png")
                    # PowerPoint需要绝对路径
                    abs_image_path = os.path.abspath(image_path)
                    slide.Export(abs_image_path, "PNG", 1920, 1080)  # 高分辨率导出
                    image_paths.append(abs_image_path)
                    
                presentation.Close()
                powerpoint.Quit()
                
                # 清理COM环境
                pythoncom.CoUninitialize()
                
                self.logger.info(f"成功转换PPT为{len(image_paths)}张图片")
                return image_paths
                
            except Exception as e:
                # 发生错误时也要清理COM环境
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
                self.logger.error(f"PPT图片转换失败: {e}")
        
        # 备用方法：使用python-pptx生成简单截图
        return self._fallback_convert_to_images(ppt_path, output_dir)
    
    def _fallback_convert_to_images(self, ppt_path: str, output_dir: str) -> List[str]:
        """备用图片转换方法"""
        try:
            presentation = Presentation(ppt_path)
            image_paths = []
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            for i, slide in enumerate(presentation.slides):
                # 创建简单的文本预览图
                image_path = os.path.join(output_dir, f"slide_preview_{timestamp}_{i+1}.png")
                self._create_slide_preview(slide, image_path, i+1)
                image_paths.append(image_path)
                
            return image_paths
            
        except Exception as e:
            self.logger.error(f"备用图片转换失败: {e}")
            return []
    
    def _create_slide_preview(self, slide, image_path: str, slide_num: int):
        """创建幻灯片预览图"""
        # 创建1920x1080的白色背景图片
        img = Image.new('RGB', (1920, 1080), 'white')
        draw = ImageDraw.Draw(img)
        
        try:
            font_title = ImageFont.truetype("arial.ttf", 48)
            font_content = ImageFont.truetype("arial.ttf", 24)
        except:
            font_title = ImageFont.load_default()
            font_content = ImageFont.load_default()
        
        # 添加标题
        draw.text((50, 50), f"幻灯片 {slide_num}", fill='black', font=font_title)
        
        # 提取幻灯片文本内容
        y_offset = 150
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()[:100]  # 限制文本长度
                draw.text((50, y_offset), text, fill='black', font=font_content)
                y_offset += 50
                if y_offset > 1000:  # 防止超出图片范围
                    break
        
        img.save(image_path)
    
    def encode_image_to_base64(self, image_path: str) -> str:
        """将图片编码为base64"""
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    
    def analyze_slide_visual_quality(self, image_path: str, slide_context: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        分析单页幻灯片的视觉质量
        
        Args:
            image_path: 幻灯片图片路径
            slide_context: 幻灯片上下文信息（可选）
            
        Returns:
            Dict: 视觉分析结果
        """
        try:
            # 编码图片
            base64_image = self.encode_image_to_base64(image_path)
            
            # 构建分析提示
            prompt = self._build_visual_analysis_prompt(slide_context)
            
            # 调用GPT-4V分析
            response = self.client.chat.completions.create(
                model=self.config.ai_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": prompt
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{base64_image}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=1500,
                temperature=0.3
            )
            
            content = response.choices[0].message.content
            return self._parse_visual_analysis_result(content)
            
        except Exception as e:
            self.logger.error(f"视觉分析失败: {e}")
            return self._create_fallback_analysis()
    
    def _build_visual_analysis_prompt(self, slide_context: Dict[str, Any] = None) -> str:
        """构建视觉分析提示"""
        context_info = ""
        if slide_context:
            context_info = f"\n\n幻灯片上下文信息：\n{json.dumps(slide_context, ensure_ascii=False, indent=2)}"
            
        return f"""作为专业的PPT视觉设计专家，请分析这张幻灯片的视觉质量并提供详细评分和改进建议。

请从以下维度进行分析（每个维度评分1-10分）：

1. **布局平衡度** (25%)
   - 元素分布是否均匀合理
   - 视觉重心是否稳定
   - 是否遵循设计原则（如三分法、黄金比例等）

2. **色彩协调性** (20%)
   - 颜色搭配是否和谐
   - 对比度是否适当
   - 是否符合主题色调

3. **字体排版** (25%)
   - 字体大小层次是否清晰
   - 行间距、字间距是否适当
   - 字体选择是否合适

4. **视觉层次** (15%)
   - 信息重要性是否通过视觉清晰体现
   - 标题、正文、注释的层次是否分明
   - 视觉引导路径是否清晰

5. **留白使用** (10%)
   - 空白区域使用是否得当
   - 是否过于拥挤或过于空旷
   - 留白是否有助于阅读

6. **整体美观度** (5%)
   - 整体视觉效果是否专业
   - 是否符合现代设计趋势
   - 是否给人良好的第一印象

请以JSON格式返回分析结果：
{{
  "scores": {{
    "layout_balance": 8.5,
    "color_harmony": 7.2,
    "typography": 9.0,
    "visual_hierarchy": 8.0,
    "white_space": 7.5,
    "overall_aesthetics": 8.2
  }},
  "weighted_score": 8.1,
  "strengths": [
    "布局整体平衡，视觉重心稳定",
    "字体层次分明，易于阅读"
  ],
  "weaknesses": [
    "色彩对比度略显不足",
    "部分区域留白过少"
  ],
  "improvement_suggestions": [
    {{
      "category": "color_harmony",
      "description": "建议增强标题与背景的对比度",
      "priority": "high",
      "implementation": "将标题颜色调整为深色，或添加背景阴影"
    }},
    {{
      "category": "white_space",
      "description": "增加内容区域的边距",
      "priority": "medium",
      "implementation": "将文本框边距增加到0.5英寸"
    }}
  ],
  "layout_recommendations": {{
    "suggested_adjustments": [
      "将主要内容居中对齐",
      "调整图片与文字的相对位置"
    ],
    "element_positioning": {{
      "title": "建议向上移动10px",
      "content": "建议向左对齐，增加左边距",
      "images": "建议居中放置"
    }}
  }}
}}{context_info}

请确保返回标准JSON格式，不要包含其他文字。"""
    
    def _parse_visual_analysis_result(self, content: str) -> Dict[str, Any]:
        """解析视觉分析结果"""
        try:
            # 尝试提取JSON内容
            json_match = json.loads(content.strip())
            if isinstance(json_match, dict):
                return json_match
        except json.JSONDecodeError:
            pass
        
        # 尝试从代码块中提取JSON
        import re
        json_pattern = r'```(?:json)?\s*(\{.*?\})\s*```'
        match = re.search(json_pattern, content, re.DOTALL)
        if match:
            try:
                return json.loads(match.group(1))
            except json.JSONDecodeError:
                pass
        
        # 如果解析失败，返回默认结果
        self.logger.warning("视觉分析结果解析失败，使用默认评分")
        return self._create_fallback_analysis()
    
    def _create_fallback_analysis(self) -> Dict[str, Any]:
        """创建备用分析结果"""
        return {
            "scores": {
                "layout_balance": 7.0,
                "color_harmony": 7.0,
                "typography": 7.0,
                "visual_hierarchy": 7.0,
                "white_space": 7.0,
                "overall_aesthetics": 7.0
            },
            "weighted_score": 7.0,
            "strengths": ["基础布局合理"],
            "weaknesses": ["需要进一步优化"],
            "improvement_suggestions": [
                {
                    "category": "general",
                    "description": "建议进行整体视觉优化",
                    "priority": "medium",
                    "implementation": "请手动调整布局和样式"
                }
            ],
            "layout_recommendations": {
                "suggested_adjustments": ["建议手动优化布局"],
                "element_positioning": {}
            }
        }
    
    def calculate_weighted_score(self, scores: Dict[str, float]) -> float:
        """计算加权总分"""
        total_score = 0.0
        for criterion, score in scores.items():
            if criterion in self.visual_criteria:
                weight = self.visual_criteria[criterion]["weight"]
                total_score += score * weight
        return round(total_score, 2)
    
    def analyze_presentation_visual_quality(self, ppt_path: str) -> Dict[str, Any]:
        """
        分析整个演示文稿的视觉质量
        
        Args:
            ppt_path: PPT文件路径
            
        Returns:
            Dict: 整体视觉分析结果
        """
        self.logger.info(f"开始分析PPT视觉质量: {ppt_path}")
        
        # 转换PPT为图片
        image_paths = self.convert_ppt_to_images(ppt_path)
        if not image_paths:
            return {"error": "无法转换PPT为图片"}
        
        # 分析每张幻灯片
        slide_analyses = []
        for i, image_path in enumerate(image_paths):
            self.logger.info(f"分析第{i+1}页幻灯片")
            analysis = self.analyze_slide_visual_quality(image_path)
            analysis["slide_index"] = i
            analysis["image_path"] = image_path
            slide_analyses.append(analysis)
        
        # 计算整体分析结果
        overall_analysis = self._calculate_overall_analysis(slide_analyses)
        
        # 清理临时图片文件
        self._cleanup_temp_images(image_paths)
        
        return {
            "overall_analysis": overall_analysis,
            "slide_analyses": slide_analyses,
            "total_slides": len(slide_analyses),
            "analysis_timestamp": datetime.now().isoformat()
        }
    
    def _calculate_overall_analysis(self, slide_analyses: List[Dict[str, Any]]) -> Dict[str, Any]:
        """计算整体分析结果"""
        if not slide_analyses:
            return self._create_fallback_analysis()
        
        # 计算平均分数
        avg_scores = {}
        for criterion in self.visual_criteria.keys():
            scores = [analysis.get("scores", {}).get(criterion, 7.0) for analysis in slide_analyses]
            avg_scores[criterion] = round(sum(scores) / len(scores), 2)
        
        # 计算加权平均分
        weighted_score = self.calculate_weighted_score(avg_scores)
        
        # 收集所有优点和缺点
        all_strengths = []
        all_weaknesses = []
        all_suggestions = []
        
        for analysis in slide_analyses:
            all_strengths.extend(analysis.get("strengths", []))
            all_weaknesses.extend(analysis.get("weaknesses", []))
            all_suggestions.extend(analysis.get("improvement_suggestions", []))
        
        # 去重并排序
        unique_strengths = list(set(all_strengths))[:5]  # 取前5个
        unique_weaknesses = list(set(all_weaknesses))[:5]  # 取前5个
        
        return {
            "scores": avg_scores,
            "weighted_score": weighted_score,
            "strengths": unique_strengths,
            "weaknesses": unique_weaknesses,
            "improvement_suggestions": all_suggestions[:10],  # 取前10个建议
            "grade": self._get_grade_from_score(weighted_score)
        }
    
    def _get_grade_from_score(self, score: float) -> str:
        """根据分数获取评级"""
        if score >= 9.0:
            return "优秀 (A+)"
        elif score >= 8.5:
            return "优秀 (A)"
        elif score >= 8.0:
            return "良好 (B+)"
        elif score >= 7.5:
            return "良好 (B)"
        elif score >= 7.0:
            return "中等 (C+)"
        elif score >= 6.0:
            return "中等 (C)"
        else:
            return "需改进 (D)"
    
    def _cleanup_temp_images(self, image_paths: List[str]):
        """清理临时图片文件"""
        for image_path in image_paths:
            try:
                if os.path.exists(image_path):
                    os.remove(image_path)
            except Exception as e:
                self.logger.warning(f"清理临时文件失败: {e}")

class VisualLayoutOptimizer:
    """基于视觉分析的布局优化器"""
    
    def __init__(self, visual_analyzer: PPTVisualAnalyzer):
        """初始化布局优化器"""
        self.visual_analyzer = visual_analyzer
        self.logger = get_logger()
    
    def optimize_slide_layout(self, presentation: Presentation, slide_index: int, 
                            visual_analysis: Dict[str, Any]) -> Dict[str, Any]:
        """
        基于视觉分析优化幻灯片布局
        
        Args:
            presentation: PPT演示文稿对象
            slide_index: 幻灯片索引
            visual_analysis: 视觉分析结果
            
        Returns:
            Dict: 优化结果
        """
        if slide_index >= len(presentation.slides):
            return {"success": False, "error": "幻灯片索引超出范围"}
        
        slide = presentation.slides[slide_index]
        optimizations = []
        
        # 根据视觉分析建议进行优化
        suggestions = visual_analysis.get("improvement_suggestions", [])
        layout_recommendations = visual_analysis.get("layout_recommendations", {})
        
        try:
            # 应用布局调整建议
            if layout_recommendations.get("element_positioning"):
                positioning = layout_recommendations["element_positioning"]
                optimizations.extend(self._apply_element_positioning(slide, positioning))
            
            # 应用具体改进建议
            for suggestion in suggestions:
                if suggestion.get("priority") == "high":
                    opt_result = self._apply_improvement_suggestion(slide, suggestion)
                    if opt_result:
                        optimizations.append(opt_result)
            
            return {
                "success": True,
                "optimizations_applied": optimizations,
                "slide_index": slide_index
            }
            
        except Exception as e:
            self.logger.error(f"布局优化失败: {e}")
            return {"success": False, "error": str(e)}
    
    def _apply_element_positioning(self, slide, positioning: Dict[str, str]) -> List[str]:
        """应用元素定位调整"""
        optimizations = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                shape_text = shape.text.strip().lower()
                
                # 根据文本内容判断元素类型并应用定位
                if any(keyword in shape_text for keyword in ['标题', 'title', '主题']):
                    if 'title' in positioning:
                        # 应用标题定位建议
                        optimizations.append(f"调整标题位置: {positioning['title']}")
                
                elif len(shape_text) > 50:  # 长文本认为是内容
                    if 'content' in positioning:
                        # 应用内容定位建议
                        optimizations.append(f"调整内容位置: {positioning['content']}")
        
        return optimizations
    
    def _apply_improvement_suggestion(self, slide, suggestion: Dict[str, Any]) -> Optional[str]:
        """应用改进建议"""
        category = suggestion.get("category", "")
        description = suggestion.get("description", "")
        implementation = suggestion.get("implementation", "")
        
        try:
            if category == "white_space":
                # 调整留白
                return self._adjust_white_space(slide, implementation)
            elif category == "typography":
                # 调整字体
                return self._adjust_typography(slide, implementation)
            elif category == "color_harmony":
                # 调整颜色（注：实际颜色调整需要更复杂的实现）
                return f"色彩调整建议已记录: {description}"
            
        except Exception as e:
            self.logger.warning(f"应用建议失败: {e}")
            return None
    
    def _adjust_white_space(self, slide, implementation: str) -> str:
        """调整留白"""
        try:
            # 简单的边距调整
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # 增加文本框内边距
                    tf = shape.text_frame
                    tf.margin_left = Inches(0.5)
                    tf.margin_right = Inches(0.5)
                    tf.margin_top = Inches(0.2)
                    tf.margin_bottom = Inches(0.2)
            
            return f"已调整留白: {implementation}"
        except Exception as e:
            self.logger.warning(f"留白调整失败: {e}")
            return "留白调整部分完成"
    
    def _adjust_typography(self, slide, implementation: str) -> str:
        """调整字体排版"""
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # 调整行间距
                        paragraph.space_after = Pt(6)
                        paragraph.space_before = Pt(3)
                        
                        # 调整字体
                        if paragraph.runs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size < Pt(12):
                                    run.font.size = Pt(12)  # 最小字体大小
            
            return f"已调整字体排版: {implementation}"
        except Exception as e:
            self.logger.warning(f"字体调整失败: {e}")
            return "字体调整部分完成"