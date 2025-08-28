#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
调试占位符清理功能
检查为什么有些占位符没有被正确识别和清理
"""

import os
import re
from pptx import Presentation
from utils import PPTProcessor, PPTAnalyzer

def debug_placeholder_cleanup(ppt_path):
    """调试占位符清理功能"""
    print(f"调试PPT文件: {ppt_path}")
    
    if not os.path.exists(ppt_path):
        print(f"文件不存在: {ppt_path}")
        return
    
    # 加载PPT
    try:
        presentation = Presentation(ppt_path)
        print(f"成功加载PPT，共{len(presentation.slides)}页")
    except Exception as e:
        print(f"加载PPT失败: {e}")
        return
    
    # 创建PPT处理器
    ppt_processor = PPTProcessor(presentation)
    
    # 1. 首先分析PPT结构，查看检测到的占位符
    print("\n=== 1. PPT结构分析 ===")
    ppt_structure = ppt_processor.ppt_structure
    
    all_placeholders = {}
    for slide_idx, slide_info in enumerate(ppt_structure['slides']):
        placeholders = slide_info.get('placeholders', {})
        if placeholders:
            all_placeholders[slide_idx] = placeholders
            print(f"第{slide_idx+1}页检测到{len(placeholders)}个占位符:")
            for placeholder, info in placeholders.items():
                placeholder_type = info.get('type', 'text_box')
                if placeholder_type == 'table_cell':
                    print(f"  - {{{placeholder}}} (表格{info['row_idx']+1},{info['col_idx']+1})")
                else:
                    print(f"  - {{{placeholder}}} (文本框)")
    
    if not all_placeholders:
        print("没有检测到任何占位符")
        return
    
    # 2. 模拟一些占位符填充（只填充部分占位符）
    print("\n=== 2. 模拟占位符填充 ===")
    # 模拟只填充部分占位符
    filled_count = 0
    for slide_idx, placeholders in all_placeholders.items():
        if slide_idx not in ppt_processor.filled_placeholders:
            ppt_processor.filled_placeholders[slide_idx] = set()
        
        # 只填充前3个占位符（模拟部分填充）
        for i, placeholder in enumerate(list(placeholders.keys())):
            if i < 3:  # 只填充前3个
                ppt_processor.filled_placeholders[slide_idx].add(placeholder)
                filled_count += 1
                print(f"  模拟填充: 第{slide_idx+1}页 {{{placeholder}}}")
    
    print(f"共模拟填充了{filled_count}个占位符")
    
    # 3. 显示填充状态
    print("\n=== 3. 占位符填充状态 ===")
    for slide_idx, filled_set in ppt_processor.filled_placeholders.items():
        print(f"第{slide_idx+1}页已填充: {', '.join([f'{{{p}}}' for p in filled_set])}")
    
    # 4. 检查每页的占位符状态
    print("\n=== 4. 逐页检查占位符状态 ===")
    for slide_idx, slide in enumerate(presentation.slides):
        print(f"\n第{slide_idx+1}页:")
        filled_placeholders_in_slide = ppt_processor.filled_placeholders.get(slide_idx, set())
        print(f"  已填充占位符: {filled_placeholders_in_slide}")
        
        # 检查文本框
        text_box_placeholders = []
        table_placeholders = []
        
        for shape in slide.shapes:
            # 检查普通文本框
            if hasattr(shape, 'text') and shape.text:
                placeholder_matches = re.findall(r'\{([^}]+)\}', shape.text)
                if placeholder_matches:
                    text_box_placeholders.extend(placeholder_matches)
                    for p in placeholder_matches:
                        is_filled = p in filled_placeholders_in_slide
                        print(f"    文本框 {{{p}}}: {'已填充' if is_filled else '未填充'}")
            
            # 检查表格
            elif hasattr(shape, 'shape_type') and shape.shape_type == 19:  # 表格
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if cell_text:
                            placeholder_matches = re.findall(r'\{([^}]+)\}', cell_text)
                            if placeholder_matches:
                                for p in placeholder_matches:
                                    table_placeholders.append(p)
                                    is_filled = p in filled_placeholders_in_slide
                                    print(f"    表格{row_idx+1},{col_idx+1} {{{p}}}: {'已填充' if is_filled else '未填充'}")
        
        print(f"  文本框占位符: {len(text_box_placeholders)}个")
        print(f"  表格占位符: {len(table_placeholders)}个")
    
    # 5. 测试清理功能
    print("\n=== 5. 测试清理功能 ===")
    
    # 手动实现清理逻辑，详细记录过程
    cleanup_count = 0
    cleaned_placeholders = []
    
    for slide_idx, slide in enumerate(presentation.slides):
        filled_placeholders_in_slide = ppt_processor.filled_placeholders.get(slide_idx, set())
        print(f"\n处理第{slide_idx+1}页，已填充: {filled_placeholders_in_slide}")
        
        for shape in slide.shapes:
            # 处理文本框
            if hasattr(shape, 'text') and shape.text:
                original_text = shape.text
                placeholder_matches = re.findall(r'\{([^}]+)\}', original_text)
                
                if placeholder_matches:
                    print(f"  文本框包含占位符: {placeholder_matches}")
                    
                    unfilled_placeholders = [
                        p for p in placeholder_matches 
                        if p not in filled_placeholders_in_slide
                    ]
                    
                    if unfilled_placeholders:
                        print(f"    需要清理的未填充占位符: {unfilled_placeholders}")
                        cleaned_text = original_text
                        for unfilled_placeholder in unfilled_placeholders:
                            pattern = f"{{{unfilled_placeholder}}}"
                            cleaned_text = cleaned_text.replace(pattern, "")
                            cleaned_placeholders.append(f"第{slide_idx+1}页(文本框): {{{unfilled_placeholder}}}")
                            print(f"    清理 {{{unfilled_placeholder}}}")
                        
                        cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()
                        if cleaned_text != original_text:
                            # 不实际修改，只记录
                            cleanup_count += 1
                            # 避免Unicode编码错误，使用repr显示特殊字符
                            print(f"    文本框内容变化: {repr(original_text)} -> {repr(cleaned_text)}")
                    else:
                        print("    所有占位符都已填充，无需清理")
            
            # 处理表格
            elif hasattr(shape, 'shape_type') and shape.shape_type == 19:
                table = shape.table
                print(f"  处理表格 ({len(table.rows)}行, {len(table.rows[0].cells) if table.rows else 0}列)")
                
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        original_cell_text = cell.text.strip()
                        if original_cell_text:
                            placeholder_matches = re.findall(r'\{([^}]+)\}', original_cell_text)
                            
                            if placeholder_matches:
                                print(f"    表格{row_idx+1},{col_idx+1}包含占位符: {placeholder_matches}")
                                
                                unfilled_placeholders = [
                                    p for p in placeholder_matches 
                                    if p not in filled_placeholders_in_slide
                                ]
                                
                                if unfilled_placeholders:
                                    print(f"      需要清理的未填充占位符: {unfilled_placeholders}")
                                    cleaned_cell_text = original_cell_text
                                    for unfilled_placeholder in unfilled_placeholders:
                                        pattern = f"{{{unfilled_placeholder}}}"
                                        cleaned_cell_text = cleaned_cell_text.replace(pattern, "")
                                        cleaned_placeholders.append(f"第{slide_idx+1}页(表格{row_idx+1},{col_idx+1}): {{{unfilled_placeholder}}}")
                                        print(f"      清理 {{{unfilled_placeholder}}}")
                                    
                                    cleaned_cell_text = re.sub(r'\s+', ' ', cleaned_cell_text).strip()
                                    if cleaned_cell_text != original_cell_text:
                                        cleanup_count += 1
                                        print(f"      单元格内容变化: {repr(original_cell_text)} -> {repr(cleaned_cell_text)}")
                                else:
                                    print("      所有占位符都已填充，无需清理")
    
    print(f"\n=== 清理总结 ===")
    print(f"需要清理的占位符数量: {len(cleaned_placeholders)}")
    print(f"修改的文本框/单元格数量: {cleanup_count}")
    
    if cleaned_placeholders:
        print("需要清理的占位符列表:")
        for item in cleaned_placeholders:
            print(f"  - {item}")
    else:
        print("没有需要清理的占位符")

def main():
    """主函数"""
    print("=== 占位符清理功能调试 ===")
    
    # 可以指定要测试的PPT文件路径
    test_files = [
        "test_placeholder_template.pptx",  # 新创建的测试文件
        "test_template.pptx",  # 如果用户有测试文件
        "templates/ppt_template/split_presentations_1.pptx",
        "templates/ppt_template/split_presentations_2.pptx",
        "templates/title_slides.pptx",
        "templates/table_of_contents_slides.pptx"
    ]
    
    for test_file in test_files:
        if os.path.exists(test_file):
            debug_placeholder_cleanup(test_file)
            break
    else:
        print("没有找到可用的测试文件")
        print("请将包含占位符的PPT文件放在当前目录，或修改脚本中的文件路径")

if __name__ == "__main__":
    main()