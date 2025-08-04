#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
测试自定义模板功能的脚本
验证新增的模板上传和测试功能是否正常工作
"""

import os
import sys
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_test_template():
    """创建一个简单的测试PPT模板"""
    print("正在创建测试PPT模板...")
    
    # 创建新的演示文稿
    prs = Presentation()
    
    # 删除默认的空白幻灯片
    xml_slides = prs.slides._sldIdLst[:]
    for slide in xml_slides:
        prs.slides._sldIdLst.remove(slide)
    
    # 添加标题页
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "{项目标题}"
    subtitle.text = "{副标题}\n{日期}"
    
    # 添加内容页1
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "{章节标题1}"
    content.text = """主要内容：
• {要点1}
• {要点2}
• {要点3}

详细描述：
{详细内容1}"""
    
    # 添加内容页2
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "{章节标题2}"
    content.text = """核心观点：
{核心观点}

支撑论据：
• {论据1}
• {论据2}

{详细内容2}"""
    
    # 添加结论页
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "{结论标题}"
    content.text = """{总结内容}

关键收获：
• {收获1}
• {收获2}
• {收获3}

{结语}"""
    
    # 保存模板
    template_path = os.path.join(os.path.dirname(__file__), "test_custom_template.pptx")
    prs.save(template_path)
    
    print(f"✅ 测试模板已创建: {template_path}")
    return template_path

def test_template_analysis(template_path):
    """测试模板分析功能"""
    print("\n正在测试模板分析功能...")
    
    try:
        from pptx import Presentation
        import re
        
        presentation = Presentation(template_path)
        
        print(f"📑 幻灯片数量: {len(presentation.slides)}")
        
        total_placeholders = 0
        for i, slide in enumerate(presentation.slides):
            slide_placeholders = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    placeholders = re.findall(r'\{([^}]+)\}', shape.text)
                    if placeholders:
                        slide_placeholders.extend(placeholders)
                        total_placeholders += len(placeholders)
            
            if slide_placeholders:
                print(f"第{i+1}页占位符: {', '.join([f'{{{p}}}' for p in slide_placeholders])}")
        
        print(f"🎯 总占位符数量: {total_placeholders}")
        print("✅ 模板分析测试通过")
        
    except Exception as e:
        print(f"❌ 模板分析测试失败: {e}")

def test_custom_template_feature():
    """测试自定义模板功能的基本组件"""
    print("=" * 50)
    print("测试自定义模板功能")
    print("=" * 50)
    
    # 创建测试模板
    template_path = create_test_template()
    
    # 测试模板分析
    test_template_analysis(template_path)
    
    # 测试文件验证
    print("\n正在测试文件验证功能...")
    try:
        from utils import FileManager
        
        is_valid, error_msg = FileManager.validate_ppt_file(template_path)
        if is_valid:
            print("✅ 文件验证通过")
        else:
            print(f"❌ 文件验证失败: {error_msg}")
    except Exception as e:
        print(f"❌ 文件验证测试失败: {e}")
    
    print("\n" + "=" * 50)
    print("测试完成！")
    print("=" * 50)
    print("\n📋 使用说明:")
    print("1. 启动用户界面: python run_user_app.py")
    print("2. 选择'自定义模板测试'标签页")
    print(f"3. 上传刚创建的测试模板: {template_path}")
    print("4. 输入测试文本并运行测试")
    print("\n💡 测试模板包含以下占位符:")
    print("   {项目标题}, {副标题}, {日期}, {章节标题1}, {要点1}, {要点2}, {要点3}")
    print("   {详细内容1}, {章节标题2}, {核心观点}, {论据1}, {论据2}, {详细内容2}")
    print("   {结论标题}, {总结内容}, {收获1}, {收获2}, {收获3}, {结语}")

if __name__ == "__main__":
    try:
        test_custom_template_feature()
    except KeyboardInterrupt:
        print("\n\n测试被用户中断")
    except Exception as e:
        print(f"\n❌ 测试过程中出现错误: {e}")
        import traceback
        traceback.print_exc()