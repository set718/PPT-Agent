#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
测试占位符清理计数修复
"""

from pptx import Presentation
from pptx.util import Inches

def create_multi_placeholder_test():
    """创建包含多个占位符在同一文本框的测试PPT"""
    presentation = Presentation()
    
    # 第一页：一个文本框包含多个占位符
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    
    # 标题
    title_shape = slide1.shapes.title
    title_shape.text = "多占位符测试"
    
    # 内容文本框 - 包含多个占位符
    content_shape = slide1.placeholders[1]
    content_shape.text = """
主要内容：
• {bullet_1_summary} - 要点1总结
• {bullet_1_content} - 要点1详细内容  
• {bullet_2_summary} - 要点2总结
• {bullet_2_content} - 要点2详细内容
结论: {conclusion}
"""
    
    # 第二页：另一个包含重复占位符的文本框
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    
    # 添加文本框
    text_box = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.text = """
重复测试：
第一次提到: {bullet_1_summary}
第二次提到: {bullet_1_content}  
第三次提到: {bullet_1_summary}
第四次提到: {bullet_1_content}
最后提到: {bullet_2}
"""
    
    # 保存文件
    test_file = "multi_placeholder_count_test.pptx"
    presentation.save(test_file)
    
    print(f"已创建测试文件: {test_file}")
    
    # 统计占位符
    total_placeholders = []
    for slide_idx, slide in enumerate(presentation.slides):
        slide_placeholders = []
        text_box_count = 0
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                text_box_count += 1
                import re
                placeholders = re.findall(r'\{([^}]+)\}', shape.text)
                slide_placeholders.extend(placeholders)
        
        if slide_placeholders:
            unique_placeholders = list(set(slide_placeholders))
            print(f"第{slide_idx+1}页：")
            print(f"  - 文本框数量: {text_box_count}")
            print(f"  - 占位符总数: {len(slide_placeholders)} (包括重复)")
            print(f"  - 唯一占位符: {len(unique_placeholders)}")
            print(f"  - 占位符列表: {slide_placeholders}")
            print(f"  - 唯一占位符列表: {unique_placeholders}")
            total_placeholders.extend(slide_placeholders)
    
    print(f"\n总计:")
    print(f"  - 所有占位符实例: {len(total_placeholders)}")
    print(f"  - 唯一占位符类型: {len(set(total_placeholders))}")
    
    return test_file

def test_cleanup_count_logic():
    """测试清理计数逻辑"""
    print("=== 测试占位符清理计数修复 ===\n")
    
    # 1. 创建测试文件
    test_file = create_multi_placeholder_test()
    
    # 2. 加载并处理
    from user_app import UserPPTGenerator
    
    try:
        generator = UserPPTGenerator("test-api-key")
        success, message = generator.load_ppt_from_path(test_file)
        
        if not success:
            print(f"加载失败: {message}")
            return
        
        print(f"\n加载成功: {message}")
        
        # 3. 模拟填充部分占位符（只填充少数几个）
        print("\n3. 模拟填充部分占位符...")
        test_assignments = {
            "assignments": [
                {"slide_index": 0, "action": "replace_placeholder", "placeholder": "bullet_1_summary", "content": "测试总结1", "reason": "测试"},
                {"slide_index": 0, "action": "replace_placeholder", "placeholder": "conclusion", "content": "测试结论", "reason": "测试"},
                # 故意不填充 bullet_1_content, bullet_2_summary, bullet_2_content, bullet_2
            ]
        }
        
        success, results = generator.apply_text_assignments(test_assignments, "测试原始文本")
        
        if success:
            print("占位符填充完成")
            
            # 显示填充记录
            filled_placeholders = generator.ppt_processor.filled_placeholders
            total_filled = sum(len(s) for s in filled_placeholders.values())
            print(f"已填充占位符数量: {total_filled}")
            
            for slide_idx, filled_set in filled_placeholders.items():
                if filled_set:
                    print(f"  第{slide_idx+1}页: {', '.join([f'{{{p}}}' for p in filled_set])}")
        
        # 4. 执行清理并检查计数
        print("\n4. 执行占位符清理...")
        cleanup_results = generator.cleanup_unfilled_placeholders()
        
        if cleanup_results.get('success'):
            cleaned_count = cleanup_results.get('cleaned_placeholders', 0)
            cleaned_list = cleanup_results.get('cleaned_placeholder_list', [])
            message = cleanup_results.get('message', '')
            
            print(f"清理完成!")
            print(f"报告的清理数量: {cleaned_count}")
            print(f"实际清理列表长度: {len(cleaned_list)}")
            print(f"消息: {message}")
            
            # 验证计数是否正确
            if cleaned_count == len(cleaned_list):
                print("计数正确！清理数量与实际清理列表一致")
            else:
                print(f"计数错误！报告{cleaned_count}个，但实际清理{len(cleaned_list)}个")
            
            print("\n清理详情:")
            for i, item in enumerate(cleaned_list, 1):
                print(f"  {i}. {item}")
                
            # 分析清理的占位符类型
            placeholder_types = {}
            for item in cleaned_list:
                # 提取占位符名称（例如从"第1页(文本框): {bullet_1_content}"中提取"bullet_1_content"）
                import re
                match = re.search(r'\{([^}]+)\}', item)
                if match:
                    placeholder_name = match.group(1)
                    placeholder_types[placeholder_name] = placeholder_types.get(placeholder_name, 0) + 1
            
            print(f"\n按类型统计清理的占位符:")
            for placeholder_type, count in placeholder_types.items():
                print(f"  - {{{placeholder_type}}}: {count}次")
            
        else:
            error_msg = cleanup_results.get('error', '未知错误')
            print(f"清理失败: {error_msg}")
    
    except Exception as e:
        print(f"测试过程中出现异常: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_cleanup_count_logic()