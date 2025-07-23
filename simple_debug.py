#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
简化的调试脚本
"""

import os
import json
from pptx import Presentation
from config import get_config
from utils import PPTAnalyzer

def main():
    print("检查PPT模板占位符...")
    
    config = get_config()
    ppt_path = config.default_ppt_template
    
    if not os.path.exists(ppt_path):
        print(f"PPT模板文件不存在: {ppt_path}")
        return
    
    presentation = Presentation(ppt_path)
    ppt_structure = PPTAnalyzer.analyze_ppt_structure(presentation)
    
    print(f"总页数: {ppt_structure['total_slides']}")
    print("\n各页占位符:")
    
    for i, slide_info in enumerate(ppt_structure['slides']):
        print(f"\n第{i+1}页:")
        print(f"  标题: {slide_info['title'] or '无'}")
        print(f"  占位符数量: {len(slide_info['placeholders'])}")
        
        if slide_info['placeholders']:
            for placeholder_name, placeholder_info in slide_info['placeholders'].items():
                original_text = placeholder_info['original_text']
                print(f"    {{{placeholder_name}}} -> '{original_text}'")
        else:
            print("    无占位符")
    
    # 直接检查所有文本内容
    print("\n=" * 40)
    print("所有文本内容检查:")
    
    for i, slide in enumerate(presentation.slides):
        print(f"\n第{i+1}页所有文本:")
        text_count = 0
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    text_count += 1
                    print(f"  文本{text_count}: '{text}'")
        
        if text_count == 0:
            print("  (无文本内容)")

if __name__ == "__main__":
    main()