#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
完整流程测试脚本
测试：文本分析 -> 内容填充 -> 清理占位符 -> 保存结果
"""

import os
import json
from pptx import Presentation
from config import get_config
from utils import AIProcessor, PPTProcessor, PPTAnalyzer

def test_complete_flow():
    print("=== 完整PPT文本填充流程测试 ===")
    
    config = get_config()
    ppt_path = config.default_ppt_template
    
    if not os.path.exists(ppt_path):
        print(f"PPT文件不存在: {ppt_path}")
        return
    
    # 跳过交互，直接使用手动测试
    api_key = ""
    
    # 测试文本
    test_text = """
    人工智能发展历程
    
    人工智能技术的发展经历了三个重要阶段：
    
    第一阶段是符号主义时期，从1950年代开始，强调逻辑推理和知识表示。
    第二阶段是专家系统的兴起，在1980年代达到高峰。
    第三阶段是深度学习的突破，带来了革命性的进展。
    
    当前的大语言模型如GPT、Claude等展现出前所未有的能力，能够进行复杂的文本理解和生成。
    
    未来，人工智能将继续向更加智能化、人性化的方向发展。
    """
    
    print(f"测试文本: {test_text.strip()}")
    print()
    
    # 步骤1: 加载和分析PPT结构
    print("步骤1: 分析PPT结构")
    print("-" * 30)
    
    presentation = Presentation(ppt_path)
    ppt_structure = PPTAnalyzer.analyze_ppt_structure(presentation)
    ppt_processor = PPTProcessor(presentation)
    
    print(f"模板包含 {len(ppt_structure['slides'])} 页，共 {sum(len(slide['placeholders']) for slide in ppt_structure['slides'])} 个占位符")
    
    # 显示主要占位符
    main_placeholders = []
    for slide_info in ppt_structure['slides']:
        for placeholder_name in slide_info['placeholders'].keys():
            if any(keyword in placeholder_name.lower() for keyword in ['title', 'content_1', 'content_2']):
                main_placeholders.append(placeholder_name)
    
    print(f"主要占位符: {main_placeholders[:5]}")  # 只显示前5个
    print()
    
    # 步骤2: AI分析和分配
    assignments = None
    if api_key:
        print("步骤2: AI分析文本")
        print("-" * 30)
        
        try:
            ai_processor = AIProcessor(api_key)
            assignments = ai_processor.analyze_text_for_ppt(test_text, ppt_structure)
            
            print("AI分配结果:")
            assignments_list = assignments.get('assignments', [])
            for i, assignment in enumerate(assignments_list):
                placeholder = assignment.get('placeholder', '')
                content = assignment.get('content', '')[:50] + '...' if len(assignment.get('content', '')) > 50 else assignment.get('content', '')
                print(f"  {i+1}. {placeholder} -> '{content}'")
            
            print(f"共生成 {len(assignments_list)} 个分配方案")
            
        except Exception as e:
            print(f"AI分析失败: {e}")
            assignments = None
    else:
        print("步骤2: 跳过AI测试，使用手动分配")
        print("-" * 30)
        
        # 手动创建简单的分配方案
        assignments = {
            "assignments": [
                {
                    "slide_index": 0,
                    "action": "replace_placeholder",
                    "placeholder": "title",
                    "content": "人工智能发展历程",
                    "reason": "提取标题"
                },
                {
                    "slide_index": 0,
                    "action": "replace_placeholder",
                    "placeholder": "content_1",
                    "content": "符号主义时期：从1950年代开始，强调逻辑推理",
                    "reason": "第一阶段内容"
                },
                {
                    "slide_index": 0,
                    "action": "replace_placeholder",
                    "placeholder": "content_2",
                    "content": "专家系统兴起：1980年代达到高峰",
                    "reason": "第二阶段内容"
                }
            ]
        }
        print("使用手动分配方案，共3个分配")
    
    print()
    
    if not assignments:
        print("没有分配方案，测试结束")
        return
    
    # 步骤3: 应用文本分配
    print("步骤3: 应用文本分配")
    print("-" * 30)
    
    results = ppt_processor.apply_assignments(assignments)
    
    print("填充结果:")
    for result in results:
        status = "SUCCESS" if "SUCCESS" in result else "ERROR"
        print(f"  {status}: {result}")
    
    print()
    
    # 步骤4: 清理剩余占位符
    print("步骤4: 清理剩余占位符")
    print("-" * 30)
    
    cleanup_count = 0
    for slide_idx, slide in enumerate(presentation.slides):
        slide_cleanup = 0
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                original_text = shape.text
                import re
                cleaned_text = re.sub(r'\{[^}]+\}', '', original_text)
                if cleaned_text != original_text:
                    shape.text = cleaned_text.strip()
                    slide_cleanup += 1
                    cleanup_count += 1
        
        if slide_cleanup > 0:
            print(f"  第{slide_idx+1}页清理了 {slide_cleanup} 个占位符")
    
    print(f"总共清理了 {cleanup_count} 个未填充的占位符")
    print()
    
    # 步骤5: 验证结果
    print("步骤5: 验证最终结果")
    print("-" * 30)
    
    for i, slide in enumerate(presentation.slides):
        print(f"第{i+1}页内容:")
        text_count = 0
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                text_count += 1
                # 检查是否还有未清理的占位符
                remaining_placeholders = re.findall(r'\{[^}]+\}', text)
                status = " (含占位符)" if remaining_placeholders else ""
                print(f"  文本{text_count}: '{text[:60]}...'{status}")
        
        if text_count == 0:
            print("  (无内容)")
    
    print()
    
    # 步骤6: 保存结果
    print("步骤6: 保存结果")
    print("-" * 30)
    
    output_path = "output/complete_flow_test.pptx"
    os.makedirs("output", exist_ok=True)
    presentation.save(output_path)
    
    print(f"已保存测试结果: {output_path}")
    print("测试完成！")

if __name__ == "__main__":
    test_complete_flow()