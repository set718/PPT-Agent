#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT高级分析器模块
包含PPT结构解析器、位置信息提取器、智能布局调整器
用于帮助AI更好地理解PPT整体结构、元素位置和布局优化
"""

import os
import re
from typing import Dict, List, Any, Optional, Tuple, TYPE_CHECKING
from dataclasses import dataclass
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from logger import get_logger

if TYPE_CHECKING:
    from pptx.presentation import Presentation
else:
    from pptx import Presentation

@dataclass
class ElementPosition:
    """元素位置信息"""
    left: float
    top: float
    width: float
    height: float
    center_x: float
    center_y: float
    right: float
    bottom: float
    
    def __post_init__(self):
        self.center_x = self.left + self.width / 2
        self.center_y = self.top + self.height / 2
        self.right = self.left + self.width
        self.bottom = self.top + self.height

@dataclass
class ElementInfo:
    """元素详细信息"""
    element_id: str
    element_type: str
    position: ElementPosition
    content: str
    placeholder_name: Optional[str]
    visual_weight: int  # 1-5，数字越大视觉权重越高
    is_title: bool
    is_content: bool
    is_decorative: bool
    font_size: Optional[float]
    shape_obj: Any  # 原始shape对象

@dataclass
class SlideLayout:
    """幻灯片布局信息"""
    slide_index: int
    layout_type: str
    visual_regions: Dict[str, List[ElementInfo]]
    content_density: float  # 0-1，内容密度
    visual_balance: float  # 0-1，视觉平衡度
    hierarchy_clarity: float  # 0-1，层次清晰度
    elements: List[ElementInfo]
    design_intent: str

class PPTStructureAnalyzer:
    """PPT结构解析器"""
    
    def __init__(self, presentation: Presentation):
        self.presentation = presentation
        self.logger = get_logger()
        self.slide_layouts = []
        self.overall_structure = {}
        
    def analyze_complete_structure(self) -> Dict[str, Any]:
        """完整分析PPT结构"""
        self.logger.info("开始完整PPT结构分析")
        
        # 分析每张幻灯片
        for i, slide in enumerate(self.presentation.slides):
            slide_layout = self._analyze_slide_layout(slide, i)
            self.slide_layouts.append(slide_layout)
        
        # 分析整体结构
        self.overall_structure = self._analyze_overall_structure()
        
        return {
            'slide_layouts': self.slide_layouts,
            'overall_structure': self.overall_structure,
            'analysis_summary': self._generate_analysis_summary()
        }
    
    def _analyze_slide_layout(self, slide, slide_index: int) -> SlideLayout:
        """分析单张幻灯片布局"""
        elements = []
        
        # 获取幻灯片尺寸
        slide_width = float(self.presentation.slide_width or 0)
        slide_height = float(self.presentation.slide_height or 0)
        
        for shape in slide.shapes:
            element_info = self._extract_element_info(shape, slide_width, slide_height)
            if element_info:
                elements.append(element_info)
        
        # 分析布局类型
        layout_type = self._detect_layout_type(elements)
        
        # 分析视觉区域
        visual_regions = self._analyze_visual_regions(elements, slide_width, slide_height)
        
        # 计算布局指标
        content_density = self._calculate_content_density(elements, slide_width, slide_height)
        visual_balance = self._calculate_visual_balance(elements, slide_width, slide_height)
        hierarchy_clarity = self._calculate_hierarchy_clarity(elements)
        
        # 分析设计意图
        design_intent = self._analyze_design_intent(elements, layout_type)
        
        return SlideLayout(
            slide_index=slide_index,
            layout_type=layout_type,
            visual_regions=visual_regions,
            content_density=content_density,
            visual_balance=visual_balance,
            hierarchy_clarity=hierarchy_clarity,
            elements=elements,
            design_intent=design_intent
        )
    
    def _extract_element_info(self, shape, slide_width: float, slide_height: float) -> Optional[ElementInfo]:
        """提取元素信息"""
        try:
            # 获取位置信息
            position = ElementPosition(
                left=float(shape.left),
                top=float(shape.top),
                width=float(shape.width),
                height=float(shape.height),
                center_x=0, center_y=0, right=0, bottom=0  # 将在__post_init__中计算
            )
            
            # 获取内容
            content = ""
            if hasattr(shape, 'text'):
                content = shape.text.strip()
            
            # 检测占位符 - 识别所有{}格式的占位符
            placeholder_name = None
            placeholder_matches = []
            if content and re.search(r'\{([^}]+)\}', content):
                placeholder_matches = re.findall(r'\{([^}]+)\}', content)
                placeholder_name = placeholder_matches[0] if placeholder_matches else None
            
            # 分析元素类型
            element_type = self._detect_element_type(shape)
            
            # 计算视觉权重
            visual_weight = self._calculate_visual_weight(shape, position, content, slide_width, slide_height)
            
            # 分析元素属性
            is_title = self._is_title_element(shape, content, position, slide_height)
            is_content = self._is_content_element(shape, content)
            is_decorative = self._is_decorative_element(shape, content)
            
            # 获取字体大小
            font_size = self._get_font_size(shape)
            
            return ElementInfo(
                element_id=f"slide_{slide_width}_{shape.shape_id if hasattr(shape, 'shape_id') else id(shape)}",
                element_type=element_type,
                position=position,
                content=content,
                placeholder_name=placeholder_name,
                placeholder_matches=placeholder_matches,  # 添加所有占位符匹配信息
                visual_weight=visual_weight,
                is_title=is_title,
                is_content=is_content,
                is_decorative=is_decorative,
                font_size=font_size,
                shape_obj=shape
            )
            
        except Exception as e:
            self.logger.error(f"提取元素信息时出错: {e}")
            return None
    
    def _detect_element_type(self, shape) -> str:
        """检测元素类型"""
        if hasattr(shape, 'shape_type'):
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return "text_box"
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                return "placeholder"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return "picture"
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return "auto_shape"
        return "unknown"
    
    def _calculate_visual_weight(self, shape, position: ElementPosition, content: str, slide_width: float, slide_height: float) -> int:
        """计算视觉权重 (1-5)"""
        weight = 1
        
        # 基于位置权重（上方和左上方权重更高）
        if position.top < slide_height * 0.3:
            weight += 1
        if position.left < slide_width * 0.3:
            weight += 1
        
        # 基于大小权重
        area_ratio = (position.width * position.height) / (slide_width * slide_height)
        if area_ratio > 0.2:
            weight += 1
        if area_ratio > 0.4:
            weight += 1
        
        # 基于内容权重
        if content:
            if len(content) < 20:  # 短文本通常是标题
                weight += 1
            if any(keyword in content.lower() for keyword in ['title', 'subtitle', 'heading']):
                weight += 1
        
        return min(weight, 5)
    
    def _is_title_element(self, shape, content: str, position: ElementPosition, slide_height: float) -> bool:
        """判断是否为标题元素"""
        # 位置判断：位于幻灯片上方
        if position.top > slide_height * 0.4:
            return False
        
        # 内容判断
        if not content:
            return False
        
        # 长度判断：标题通常较短
        if len(content) > 100:
            return False
        
        # 关键词判断
        title_keywords = ['title', 'subtitle', 'heading', 'topic']
        if any(keyword in content.lower() for keyword in title_keywords):
            return True
        
        # 占位符判断 - 智能识别标题类占位符
        placeholder_matches = re.findall(r'\{([^}]+)\}', content)
        if placeholder_matches:
            placeholder_name = placeholder_matches[0].lower()
            title_keywords = ['title', 'heading', '主题', 'topic', '标题', 'header']
            if any(keyword in placeholder_name for keyword in title_keywords):
                return True
        
        return False
    
    def _is_content_element(self, shape, content: str) -> bool:
        """判断是否为内容元素"""
        if not content:
            return False
        
        # 占位符判断 - 智能识别内容类占位符
        placeholder_matches = re.findall(r'\{([^}]+)\}', content)
        if placeholder_matches:
            placeholder_name = placeholder_matches[0].lower()
            content_keywords = ['content', 'bullet', 'description', 'text', '介绍', '内容', '描述', 'detail', 'point', 'list', 'item', '要点', '列表', '项目']
            return any(keyword in placeholder_name for keyword in content_keywords)
        
        # 长度判断：内容通常较长
        return len(content) > 20
    
    def _is_decorative_element(self, shape, content: str) -> bool:
        """判断是否为装饰元素"""
        if hasattr(shape, 'shape_type'):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not content:
                return True
        return False
    
    def _get_font_size(self, shape) -> Optional[float]:
        """获取字体大小"""
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            return float(run.font.size.pt)
        except:
            pass
        return None
    
    def _detect_layout_type(self, elements: List[ElementInfo]) -> str:
        """检测布局类型"""
        title_count = sum(1 for elem in elements if elem.is_title)
        content_count = sum(1 for elem in elements if elem.is_content)
        total_elements = len(elements)
        
        if title_count >= 1 and content_count >= 3:
            return "title_with_bullets"
        elif title_count >= 1 and content_count >= 1:
            return "title_with_content"
        elif content_count >= 4:
            return "content_grid"
        elif total_elements <= 2:
            return "simple"
        else:
            return "complex"
    
    def _analyze_visual_regions(self, elements: List[ElementInfo], slide_width: float, slide_height: float) -> Dict[str, List[ElementInfo]]:
        """分析视觉区域"""
        regions = {
            'header': [],
            'main_left': [],
            'main_center': [],
            'main_right': [],
            'footer': []
        }
        
        for element in elements:
            pos = element.position
            
            # 头部区域（上方20%）
            if pos.top < slide_height * 0.2:
                regions['header'].append(element)
            # 底部区域（下方20%）
            elif pos.top > slide_height * 0.8:
                regions['footer'].append(element)
            # 主要区域
            else:
                # 左侧区域
                if pos.center_x < slide_width * 0.33:
                    regions['main_left'].append(element)
                # 右侧区域
                elif pos.center_x > slide_width * 0.67:
                    regions['main_right'].append(element)
                # 中央区域
                else:
                    regions['main_center'].append(element)
        
        return regions
    
    def _calculate_content_density(self, elements: List[ElementInfo], slide_width: float, slide_height: float) -> float:
        """计算内容密度"""
        if not elements:
            return 0.0
        
        total_content_area = sum(elem.position.width * elem.position.height for elem in elements if elem.is_content)
        slide_area = slide_width * slide_height
        
        return min(total_content_area / slide_area, 1.0)
    
    def _calculate_visual_balance(self, elements: List[ElementInfo], slide_width: float, slide_height: float) -> float:
        """计算视觉平衡度"""
        if not elements:
            return 1.0
        
        # 计算重心
        total_weight = sum(elem.visual_weight for elem in elements)
        if total_weight == 0:
            return 1.0
        
        center_x = sum(elem.position.center_x * elem.visual_weight for elem in elements) / total_weight
        center_y = sum(elem.position.center_y * elem.visual_weight for elem in elements) / total_weight
        
        # 计算偏离中心的程度
        slide_center_x = slide_width / 2
        slide_center_y = slide_height / 2
        
        deviation_x = abs(center_x - slide_center_x) / slide_center_x
        deviation_y = abs(center_y - slide_center_y) / slide_center_y
        
        # 平衡度 = 1 - 平均偏离度
        balance = 1.0 - (deviation_x + deviation_y) / 2
        return max(balance, 0.0)
    
    def _calculate_hierarchy_clarity(self, elements: List[ElementInfo]) -> float:
        """计算层次清晰度"""
        if not elements:
            return 1.0
        
        # 统计不同视觉权重的元素数量
        weight_distribution = {}
        for elem in elements:
            weight_distribution[elem.visual_weight] = weight_distribution.get(elem.visual_weight, 0) + 1
        
        # 层次清晰度 = 权重分布的多样性
        unique_weights = len(weight_distribution)
        max_possible_weights = 5
        
        return min(unique_weights / max_possible_weights, 1.0)
    
    def _analyze_design_intent(self, elements: List[ElementInfo], layout_type: str) -> str:
        """分析设计意图"""
        title_elements = [elem for elem in elements if elem.is_title]
        content_elements = [elem for elem in elements if elem.is_content]
        
        if layout_type == "title_with_bullets":
            return f"标题要点型设计，适合概要展示。包含{len(title_elements)}个标题元素和{len(content_elements)}个内容要点。"
        elif layout_type == "title_with_content":
            return f"标题内容型设计，适合主题阐述。包含{len(title_elements)}个标题元素和{len(content_elements)}个内容区域。"
        elif layout_type == "content_grid":
            return f"内容网格型设计，适合信息并列展示。包含{len(content_elements)}个内容区域。"
        elif layout_type == "simple":
            return f"简洁型设计，适合重点突出。包含{len(elements)}个核心元素。"
        else:
            return f"复合型设计，包含多种元素类型。需要灵活处理{len(elements)}个元素。"
    
    def _analyze_overall_structure(self) -> Dict[str, Any]:
        """分析整体结构"""
        total_slides = len(self.slide_layouts)
        
        # 统计布局类型分布
        layout_distribution = {}
        for layout in self.slide_layouts:
            layout_type = layout.layout_type
            layout_distribution[layout_type] = layout_distribution.get(layout_type, 0) + 1
        
        # 计算平均指标
        avg_density = sum(layout.content_density for layout in self.slide_layouts) / total_slides if total_slides > 0 else 0
        avg_balance = sum(layout.visual_balance for layout in self.slide_layouts) / total_slides if total_slides > 0 else 0
        avg_clarity = sum(layout.hierarchy_clarity for layout in self.slide_layouts) / total_slides if total_slides > 0 else 0
        
        # 分析整体风格
        overall_style = self._determine_overall_style(layout_distribution)
        
        return {
            'total_slides': total_slides,
            'layout_distribution': layout_distribution,
            'average_metrics': {
                'content_density': avg_density,
                'visual_balance': avg_balance,
                'hierarchy_clarity': avg_clarity
            },
            'overall_style': overall_style,
            'design_consistency': self._calculate_design_consistency()
        }
    
    def _determine_overall_style(self, layout_distribution: Dict[str, int]) -> str:
        """确定整体风格"""
        dominant_layout = max(layout_distribution.items(), key=lambda x: x[1])[0]
        
        style_mapping = {
            'title_with_bullets': '商务演示风格',
            'title_with_content': '教学报告风格',
            'content_grid': '信息展示风格',
            'simple': '简约设计风格',
            'complex': '复合多元风格'
        }
        
        return style_mapping.get(dominant_layout, '混合风格')
    
    def _calculate_design_consistency(self) -> float:
        """计算设计一致性"""
        if len(self.slide_layouts) <= 1:
            return 1.0
        
        # 计算布局类型的一致性
        layout_types = [layout.layout_type for layout in self.slide_layouts]
        unique_types = len(set(layout_types))
        total_slides = len(layout_types)
        
        # 一致性 = 1 - 类型多样性
        consistency = 1.0 - (unique_types - 1) / (total_slides - 1) if total_slides > 1 else 1.0
        return max(consistency, 0.0)
    
    def _generate_analysis_summary(self) -> str:
        """生成分析摘要"""
        structure = self.overall_structure
        
        summary = f"""
PPT结构分析摘要：
• 总计 {structure['total_slides']} 张幻灯片
• 整体风格：{structure['overall_style']}
• 设计一致性：{structure['design_consistency']:.2f}
• 平均内容密度：{structure['average_metrics']['content_density']:.2f}
• 平均视觉平衡度：{structure['average_metrics']['visual_balance']:.2f}
• 平均层次清晰度：{structure['average_metrics']['hierarchy_clarity']:.2f}

布局类型分布：
"""
        
        for layout_type, count in structure['layout_distribution'].items():
            summary += f"• {layout_type}: {count} 张幻灯片\n"
        
        return summary.strip()

class PositionExtractor:
    """位置信息提取器"""
    
    def __init__(self, presentation: Presentation):
        self.presentation = presentation
        self.logger = get_logger()
    
    def extract_all_positions(self) -> Dict[str, Any]:
        """提取所有位置信息"""
        self.logger.info("开始提取位置信息")
        
        positions = {}
        
        for i, slide in enumerate(self.presentation.slides):
            slide_positions = self._extract_slide_positions(slide, i)
            positions[f'slide_{i}'] = slide_positions
        
        return {
            'slide_positions': positions,
            'spatial_relationships': self._analyze_spatial_relationships(positions)
        }
    
    def _extract_slide_positions(self, slide, slide_index: int) -> Dict[str, Any]:
        """提取单张幻灯片的位置信息"""
        elements = []
        
        for j, shape in enumerate(slide.shapes):
            element_pos = self._extract_element_position(shape, f"slide_{slide_index}_element_{j}")
            if element_pos:
                elements.append(element_pos)
        
        return {
            'elements': elements,
            'bounding_box': self._calculate_content_bounding_box(elements),
            'spatial_layout': self._analyze_spatial_layout(elements)
        }
    
    def _extract_element_position(self, shape, element_id: str) -> Optional[Dict[str, Any]]:
        """提取元素位置信息"""
        try:
            position = {
                'element_id': element_id,
                'left': float(shape.left),
                'top': float(shape.top),
                'width': float(shape.width),
                'height': float(shape.height),
                'center_x': float(shape.left + shape.width / 2),
                'center_y': float(shape.top + shape.height / 2),
                'right': float(shape.left + shape.width),
                'bottom': float(shape.top + shape.height),
                'area': float(shape.width * shape.height)
            }
            
            # 添加内容信息
            if hasattr(shape, 'text'):
                position['content'] = shape.text.strip()
                position['has_content'] = bool(shape.text.strip())
            
            return position
            
        except Exception as e:
            self.logger.error(f"提取元素位置信息时出错: {e}")
            return None
    
    def _calculate_content_bounding_box(self, elements: List[Dict[str, Any]]) -> Dict[str, float]:
        """计算内容边界框"""
        if not elements:
            return {'left': 0, 'top': 0, 'right': 0, 'bottom': 0, 'width': 0, 'height': 0}
        
        left = min(elem['left'] for elem in elements)
        top = min(elem['top'] for elem in elements)
        right = max(elem['right'] for elem in elements)
        bottom = max(elem['bottom'] for elem in elements)
        
        return {
            'left': left,
            'top': top,
            'right': right,
            'bottom': bottom,
            'width': right - left,
            'height': bottom - top
        }
    
    def _analyze_spatial_layout(self, elements: List[Dict[str, Any]]) -> Dict[str, Any]:
        """分析空间布局"""
        if not elements:
            return {}
        
        # 分析排列方式
        arrangement = self._detect_arrangement_pattern(elements)
        
        # 分析对齐方式
        alignment = self._detect_alignment_pattern(elements)
        
        # 分析间距
        spacing = self._analyze_spacing(elements)
        
        return {
            'arrangement': arrangement,
            'alignment': alignment,
            'spacing': spacing,
            'element_count': len(elements)
        }
    
    def _detect_arrangement_pattern(self, elements: List[Dict[str, Any]]) -> str:
        """检测排列模式"""
        if len(elements) <= 1:
            return "single"
        
        # 按位置排序
        sorted_by_y = sorted(elements, key=lambda x: x['top'])
        sorted_by_x = sorted(elements, key=lambda x: x['left'])
        
        # 检测是否为行排列
        y_variance = self._calculate_position_variance([elem['center_y'] for elem in elements])
        x_variance = self._calculate_position_variance([elem['center_x'] for elem in elements])
        
        if y_variance < x_variance * 0.5:
            return "horizontal"
        elif x_variance < y_variance * 0.5:
            return "vertical"
        else:
            return "grid"
    
    def _detect_alignment_pattern(self, elements: List[Dict[str, Any]]) -> Dict[str, Any]:
        """检测对齐模式"""
        if len(elements) <= 1:
            return {}
        
        # 检测左对齐
        left_positions = [elem['left'] for elem in elements]
        left_aligned = self._is_aligned(left_positions)
        
        # 检测右对齐
        right_positions = [elem['right'] for elem in elements]
        right_aligned = self._is_aligned(right_positions)
        
        # 检测中心对齐
        center_x_positions = [elem['center_x'] for elem in elements]
        center_x_aligned = self._is_aligned(center_x_positions)
        
        center_y_positions = [elem['center_y'] for elem in elements]
        center_y_aligned = self._is_aligned(center_y_positions)
        
        return {
            'left_aligned': left_aligned,
            'right_aligned': right_aligned,
            'center_x_aligned': center_x_aligned,
            'center_y_aligned': center_y_aligned
        }
    
    def _is_aligned(self, positions: List[float], tolerance: float = 1000) -> bool:
        """判断位置是否对齐"""
        if len(positions) <= 1:
            return True
        
        variance = self._calculate_position_variance(positions)
        return variance < tolerance
    
    def _calculate_position_variance(self, positions: List[float]) -> float:
        """计算位置方差"""
        if len(positions) <= 1:
            return 0.0
        
        mean = sum(positions) / len(positions)
        variance = sum((pos - mean) ** 2 for pos in positions) / len(positions)
        return variance
    
    def _analyze_spacing(self, elements: List[Dict[str, Any]]) -> Dict[str, Any]:
        """分析间距"""
        if len(elements) <= 1:
            return {}
        
        # 计算水平间距
        horizontal_gaps = []
        elements_by_x = sorted(elements, key=lambda x: x['left'])
        for i in range(len(elements_by_x) - 1):
            gap = elements_by_x[i + 1]['left'] - elements_by_x[i]['right']
            if gap > 0:
                horizontal_gaps.append(gap)
        
        # 计算垂直间距
        vertical_gaps = []
        elements_by_y = sorted(elements, key=lambda x: x['top'])
        for i in range(len(elements_by_y) - 1):
            gap = elements_by_y[i + 1]['top'] - elements_by_y[i]['bottom']
            if gap > 0:
                vertical_gaps.append(gap)
        
        return {
            'horizontal_gaps': horizontal_gaps,
            'vertical_gaps': vertical_gaps,
            'avg_horizontal_gap': sum(horizontal_gaps) / len(horizontal_gaps) if horizontal_gaps else 0,
            'avg_vertical_gap': sum(vertical_gaps) / len(vertical_gaps) if vertical_gaps else 0
        }
    
    def _analyze_spatial_relationships(self, positions: Dict[str, Any]) -> Dict[str, Any]:
        """分析空间关系"""
        relationships = {}
        
        # 分析跨幻灯片的一致性
        slide_keys = list(positions.keys())
        
        for i, slide_key in enumerate(slide_keys):
            slide_data = positions[slide_key]
            
            # 分析与其他幻灯片的关系
            for j, other_slide_key in enumerate(slide_keys):
                if i != j:
                    other_slide_data = positions[other_slide_key]
                    relationship = self._compare_slide_layouts(slide_data, other_slide_data)
                    relationships[f"{slide_key}_vs_{other_slide_key}"] = relationship
        
        return relationships
    
    def _compare_slide_layouts(self, slide1: Dict[str, Any], slide2: Dict[str, Any]) -> Dict[str, Any]:
        """比较两张幻灯片的布局"""
        elements1 = slide1.get('elements', [])
        elements2 = slide2.get('elements', [])
        
        # 比较元素数量
        count_similarity = 1.0 - abs(len(elements1) - len(elements2)) / max(len(elements1), len(elements2), 1)
        
        # 比较空间布局
        layout1 = slide1.get('spatial_layout', {})
        layout2 = slide2.get('spatial_layout', {})
        
        arrangement_similarity = 1.0 if layout1.get('arrangement') == layout2.get('arrangement') else 0.0
        
        return {
            'count_similarity': count_similarity,
            'arrangement_similarity': arrangement_similarity,
            'overall_similarity': (count_similarity + arrangement_similarity) / 2
        }

class SmartLayoutAdjuster:
    """智能布局调整器"""
    
    def __init__(self, presentation: Presentation, structure_analyzer: PPTStructureAnalyzer, position_extractor: PositionExtractor):
        self.presentation = presentation
        self.structure_analyzer = structure_analyzer
        self.position_extractor = position_extractor
        self.logger = get_logger()
        
    def suggest_optimal_layout(self, slide_index: int, new_content: Dict[str, str]) -> Dict[str, Any]:
        """为新内容建议最优布局"""
        self.logger.info(f"为第{slide_index+1}张幻灯片建议最优布局")
        
        # 获取当前幻灯片结构
        current_layout = self.structure_analyzer.slide_layouts[slide_index]
        
        # 分析新内容特征
        content_analysis = self._analyze_content_characteristics(new_content)
        
        # 生成布局建议
        layout_suggestions = self._generate_layout_suggestions(current_layout, content_analysis)
        
        return {
            'current_layout': current_layout,
            'content_analysis': content_analysis,
            'layout_suggestions': layout_suggestions,
            'recommended_adjustments': self._recommend_adjustments(current_layout, content_analysis)
        }
    
    def _analyze_content_characteristics(self, content: Dict[str, str]) -> Dict[str, Any]:
        """分析内容特征"""
        characteristics = {
            'content_items': [],
            'total_length': 0,
            'avg_length': 0,
            'has_title': False,
            'has_bullets': False,
            'content_density': 0,
            'hierarchy_levels': 0
        }
        
        total_length = 0
        lengths = []
        
        for placeholder, text in content.items():
            item_info = {
                'placeholder': placeholder,
                'text': text,
                'length': len(text),
                'type': self._classify_content_type(placeholder, text)
            }
            characteristics['content_items'].append(item_info)
            
            total_length += len(text)
            lengths.append(len(text))
            
            if 'title' in placeholder.lower():
                characteristics['has_title'] = True
            if 'bullet' in placeholder.lower():
                characteristics['has_bullets'] = True
        
        characteristics['total_length'] = total_length
        characteristics['avg_length'] = sum(lengths) / len(lengths) if lengths else 0
        characteristics['content_density'] = self._calculate_content_density_score(lengths)
        characteristics['hierarchy_levels'] = self._count_hierarchy_levels(content)
        
        return characteristics
    
    def _classify_content_type(self, placeholder: str, text: str) -> str:
        """分类内容类型"""
        placeholder_lower = placeholder.lower()
        
        if 'title' in placeholder_lower:
            return 'title'
        elif 'subtitle' in placeholder_lower:
            return 'subtitle'
        elif 'bullet' in placeholder_lower:
            return 'bullet'
        elif 'description' in placeholder_lower:
            return 'description'
        elif 'conclusion' in placeholder_lower:
            return 'conclusion'
        elif len(text) < 20:
            return 'short_text'
        elif len(text) > 100:
            return 'long_text'
        else:
            return 'medium_text'
    
    def _calculate_content_density_score(self, lengths: List[int]) -> float:
        """计算内容密度分数"""
        if not lengths:
            return 0.0
        
        total_chars = sum(lengths)
        # 假设每个字符占用一定的显示空间
        # 这里简化计算，实际应该考虑字体大小、行距等因素
        density_score = min(total_chars / 1000, 1.0)  # 1000字符为满密度
        return density_score
    
    def _count_hierarchy_levels(self, content: Dict[str, str]) -> int:
        """计算层次级别数"""
        levels = set()
        
        for placeholder in content.keys():
            if 'title' in placeholder.lower():
                levels.add(1)
            elif 'subtitle' in placeholder.lower():
                levels.add(2)
            elif 'content' in placeholder.lower():
                levels.add(3)
            elif 'bullet' in placeholder.lower():
                levels.add(4)
            elif 'description' in placeholder.lower():
                levels.add(5)
        
        return len(levels)
    
    def _generate_layout_suggestions(self, current_layout: SlideLayout, content_analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
        """生成布局建议"""
        suggestions = []
        
        # 建议1: 基于内容密度的布局调整
        if content_analysis['content_density'] > 0.7:
            suggestions.append({
                'type': 'reduce_density',
                'description': '内容密度过高，建议分拆到多张幻灯片或使用更大的字体',
                'priority': 'high',
                'adjustments': ['increase_font_size', 'add_more_whitespace', 'split_content']
            })
        
        # 建议2: 基于层次结构的布局调整
        if content_analysis['hierarchy_levels'] > 3:
            suggestions.append({
                'type': 'improve_hierarchy',
                'description': '层次结构复杂，建议优化视觉层次',
                'priority': 'medium',
                'adjustments': ['adjust_font_sizes', 'use_different_colors', 'improve_spacing']
            })
        
        # 建议3: 基于现有布局的优化
        if current_layout.visual_balance < 0.6:
            suggestions.append({
                'type': 'improve_balance',
                'description': '视觉平衡度不足，建议调整元素位置',
                'priority': 'high',
                'adjustments': ['reposition_elements', 'adjust_element_sizes', 'add_visual_anchors']
            })
        
        return suggestions
    
    def _recommend_adjustments(self, current_layout: SlideLayout, content_analysis: Dict[str, Any]) -> Dict[str, Any]:
        """推荐具体调整"""
        adjustments = {
            'font_size_adjustments': {},
            'position_adjustments': {},
            'spacing_adjustments': {},
            'visual_improvements': []
        }
        
        # 字体大小调整
        for item in content_analysis['content_items']:
            if item['type'] == 'title':
                adjustments['font_size_adjustments'][item['placeholder']] = {'min_size': 24, 'max_size': 36}
            elif item['type'] == 'bullet':
                adjustments['font_size_adjustments'][item['placeholder']] = {'min_size': 16, 'max_size': 20}
            elif item['type'] == 'description':
                adjustments['font_size_adjustments'][item['placeholder']] = {'min_size': 12, 'max_size': 16}
        
        # 视觉改进建议
        if current_layout.hierarchy_clarity < 0.7:
            adjustments['visual_improvements'].append("增强层次对比，使用不同的字体大小和颜色")
        
        if current_layout.content_density > 0.8:
            adjustments['visual_improvements'].append("增加留白空间，提高可读性")
        
        if current_layout.visual_balance < 0.6:
            adjustments['visual_improvements'].append("调整元素位置，改善视觉平衡")
        
        return adjustments
    
    def apply_smart_adjustments(self, slide_index: int, adjustments: Dict[str, Any]) -> Dict[str, Any]:
        """应用智能调整"""
        self.logger.info(f"为第{slide_index+1}张幻灯片应用智能调整")
        
        slide = self.presentation.slides[slide_index]
        applied_adjustments = []
        
        try:
            # 应用字体大小调整
            font_adjustments = adjustments.get('font_size_adjustments', {})
            for placeholder_name, size_info in font_adjustments.items():
                success = self._apply_font_size_adjustment(slide, placeholder_name, size_info)
                if success:
                    applied_adjustments.append(f"调整了{placeholder_name}的字体大小")
            
            # 应用位置调整
            position_adjustments = adjustments.get('position_adjustments', {})
            for placeholder_name, position_info in position_adjustments.items():
                success = self._apply_position_adjustment(slide, placeholder_name, position_info)
                if success:
                    applied_adjustments.append(f"调整了{placeholder_name}的位置")
            
            # 应用间距调整
            spacing_adjustments = adjustments.get('spacing_adjustments', {})
            if spacing_adjustments:
                success = self._apply_spacing_adjustments(slide, spacing_adjustments)
                if success:
                    applied_adjustments.append("调整了元素间距")
            
            return {
                'success': True,
                'applied_adjustments': applied_adjustments,
                'message': f"成功应用了{len(applied_adjustments)}项调整"
            }
            
        except Exception as e:
            self.logger.error(f"应用智能调整时出错: {e}")
            return {
                'success': False,
                'error': str(e),
                'applied_adjustments': applied_adjustments
            }
    
    def _apply_font_size_adjustment(self, slide, placeholder_name: str, size_info: Dict[str, Any]) -> bool:
        """应用字体大小调整"""
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    # 检查是否包含目标占位符 - 支持所有{}格式占位符
                    if f"{{{placeholder_name}}}" in shape.text:
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            # 计算建议的字体大小
                            content_length = len(shape.text)
                            if content_length < 20:
                                font_size = size_info.get('max_size', 18)
                            elif content_length > 100:
                                font_size = size_info.get('min_size', 12)
                            else:
                                # 根据内容长度线性插值
                                min_size = size_info.get('min_size', 12)
                                max_size = size_info.get('max_size', 18)
                                font_size = max_size - (content_length - 20) * (max_size - min_size) / 80
                            
                            # 应用字体大小
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(font_size)
                            
                            return True
            return False
        except Exception as e:
            self.logger.error(f"应用字体大小调整失败: {e}")
            return False
    
    def _apply_position_adjustment(self, slide, placeholder_name: str, position_info: Dict[str, Any]) -> bool:
        """应用位置调整"""
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    if f"{{{placeholder_name}}}" in shape.text:
                        # 应用位置调整
                        if 'left' in position_info:
                            shape.left = Inches(position_info['left'])
                        if 'top' in position_info:
                            shape.top = Inches(position_info['top'])
                        if 'width' in position_info:
                            shape.width = Inches(position_info['width'])
                        if 'height' in position_info:
                            shape.height = Inches(position_info['height'])
                        return True
            return False
        except Exception as e:
            self.logger.error(f"应用位置调整失败: {e}")
            return False
    
    def _apply_spacing_adjustments(self, slide, spacing_info: Dict[str, Any]) -> bool:
        """应用间距调整"""
        try:
            # 这里可以实现更复杂的间距调整逻辑
            # 例如重新排列所有元素以改善间距
            return True
        except Exception as e:
            self.logger.error(f"应用间距调整失败: {e}")
            return False

def create_advanced_ppt_analysis(presentation: Presentation) -> Dict[str, Any]:
    """创建高级PPT分析"""
    logger = get_logger()
    logger.info("开始创建高级PPT分析")
    
    try:
        # 创建分析器
        structure_analyzer = PPTStructureAnalyzer(presentation)
        position_extractor = PositionExtractor(presentation)
        layout_adjuster = SmartLayoutAdjuster(presentation, structure_analyzer, position_extractor)
        
        # 执行分析
        structure_analysis = structure_analyzer.analyze_complete_structure()
        position_analysis = position_extractor.extract_all_positions()
        
        return {
            'structure_analysis': structure_analysis,
            'position_analysis': position_analysis,
            'analyzers': {
                'structure_analyzer': structure_analyzer,
                'position_extractor': position_extractor,
                'layout_adjuster': layout_adjuster
            }
        }
        
    except Exception as e:
        logger.error(f"创建高级PPT分析时出错: {e}")
        return {
            'error': str(e),
            'structure_analysis': None,
            'position_analysis': None
        }