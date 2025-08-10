#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
格式基准合并器
以split_presentations_1.pptx的格式作为整个PPT的设计基准
"""

import os
import io
import sys
from typing import List, Dict, Any, Optional
from pptx import Presentation
from logger import get_logger, log_user_action

logger = get_logger()

# 导入其他合并器
try:
    from ppt_merger_win32 import WIN32_AVAILABLE
    if WIN32_AVAILABLE:
        import win32com.client
except ImportError:
    WIN32_AVAILABLE = False

class FormatBaseMerger:
    """基于split_presentations_1.pptx格式的合并器"""
    
    def __init__(self):
        """初始化合并器"""
        self.base_format_template = os.path.join("templates", "ppt_template", "split_presentations_1.pptx")
        self.ppt_app = None
        self.merged_presentation = None
        self.temp_presentations = []
        
    def __enter__(self):
        """上下文管理器入口"""
        if WIN32_AVAILABLE and sys.platform.startswith('win'):
            self._start_powerpoint()
        return self
        
    def __exit__(self, exc_type, exc_val, exc_tb):
        """上下文管理器退出，清理资源"""
        if WIN32_AVAILABLE and sys.platform.startswith('win'):
            self._cleanup()
    
    def _start_powerpoint(self):
        """启动PowerPoint应用程序"""
        try:
            import pythoncom
            pythoncom.CoInitialize()
            
            self.ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            logger.info("PowerPoint COM接口初始化成功")
            
        except Exception as e:
            logger.error(f"PowerPoint启动失败: {str(e)}")
            raise
    
    def merge_with_format_base(self, page_results: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        使用split_presentations_1.pptx的格式作为基准合并所有页面
        
        Args:
            page_results: 页面处理结果列表
            
        Returns:
            Dict: 整合结果
        """
        log_user_action("格式基准合并", f"以split_presentations_1.pptx格式为基准，合并{len(page_results)}个页面")
        
        result = {
            "success": False,
            "total_pages": 0,
            "processed_pages": 0,
            "skipped_pages": 0,
            "errors": [],
            "presentation_bytes": None,
            "output_path": None,
            "error": None,
            "format_base": "split_presentations_1.pptx"
        }
        
        if not page_results:
            result["error"] = "没有页面需要合并"
            return result
        
        # 验证格式基准模板是否存在
        if not os.path.exists(self.base_format_template):
            result["error"] = f"格式基准模板不存在: {self.base_format_template}"
            return result
        
        try:
            if WIN32_AVAILABLE and sys.platform.startswith('win'):
                return self._merge_with_win32com(page_results, result)
            else:
                return self._merge_with_python_pptx(page_results, result)
                
        except Exception as e:
            result["error"] = f"格式基准合并异常: {str(e)}"
            logger.error(f"格式基准合并异常: {str(e)}")
            return result
    
    def _merge_with_win32com(self, page_results: List[Dict[str, Any]], result: Dict[str, Any]) -> Dict[str, Any]:
        """使用Win32COM进行格式基准合并"""
        try:
            # 以split_presentations_1.pptx作为基础演示文稿
            abs_base_path = os.path.abspath(self.base_format_template)
            self.merged_presentation = self.ppt_app.Presentations.Open(abs_base_path)
            
            logger.info(f"使用格式基准模板: {os.path.basename(self.base_format_template)}")
            
            # 删除基准模板的内容页，只保留格式框架
            if self.merged_presentation.Slides.Count > 0:
                self.merged_presentation.Slides(1).Delete()
            
            result["processed_pages"] = 0
            
            # 逐个添加每个页面，但应用基准格式
            for i, page_result in enumerate(page_results):
                template_path = page_result.get('template_path')
                
                if not template_path or not os.path.exists(template_path):
                    result["errors"].append(f"第{i+1}页模板文件不存在")
                    result["skipped_pages"] += 1
                    continue
                
                try:
                    abs_path = os.path.abspath(template_path)
                    
                    # 打开源模板
                    temp_ppt = self.ppt_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)
                    self.temp_presentations.append(temp_ppt)
                    
                    if temp_ppt.Slides.Count > 0:
                        # 复制源页面内容
                        source_slide = temp_ppt.Slides(1)
                        source_slide.Copy()
                        
                        # 粘贴到基准演示文稿中
                        slide_index = self.merged_presentation.Slides.Count + 1
                        self.merged_presentation.Slides.Paste(slide_index)
                        
                        # 获取刚粘贴的幻灯片
                        pasted_slide = self.merged_presentation.Slides(slide_index)
                        
                        # 应用基准格式到粘贴的幻灯片
                        self._apply_base_format_to_slide(pasted_slide)
                        
                        result["processed_pages"] += 1
                        logger.info(f"成功合并并应用基准格式: {os.path.basename(template_path)}")
                    else:
                        result["skipped_pages"] += 1
                        result["errors"].append(f"第{i+1}页模板文件为空")
                    
                    # 关闭临时演示文稿
                    temp_ppt.Close()
                    
                except Exception as e:
                    error_msg = f"第{i+1}页格式基准合并失败: {str(e)}"
                    result["errors"].append(error_msg)
                    result["skipped_pages"] += 1
                    logger.error(error_msg)
            
            # 保存合并结果
            output_path = self._save_presentation()
            if output_path:
                result["output_path"] = output_path
                result["success"] = True
                result["total_pages"] = self.merged_presentation.Slides.Count
                
                # 读取文件字节数据
                try:
                    with open(output_path, 'rb') as f:
                        result["presentation_bytes"] = f.read()
                except Exception as read_e:
                    logger.warning(f"读取PPT字节数据失败: {str(read_e)}")
                
                log_user_action("格式基准合并完成", 
                               f"成功: {result['processed_pages']}页, 跳过: {result['skipped_pages']}页")
            else:
                result["error"] = "PPT文件保存失败"
                
        except Exception as e:
            result["error"] = f"Win32COM格式基准合并异常: {str(e)}"
            logger.error(f"Win32COM格式基准合并异常: {str(e)}")
        
        return result
    
    def _apply_base_format_to_slide(self, slide):
        """
        将基准格式应用到幻灯片
        这个方法尝试将split_presentations_1.pptx的主题设计应用到当前幻灯片
        """
        try:
            # 由于我们已经基于split_presentations_1.pptx创建演示文稿
            # PowerPoint会自动继承其主题和设计
            # 这里可以添加额外的格式调整逻辑
            
            # 尝试应用基准演示文稿的设计模板
            if hasattr(slide, 'ApplyTemplate'):
                try:
                    # 应用基准模板的设计
                    slide.ApplyTemplate(self.base_format_template)
                except:
                    pass
            
            logger.debug(f"已应用基准格式到幻灯片")
            
        except Exception as e:
            logger.warning(f"应用基准格式失败: {str(e)}")
    
    def _merge_with_python_pptx(self, page_results: List[Dict[str, Any]], result: Dict[str, Any]) -> Dict[str, Any]:
        """使用python-pptx进行格式基准合并（备用方案）"""
        try:
            # 加载基准模板
            base_ppt = Presentation(self.base_format_template)
            
            # 清空基准模板的幻灯片，但保留主题设计
            while len(base_ppt.slides) > 0:
                slides = base_ppt.slides
                slides._sldIdLst.remove(slides._sldIdLst[0])
            
            # 逐个添加页面
            for i, page_result in enumerate(page_results):
                template_path = page_result.get('template_path')
                
                if not template_path or not os.path.exists(template_path):
                    result["errors"].append(f"第{i+1}页模板文件不存在")
                    result["skipped_pages"] += 1
                    continue
                
                try:
                    # 加载源模板
                    source_ppt = Presentation(template_path)
                    if len(source_ppt.slides) == 0:
                        result["skipped_pages"] += 1
                        continue
                    
                    source_slide = source_ppt.slides[0]
                    
                    # 在基准PPT中创建新幻灯片
                    slide_layout = base_ppt.slide_layouts[6]  # 空白布局
                    new_slide = base_ppt.slides.add_slide(slide_layout)
                    
                    # 复制内容但保持基准格式
                    self._copy_content_with_base_format(source_slide, new_slide)
                    
                    result["processed_pages"] += 1
                    logger.info(f"成功合并(python-pptx): {os.path.basename(template_path)}")
                    
                except Exception as e:
                    error_msg = f"第{i+1}页python-pptx合并失败: {str(e)}"
                    result["errors"].append(error_msg)
                    result["skipped_pages"] += 1
                    logger.error(error_msg)
            
            # 生成PPT字节数据
            ppt_stream = io.BytesIO()
            base_ppt.save(ppt_stream)
            result["presentation_bytes"] = ppt_stream.getvalue()
            ppt_stream.close()
            
            result["success"] = True
            result["total_pages"] = len(base_ppt.slides)
            
            log_user_action("格式基准合并完成(python-pptx)", 
                           f"成功: {result['processed_pages']}页, 跳过: {result['skipped_pages']}页")
            
        except Exception as e:
            result["error"] = f"python-pptx格式基准合并异常: {str(e)}"
            logger.error(f"python-pptx格式基准合并异常: {str(e)}")
        
        return result
    
    def _copy_content_with_base_format(self, source_slide, target_slide):
        """复制内容但保持基准格式"""
        try:
            # 简单的内容复制，保持目标的格式基准
            for shape in source_slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    try:
                        textbox = target_slide.shapes.add_textbox(
                            shape.left, shape.top, shape.width, shape.height
                        )
                        textbox.text = shape.text
                        
                        # 不复制源格式，让其使用基准PPT的默认格式
                        
                    except Exception as text_e:
                        logger.warning(f"文本复制失败: {str(text_e)}")
                        
        except Exception as e:
            logger.warning(f"内容复制失败: {str(e)}")
    
    def _save_presentation(self) -> Optional[str]:
        """保存演示文稿到文件"""
        try:
            import tempfile
            import datetime
            
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_dir = tempfile.gettempdir()
            output_filename = f"format_base_merged_{timestamp}.pptx"
            output_path = os.path.join(output_dir, output_filename)
            
            # 保存为PowerPoint格式
            self.merged_presentation.SaveAs(output_path, 24)
            
            logger.info(f"格式基准PPT文件已保存: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"保存格式基准PPT文件失败: {str(e)}")
            return None
    
    def _cleanup(self):
        """清理资源"""
        try:
            # 关闭所有临时打开的演示文稿
            for temp_ppt in self.temp_presentations:
                try:
                    temp_ppt.Close()
                except:
                    pass
            
            # 关闭合并的演示文稿
            if self.merged_presentation:
                try:
                    self.merged_presentation.Close()
                except:
                    pass
            
            # 退出PowerPoint应用程序
            if self.ppt_app:
                try:
                    self.ppt_app.Quit()
                except:
                    pass
            
            # 清理COM
            try:
                import pythoncom
                pythoncom.CoUninitialize()
            except:
                pass
                    
            logger.info("格式基准合并器资源已清理")
            
        except Exception as e:
            logger.warning(f"格式基准合并器资源清理异常: {str(e)}")

def merge_with_split_presentations_1_format(page_results: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    便捷函数：使用split_presentations_1.pptx的格式作为基准合并PPT
    
    Args:
        page_results: 页面处理结果列表
        
    Returns:
        Dict: 整合结果
    """
    try:
        with FormatBaseMerger() as merger:
            return merger.merge_with_format_base(page_results)
    except Exception as e:
        return {
            "success": False,
            "error": f"格式基准合并失败: {str(e)}"
        }

if __name__ == "__main__":
    # 测试用例
    test_results = [
        {
            'page_number': 1,
            'template_number': 'title',
            'template_path': os.path.join("templates", "title_slides.pptx"),
            'template_filename': "title_slides.pptx"
        },
        {
            'page_number': 2,
            'template_number': 5,
            'template_path': os.path.join("templates", "ppt_template", "split_presentations_5.pptx"),
            'template_filename': "split_presentations_5.pptx"
        }
    ]
    
    result = merge_with_split_presentations_1_format(test_results)
    
    if result["success"]:
        print("✅ 格式基准合并成功！")
        print(f"总页数: {result['total_pages']}")
        print(f"处理页数: {result['processed_pages']}")
        print(f"跳过页数: {result['skipped_pages']}")
        print(f"格式基准: {result['format_base']}")
        
        if result["presentation_bytes"]:
            with open("test_format_base_merger.pptx", "wb") as f:
                f.write(result["presentation_bytes"])
            print("📄 测试PPT已保存为: test_format_base_merger.pptx")
            
        if result.get("output_path"):
            print(f"输出文件: {result['output_path']}")
    else:
        print(f"❌ 格式基准合并失败: {result['error']}")
        if result.get("errors"):
            for error in result["errors"]:
                print(f"  - {error}")
