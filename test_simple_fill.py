#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
简单的文本填充测试
"""

import os
from pptx import Presentation
from config import get_config
from utils import PPTAnalyzer, PPTProcessor

def test_simple_fill():
    print("简单文本填充测试")
    print("=" * 40)
    
    config = get_config()
    ppt_path = config.default_ppt_template
    
    if not os.path.exists(ppt_path):
        print(f"PPT文件不存在: {ppt_path}")
        return
    
    # 加载PPT
    presentation = Presentation(ppt_path)
    ppt_processor = PPTProcessor(presentation)
    
    # 手动创建一个简单的分配方案进行测试
    test_assignments = {
        "assignments": [
            {
                "slide_index": 0,
                "action": "replace_placeholder",
                "placeholder": "title",
                "content": "测试标题",
                "reason": "手动测试"
            },
            {
                "slide_index": 0,
                "action": "replace_placeholder", 
                "placeholder": "content_1",
                "content": "测试内容1",
                "reason": "手动测试"
            }
        ]
    }
    
    print("测试分配方案:")
    for assignment in test_assignments["assignments"]:
        print(f"  {assignment['placeholder']} -> '{assignment['content']}'")
    
    print("\n执行文本填充...")
    
    # 应用分配方案（不进行美化）
    results = ppt_processor.apply_assignments(test_assignments)
    
    print("\n填充结果:")
    for result in results:
        print(f"  {result}")
    
    # 检查填充后的实际内容
    print("\n验证填充结果:")
    for i, slide in enumerate(presentation.slides):
        print(f"\n第{i+1}页当前内容:")
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                print(f"  文本框{j+1}: '{text}'")
    
    # 保存测试结果
    output_path = "output/simple_fill_test.pptx"
    os.makedirs("output", exist_ok=True)
    presentation.save(output_path)
    print(f"\n已保存测试结果: {output_path}")

if __name__ == "__main__":
    test_simple_fill()