#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文本转PPT填充器 - Streamlit Web界面
使用OpenAI GPT-4V将文本填入现有PPT文件
"""

import streamlit as st
import os
from datetime import datetime
from typing import TYPE_CHECKING
from pptx.util import Inches, Pt
import json
import re

if TYPE_CHECKING:
    from pptx.presentation import Presentation
else:
    from pptx import Presentation
from config import get_config
from utils import AIProcessor, PPTProcessor, FileManager, PPTAnalyzer
from logger import get_logger, log_user_action, log_file_operation, LogContext

# 获取配置
config = get_config()
logger = get_logger()

# 页面配置
st.set_page_config(
    page_title=config.web_title,
    page_icon=config.web_icon,
    layout=config.web_layout if config.web_layout in ("centered", "wide") else "centered",
    initial_sidebar_state="expanded"
)

# 自定义CSS样式
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.2rem;
        color: #666;
        text-align: center;
        margin-bottom: 2rem;
    }
    .success-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        color: #155724;
    }
    .info-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #d1ecf1;
        border: 1px solid #b6d4ea;
        color: #0c5460;
    }
    .warning-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #fff3cd;
        border: 1px solid #ffeaa7;
        color: #856404;
    }
    .error-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #f8d7da;
        border: 1px solid #f5c6cb;
        color: #721c24;
    }
</style>
""", unsafe_allow_html=True)

class StreamlitPPTGenerator:
    def __init__(self, api_key):
        """初始化生成器"""
        self.api_key = api_key
        self.ai_processor = AIProcessor(api_key)
        self.presentation = None
        self.ppt_processor = None
        self.ppt_structure = None
        logger.info(f"初始化PPT生成器，API密钥: {'已设置' if api_key else '未设置'}")
    
    def load_ppt_from_path(self, ppt_path):
        """从文件路径加载PPT"""
        with LogContext(f"加载PPT文件: {ppt_path}"):
            try:
                # 验证文件
                is_valid, error_msg = FileManager.validate_ppt_file(ppt_path)
                if not is_valid:
                    st.error(f"PPT文件验证失败: {error_msg}")
                    log_file_operation("load_ppt", ppt_path, "error", error_msg)
                    return False
                
                self.presentation = Presentation(ppt_path)  # type: ignore
                self.ppt_processor = PPTProcessor(self.presentation)
                self.ppt_structure = self.ppt_processor.ppt_structure
                
                log_file_operation("load_ppt", ppt_path, "success")
                return True
            except Exception as e:
                st.error(f"加载PPT文件失败: {e}")
                log_file_operation("load_ppt", ppt_path, "error", str(e))
                return False
    
    
    def process_text_with_deepseek(self, user_text):
        """使用OpenAI API分析如何将用户文本填入PPT模板的占位符"""
        if not self.ppt_structure or not self.ppt_processor:
            return {"assignments": []}
        
        log_user_action("AI文本分析", f"文本长度: {len(user_text)}字符")
        
        # 获取增强的结构信息
        enhanced_info = self.ppt_processor.get_enhanced_structure_info()
        
        # 使用增强信息进行分析
        return self.ai_processor.analyze_text_for_ppt(user_text, self.ppt_structure, enhanced_info)
    
    def apply_text_assignments(self, assignments, user_text: str = ""):
        """根据分配方案替换PPT模板中的占位符，并将原始文本添加到备注"""
        if not self.presentation or not self.ppt_processor:
            return ["❌ PPT文件未正确加载"]
        
        log_user_action("应用文本分配", f"分配数量: {len(assignments.get('assignments', []))}")
        # 传递用户原始文本，用于添加到幻灯片备注
        results = self.ppt_processor.apply_assignments(assignments, user_text)
        
        # 美化演示文稿
        st.info("正在美化PPT布局...")
        beautify_results = self.ppt_processor.beautify_presentation()
        
        # 显示美化结果
        summary = beautify_results['summary']
        st.success("PPT美化完成！")
        
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("删除占位符", summary['removed_placeholders_count'])
        with col2:
            st.metric("重新排版", summary['reorganized_slides_count'])
        with col3:
            st.metric("删除空页", summary['removed_empty_slides_count'])
        with col4:
            st.metric("最终页数", summary['final_slide_count'])
        
        # 显示详细结果
        if summary['removed_placeholders_count'] > 0:
            with st.expander("🧹 查看清理详情"):
                for item in beautify_results['beautify_results']['removed_placeholders']:
                    st.write(f"• 第{item['slide_index']+1}页: 删除了 {item['removed_count']} 个未填充占位符")
                    for placeholder in item['removed_placeholders']:
                        st.write(f"  - {{{placeholder}}}")
        
        if summary['reorganized_slides_count'] > 0:
            with st.expander("🎨 查看重排版详情"):
                for item in beautify_results['beautify_results']['reorganized_slides']:
                    layout_change = item['layout_change']
                    st.write(f"• 第{item['slide_index']+1}页: 使用 {layout_change['layout_type']} 布局重新排版了 {layout_change['shape_count']} 个元素")
        
        return results
    
    
    
    
    def get_ppt_bytes(self):
        """获取修改后的PPT字节数据"""
        if not self.presentation:
            raise ValueError("PPT文件未正确加载")
        
        log_user_action("获取PPT字节数据")
        return FileManager.save_ppt_to_bytes(self.presentation)

def main():
    # 页面标题
    st.markdown('<div class="main-header">📊 文本转PPT填充器</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">使用OpenAI GPT-4V智能将您的文本填入预设PPT模板</div>', unsafe_allow_html=True)
    
    # 侧边栏配置
    with st.sidebar:
        st.header("⚙️ 配置")
        
        # 模型选择
        st.subheader("🤖 AI模型选择")
        
        available_models = config.available_models
        model_options = {}
        for model_key, model_info in available_models.items():
            display_name = f"{model_info['name']} ({model_info['cost']}成本)"
            if not model_info['supports_vision']:
                display_name += " - ⚠️ 无视觉分析"
            model_options[display_name] = model_key
        
        selected_display = st.selectbox(
            "选择AI模型",
            options=list(model_options.keys()),
            index=0,
            help="不同模型有不同的功能和成本"
        )
        
        selected_model = model_options[selected_display]
        model_info = available_models[selected_model]
        
        # 显示模型信息
        st.info(f"**{model_info['name']}**: {model_info['description']}")
        
        if not model_info['supports_vision']:
            st.warning("⚠️ 注意：此模型不支持视觉分析功能，将跟过PPT美观度分析步骤")
        
        # 动态更新配置
        if selected_model != config.ai_model:
            config.set_model(selected_model)
        
        st.markdown("---")
        
        # API密钥输入（根据选择的模型动态显示）
        current_model_info = config.get_model_info()
        api_provider = current_model_info.get('api_provider', 'OpenRouter')
        api_key_url = current_model_info.get('api_key_url', 'https://openrouter.ai/keys')
        
        api_key = st.text_input(
            f"{api_provider} API密钥",
            type="password",
            help=f"请输入您的{api_provider} API密钥",
            placeholder="sk-..." if api_provider == "OpenRouter" else "请输入API密钥"
        )
        
        if not api_key:
            st.markdown('<div class="warning-box">⚠️ 请先输入API密钥才能使用功能</div>', unsafe_allow_html=True)
            st.markdown(f"获取API密钥：[{api_provider}平台]({api_key_url})")
        else:
            # 验证API密钥格式
            if api_provider == "OpenRouter" and not api_key.startswith('sk-'):
                st.markdown('<div class="warning-box">⚠️ OpenRouter API密钥格式可能不正确，请确认是否以"sk-"开头</div>', unsafe_allow_html=True)
            elif api_provider == "DeepSeek" and not api_key.startswith('sk-'):
                st.markdown('<div class="warning-box">⚠️ DeepSeek API密钥格式可能不正确，请确认格式是否正确</div>', unsafe_allow_html=True)
        
        st.markdown("---")
        
        # 模板信息
        st.subheader("📄 PPT模板")
        st.markdown(f"**当前模板：** `{os.path.basename(config.default_ppt_template)}`")
        st.markdown(f"**模板路径：** `{config.default_ppt_template}`")
        
        # 检查模板文件状态
        is_valid, error_msg = FileManager.validate_ppt_file(config.default_ppt_template)
        if is_valid:
            st.markdown('<div class="success-box">✅ 模板文件存在</div>', unsafe_allow_html=True)
        else:
            st.markdown(f'<div class="error-box">❌ 模板文件问题: {error_msg}</div>', unsafe_allow_html=True)
        
        st.markdown("---")
        
        # 使用说明
        st.subheader("📖 使用说明")
        st.markdown("""
        1. **选择AI模型**：选择适合您需求的模型
           - GPT-4o：功能完整，支持视觉分析，成本较高
           - DeepSeek R1：成本较低，专注推理，但不支持视觉分析
        2. **输入API密钥**：根据选择的模型输入相应的API密钥
           - GPT-4o：需要OpenRouter API密钥
           - DeepSeek R1：需要DeepSeek API密钥
        3. **确认模板**：确保PPT模板文件存在
        4. **输入文本**：输入要填入PPT的文本内容
        5. **开始处理**：点击处理按钮
        6. **下载PPT**：下载更新后的PPT文件
        
        💡 **模型选择建议**：
        - 如果追求最佳效果且预算充足，选择GPT-4o
        - 如果预算有限或主要做推理处理，选择DeepSeek R1
        """)
    
    # 主界面 - 只有输入API密钥后才显示功能
    if api_key and api_key.strip():
        # 检查模板文件
        is_valid, error_msg = FileManager.validate_ppt_file(config.default_ppt_template)
        if not is_valid:
            st.markdown('<div class="error-box">❌ PPT模板文件问题</div>', unsafe_allow_html=True)
            st.error(f"模板文件验证失败: {error_msg}")
            st.info("请确保模板文件存在且格式正确")
            return
        
        # 初始化生成器
        try:
            with st.spinner("正在验证API密钥..."):
                generator = StreamlitPPTGenerator(api_key)
        except ValueError as e:
            if "API密钥" in str(e):
                st.error("❌ API密钥验证失败，请检查密钥是否正确")
            else:
                st.error(f"❌ 初始化失败: {str(e)}")
            return
        except Exception as e:
            error_msg = str(e)
            if "authentication" in error_msg.lower() or "unauthorized" in error_msg.lower():
                st.error("❌ API密钥认证失败，请检查密钥是否正确或是否有足够余额")
            elif "network" in error_msg.lower() or "connection" in error_msg.lower():
                st.error("❌ 网络连接异常，请检查网络连接后重试")
            else:
                st.error("❌ 系统初始化异常，请稍后重试")
            st.error(f"详细错误: {error_msg}")
            return
        
        # 加载PPT模板
        with st.spinner("正在加载PPT模板..."):
            if generator.load_ppt_from_path(config.default_ppt_template):
                st.success("✅ PPT模板加载成功！")
                
                # 显示PPT信息
                ppt_info = generator.ppt_structure
                if ppt_info:
                    st.markdown('<div class="info-box">', unsafe_allow_html=True)
                    st.markdown(f"**📊 PPT信息：** 共有 {ppt_info['total_slides']} 张幻灯片")
                    
                    # 显示幻灯片和占位符信息
                    total_placeholders = 0
                    for i, slide in enumerate(ppt_info['slides']):
                        title = slide['title'] if slide['title'] else "（无标题）"
                        placeholders = slide.get('placeholders', {})
                        total_placeholders += len(placeholders)
                        
                        if placeholders:
                            placeholder_list = ', '.join([f"{{{name}}}" for name in placeholders.keys()])
                            st.markdown(f"• 第{slide['slide_index']+1}页: {title} - 占位符: {placeholder_list}")
                        else:
                            st.markdown(f"• 第{slide['slide_index']+1}页: {title} - 无占位符")
                    
                    st.markdown(f"**🎯 总共找到 {total_placeholders} 个占位符**")
                    st.markdown('</div>', unsafe_allow_html=True)
                
                st.markdown("---")
                
                # 文本输入
                st.subheader("✏️ 输入文本内容")
                user_text = st.text_area(
                    "请输入您想要填入PPT的文本内容",
                    height=200,
                    placeholder="请在这里输入您的文本内容...\n\n例如：\n人工智能技术的发展经历了多个重要阶段。从1950年代的符号主义开始，强调逻辑推理和知识表示...",
                    help="保持您的原文不变，AI会智能分析并填入合适的位置"
                )
                
                # 处理按钮
                col1, col2, col3 = st.columns([1, 2, 1])
                with col2:
                    process_button = st.button(
                        "🚀 开始处理",
                        type="primary",
                        use_container_width=True,
                        disabled=not user_text.strip()
                    )
                
                # 处理文本
                if process_button and user_text.strip():
                    # 根据选择的模型显示不同的提示信息
                    current_model_info = config.get_model_info()
                    model_name = current_model_info.get('name', 'AI模型')
                    
                    spinner_text = f"正在使用{model_name}分析文本结构..."
                    if not current_model_info.get('supports_vision', False):
                        spinner_text += "（跳过视觉分析步骤）"
                    
                    try:
                        with st.spinner(spinner_text):
                            assignments = generator.process_text_with_deepseek(user_text)
                    except ValueError as e:
                        if "API密钥" in str(e):
                            st.error("❌ API密钥验证失败，请检查密钥是否正确")
                        else:
                            st.error(f"❌ AI分析失败: {str(e)}")
                        return
                    except Exception as e:
                        error_msg = str(e)
                        if "rate limit" in error_msg.lower():
                            st.error("❌ API请求频率超限，请稍后重试")
                        elif "insufficient" in error_msg.lower() or "quota" in error_msg.lower():
                            st.error("❌ API额度不足，请检查账户余额")
                        else:
                            st.error("❌ AI分析过程出现异常，请稍后重试")
                        st.error(f"详细错误: {error_msg}")
                        return
                    
                    # 显示AI分析结果（调试信息）
                    with st.expander("🔍 查看AI分析结果", expanded=True):
                        st.json(assignments)
                    
                    with st.spinner("正在将文本填入PPT并添加原始文本到备注..."):
                        results = generator.apply_text_assignments(assignments, user_text)
                    
                    # 显示处理结果
                    st.markdown('<div class="success-box">', unsafe_allow_html=True)
                    st.markdown("**✅ 处理完成！**")
                    for result in results:
                        st.markdown(f"• {result}")
                    st.markdown('</div>', unsafe_allow_html=True)
                    
                    # 提供下载
                    st.markdown("---")
                    st.subheader("💾 下载更新后的PPT")
                    
                    try:
                        updated_ppt_bytes = generator.get_ppt_bytes()
                        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                        filename = f"updated_ppt_{timestamp}.pptx"
                        
                        st.download_button(
                            label="📥 下载更新后的PPT",
                            data=updated_ppt_bytes,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            type="primary",
                            use_container_width=True
                        )
                        
                        st.success(f"📁 文件准备完成：{filename}")
                        
                    except Exception as e:
                        st.error(f"生成下载文件时出错: {e}")
            
            else:
                st.error("❌ PPT模板加载失败")
    
    else:
        # 未输入API密钥时的提示
        st.info("👈 请在左侧输入您的OpenRouter API密钥开始使用")
        st.markdown("### 💡 如何获取API密钥")
        st.markdown("""
        1. 访问 [OpenRouter平台](https://openrouter.ai/keys)
        2. 注册或登录账号
        3. 在API密钥管理页面创建新的API密钥
        4. 复制API密钥（格式：sk-xxxxxxxxxxxxx）
        5. 粘贴到左侧输入框中
        """)
        
        # 功能介绍
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### 🎯 核心功能")
            st.markdown("""
            - **预设模板**：使用指定的PPT模板文件
            - **保持原文**：完全保留您的文本内容
            - **智能分析**：AI分析PPT结构和文本逻辑
            - **合理分配**：将文本填入最适合的位置
            - **灵活处理**：更新现有或新增幻灯片
            """)
        
        with col2:
            st.markdown("### 📝 适用场景")
            st.markdown("""
            - **学术报告**：研究内容填入模板
            - **商业计划**：项目信息填入格式
            - **教学课件**：课程内容填入框架
            - **工作汇报**：数据结果填入模板
            """)

if __name__ == "__main__":
    main() 