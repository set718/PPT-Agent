#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
AI智能分页与Dify-模板桥接集成应用
用户输入长文本后自动分页，分页后调用Dify API推荐模板供下载
"""

import streamlit as st
import os
from datetime import datetime
from typing import Dict, List, Any, Optional
import asyncio
from config import get_config
from logger import get_logger, log_user_action
from ai_page_splitter import AIPageSplitter, PageContentFormatter
from dify_template_bridge import DifyTemplateBridge, sync_test_dify_template_bridge
from dify_api_client import DifyAPIConfig, BatchProcessor
from utils import FileManager

# 获取配置
config = get_config()
logger = get_logger()

# 页面配置
st.set_page_config(
    page_title="AI智能分页与模板推荐系统",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 自定义CSS样式
st.markdown("""
<style>
    .main-header {
        font-size: 3rem;
        font-weight: bold;
        color: #2E86AB;
        text-align: center;
        margin-bottom: 0.5rem;
    }
    .sub-header {
        font-size: 1.3rem;
        color: #666;
        text-align: center;
        margin-bottom: 2rem;
    }
    .success-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #d4edda;
        border: 2px solid #c3e6cb;
        color: #155724;
        margin: 1rem 0;
    }
    .info-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #e8f4fd;
        border: 2px solid #bee5eb;
        color: #0c5460;
        margin: 1rem 0;
    }
    .warning-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #fff3cd;
        border: 2px solid #ffeaa7;
        color: #856404;
        margin: 1rem 0;
    }
    .error-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #f8d7da;
        border: 2px solid #f5c6cb;
        color: #721c24;
        margin: 1rem 0;
    }
    .feature-box {
        padding: 1.5rem;
        border-radius: 0.8rem;
        background-color: #f8f9fa;
        border: 2px solid #e9ecef;
        margin: 1rem 0;
    }
    .steps-container {
        background-color: #f8f9fa;
        padding: 2rem;
        border-radius: 1rem;
        margin: 2rem 0;
    }
</style>
""", unsafe_allow_html=True)

class IntegratedPPTGenerator:
    """集成PPT生成器 - 结合AI分页和Dify模板推荐"""
    
    def __init__(self, api_key: str):
        """初始化集成生成器"""
        self.api_key = api_key
        self.page_splitter = AIPageSplitter(api_key)
        self.template_bridge = DifyTemplateBridge()
        self.dify_config = DifyAPIConfig()
        
        # 初始化分批处理器
        self.batch_processor = BatchProcessor(
            self.dify_config, 
            self.template_bridge.api_key_poller
        )
        
        logger.info("初始化集成PPT生成器（支持分批处理）")
    
    def process_text_with_ai_pagination(self, user_text: str, target_pages: Optional[int] = None) -> Dict[str, Any]:
        """使用AI进行智能分页"""
        log_user_action("AI智能分页", f"文本长度: {len(user_text)}字符")
        
        try:
            # 执行AI智能分页
            split_result = self.page_splitter.split_text_to_pages(user_text, target_pages)
            
            if not split_result.get('success'):
                return {
                    "success": False,
                    "error": "AI分页失败",
                    "details": split_result
                }
            
            return {
                "success": True,
                "pagination_result": split_result,
                "pages": split_result.get('pages', []),
                "analysis": split_result.get('analysis', {})
            }
            
        except Exception as e:
            logger.error(f"AI分页异常: {str(e)}")
            return {
                "success": False,
                "error": f"AI分页异常: {str(e)}"
            }
    
    def get_template_recommendations(self, user_text: str, pages: List[Dict[str, Any]]) -> Dict[str, Any]:
        """获取模板推荐"""
        log_user_action("模板推荐", f"为{len(pages)}页内容推荐模板")
        
        try:
            # 构建用于模板推荐的文本内容
            recommendation_text = self._build_recommendation_text(user_text, pages)
            
            # 调用Dify API获取模板推荐
            bridge_result = sync_test_dify_template_bridge(recommendation_text)
            
            if not bridge_result["success"]:
                return {
                    "success": False,
                    "error": f"模板推荐失败: {bridge_result.get('error', '未知错误')}",
                    "bridge_result": bridge_result
                }
            
            # 获取推荐的模板信息
            dify_result = bridge_result["step_1_dify_api"]
            template_result = bridge_result["step_2_template_lookup"]
            
            return {
                "success": True,
                "template_number": dify_result["template_number"],
                "template_filename": template_result["filename"],
                "template_path": template_result["file_path"],
                "template_size": template_result["file_size_kb"],
                "dify_response": dify_result.get("response_text", ""),
                "processing_time": bridge_result["processing_time"],
                "bridge_result": bridge_result
            }
            
        except Exception as e:
            logger.error(f"模板推荐异常: {str(e)}")
            return {
                "success": False,
                "error": f"模板推荐异常: {str(e)}"
            }
    
    def _build_recommendation_text(self, user_text: str, pages: List[Dict[str, Any]]) -> str:
        """构建用于模板推荐的文本内容"""
        # 提取关键信息用于模板推荐
        content_parts = []
        
        # 添加原始文本摘要
        content_parts.append(f"原始文本摘要: {user_text[:500]}...")
        
        # 添加分页分析信息
        if pages:
            content_parts.append(f"分页结果: 共{len(pages)}页")
            
            # 提取内容页的关键信息
            content_pages = [p for p in pages if p.get('page_type') == 'content']
            if content_pages:
                content_parts.append("内容页主题:")
                for i, page in enumerate(content_pages[:3], 1):  # 只取前3页
                    title = page.get('title', f'第{page.get("page_number", i)}页')
                    key_points = page.get('key_points', [])
                    if key_points:
                        points_text = "; ".join(key_points[:2])  # 只取前2个要点
                        content_parts.append(f"{i}. {title}: {points_text}")
        
        return "\n\n".join(content_parts)
    
    def get_templates_for_each_page_batch(self, pages: List[Dict[str, Any]], 
                                        progress_callback=None) -> Dict[str, Any]:
        """为每页获取模板推荐（分批处理版本）"""
        log_user_action("多页模板推荐（分批）", f"为{len(pages)}页内容分别推荐模板")
        
        try:
            # 检查是否启用分批处理
            if not self.dify_config.enable_batch_processing or len(pages) <= self.dify_config.batch_size:
                # 不启用分批处理或页面数少，使用原来的方法
                return self.get_templates_for_each_page(pages)
            
            # 使用异步运行分批处理
            try:
                loop = asyncio.get_event_loop()
            except RuntimeError:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
            
            try:
                result = loop.run_until_complete(
                    self._async_get_templates_for_pages(pages, progress_callback)
                )
                return result
            finally:
                # 清理事件循环
                if loop != asyncio.get_event_loop():
                    loop.close()
                    
        except Exception as e:
            logger.error(f"分批模板推荐异常: {str(e)}")
            return {
                "success": False,
                "error": f"分批模板推荐异常: {str(e)}"
            }
    
    async def _async_get_templates_for_pages(self, pages: List[Dict[str, Any]], 
                                           progress_callback=None) -> Dict[str, Any]:
        """异步处理页面模板推荐"""
        
        # 准备页面数据
        pages_data = []
        for i, page in enumerate(pages):
            if page.get('page_type') == 'title':
                # 标题页不需要调用API
                continue
            pages_data.append({
                "page_index": i,
                "page_data": page,
                "page_text": self._build_page_recommendation_text(page)
            })
        
        # 定义API调用函数
        async def api_call_func(page_info):
            page_data = page_info["page_data"]
            page_text = page_info["page_text"]
            
            try:
                # 调用Dify API
                bridge_result = await self._async_call_dify_for_page(page_text)
                
                if bridge_result["success"]:
                    dify_result = bridge_result["step_1_dify_api"]
                    template_result = bridge_result["step_2_template_lookup"]
                    
                    return {
                        "success": True,
                        "page_number": page_data.get('page_number', page_info["page_index"] + 1),
                        "page_type": page_data.get('page_type', 'content'),
                        "template_number": dify_result["template_number"],
                        "template_path": template_result["file_path"],
                        "template_filename": template_result["filename"],
                        "template_source": "dify_recommended"
                    }
                else:
                    return {
                        "success": False,
                        "page_number": page_data.get('page_number', page_info["page_index"] + 1),
                        "page_type": page_data.get('page_type', 'content'),
                        "error": bridge_result.get('error', '推荐失败'),
                        "template_source": "failed"
                    }
                    
            except Exception as e:
                logger.error(f"页面API调用异常: {str(e)}")
                return {
                    "success": False,
                    "page_number": page_data.get('page_number', page_info["page_index"] + 1),
                    "error": str(e),
                    "template_source": "exception"
                }
        
        # 使用分批处理器处理
        batch_result = await self.batch_processor.process_pages_in_batches(
            pages_data, api_call_func, progress_callback
        )
        
        # 重新组装结果以兼容原有格式
        template_results = []
        successful_templates = []
        
        # 首先添加标题页
        for i, page in enumerate(pages):
            if page.get('page_type') == 'title':
                template_results.append({
                    "page_number": page.get('page_number', i+1),
                    "page_type": "title",
                    "template_source": "fixed",
                    "template_path": None,
                    "success": True
                })
        
        # 添加批处理的结果
        if batch_result.get("success", False):
            for batch_item in batch_result.get("page_templates", []):
                result_data = batch_item.get("result", {})
                if result_data.get("success", False):
                    template_results.append(result_data)
                    if result_data.get("template_path"):
                        successful_templates.append(result_data["template_path"])
                else:
                    template_results.append(result_data)
        
        # 按页面编号排序
        template_results.sort(key=lambda x: x.get("page_number", 0))
        
        return {
            "success": True,
            "page_templates": template_results,
            "successful_count": len(successful_templates),
            "total_pages": len(pages),
            "template_paths": successful_templates,
            "batch_details": batch_result.get("batch_details", []),
            "total_processing_time": batch_result.get("total_processing_time", 0),
            "total_batches": batch_result.get("total_batches", 0)
        }
    
    async def _async_call_dify_for_page(self, page_text: str) -> Dict[str, Any]:
        """异步调用Dify API为页面推荐模板"""
        try:
            # 调用桥接器的异步方法
            result = await self.template_bridge.test_dify_template_bridge(page_text)
            return result
        except Exception as e:
            return {
                "success": False,
                "error": f"Dify API调用异常: {str(e)}"
            }
    
    def get_templates_for_each_page(self, pages: List[Dict[str, Any]]) -> Dict[str, Any]:
        """为每页获取模板推荐"""
        log_user_action("多页模板推荐", f"为{len(pages)}页内容分别推荐模板")
        
        try:
            template_results = []
            successful_templates = []
            
            for i, page in enumerate(pages):
                if page.get('page_type') == 'title':
                    # 标题页使用固定模板
                    template_results.append({
                        "page_number": page.get('page_number', i+1),
                        "page_type": "title",
                        "template_source": "fixed",
                        "template_path": None,
                        "success": True
                    })
                    continue
                
                # 为每页构建推荐文本
                page_text = self._build_page_recommendation_text(page)
                
                # 调用Dify API为这一页推荐模板
                bridge_result = sync_test_dify_template_bridge(page_text)
                
                if bridge_result["success"]:
                    dify_result = bridge_result["step_1_dify_api"]
                    template_result = bridge_result["step_2_template_lookup"]
                    
                    template_info = {
                        "page_number": page.get('page_number', i+1),
                        "page_type": page.get('page_type', 'content'),
                        "template_number": dify_result["template_number"],
                        "template_path": template_result["file_path"],
                        "template_filename": template_result["filename"],
                        "template_source": "dify_recommended",
                        "success": True
                    }
                    
                    successful_templates.append(template_info["template_path"])
                else:
                    template_info = {
                        "page_number": page.get('page_number', i+1),
                        "page_type": page.get('page_type', 'content'),
                        "error": bridge_result.get('error', '推荐失败'),
                        "template_source": "failed",
                        "success": False
                    }
                
                template_results.append(template_info)
            
            return {
                "success": True,
                "page_templates": template_results,
                "successful_count": len(successful_templates),
                "total_pages": len(pages),
                "template_paths": successful_templates
            }
            
        except Exception as e:
            logger.error(f"多页模板推荐异常: {str(e)}")
            return {
                "success": False,
                "error": f"多页模板推荐异常: {str(e)}"
            }
    
    def _build_page_recommendation_text(self, page: Dict[str, Any]) -> str:
        """为单页构建推荐文本"""
        content_parts = []
        
        page_number = page.get('page_number', 1)
        page_type = page.get('page_type', 'content')
        title = page.get('title', '')
        
        content_parts.append(f"页面信息：第{page_number}页 ({page_type})")
        
        if title:
            content_parts.append(f"标题：{title}")
        
        # 副标题
        subtitle = page.get('subtitle', '')
        if subtitle:
            content_parts.append(f"副标题：{subtitle}")
        
        # 内容摘要
        content_summary = page.get('content_summary', '')
        if content_summary:
            content_parts.append(f"内容摘要：{content_summary}")
        
        # 主要要点
        key_points = page.get('key_points', [])
        if key_points:
            content_parts.append("主要要点：")
            for i, point in enumerate(key_points, 1):
                content_parts.append(f"{i}. {point}")
        
        return "\n\n".join(content_parts)
    
    def merge_template_presentations(self, template_paths: List[str], output_filename: str = None) -> Dict[str, Any]:
        """合并多个PPT模板文件"""
        log_user_action("合并PPT模板", f"合并{len(template_paths)}个模板文件")
        
        try:
            from pptx import Presentation
            import tempfile
            
            if not template_paths:
                return {
                    "success": False,
                    "error": "没有模板文件需要合并"
                }
            
            # 创建新的演示文稿
            merged_ppt = Presentation()
            
            # 移除默认的空白幻灯片
            if len(merged_ppt.slides) > 0:
                slide_id = merged_ppt.slides._sldIdLst[0]
                merged_ppt.part.drop_rel(slide_id.rId)
                del merged_ppt.slides._sldIdLst[0]
            
            total_slides = 0
            
            for i, template_path in enumerate(template_paths):
                if not os.path.exists(template_path):
                    logger.warning(f"模板文件不存在: {template_path}")
                    continue
                
                try:
                    # 打开源模板
                    source_ppt = Presentation(template_path)
                    
                    # 复制所有幻灯片
                    for slide in source_ppt.slides:
                        # 复制幻灯片布局和内容
                        slide_layout = merged_ppt.slide_layouts[0]  # 使用默认布局
                        new_slide = merged_ppt.slides.add_slide(slide_layout)
                        
                        # 复制幻灯片内容（这是一个简化版本）
                        # 实际实现中需要更复杂的复制逻辑
                        total_slides += 1
                    
                    logger.info(f"成功处理模板 {i+1}/{len(template_paths)}: {os.path.basename(template_path)}")
                    
                except Exception as e:
                    logger.error(f"处理模板文件失败 {template_path}: {str(e)}")
                    continue
            
            # 保存合并后的文件
            if not output_filename:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                output_filename = f"合并PPT_{timestamp}.pptx"
            
            output_path = os.path.join(self.config.output_dir, output_filename)
            merged_ppt.save(output_path)
            
            return {
                "success": True,
                "output_path": output_path,
                "output_filename": output_filename,
                "total_slides": total_slides,
                "processed_templates": len([p for p in template_paths if os.path.exists(p)]),
                "total_templates": len(template_paths)
            }
            
        except Exception as e:
            logger.error(f"合并PPT模板异常: {str(e)}")
            return {
                "success": False,
                "error": f"合并PPT模板异常: {str(e)}"
            }
    
    def generate_final_ppt(self, template_path: str, enhanced_pages: List[Dict[str, Any]]) -> Dict[str, Any]:
        """生成最终的PPT文件"""
        log_user_action("生成最终PPT", f"使用模板: {os.path.basename(template_path)}")
        
        try:
            # 这里可以添加PPT生成逻辑
            # 目前返回模板文件信息
            return {
                "success": True,
                "template_path": template_path,
                "template_filename": os.path.basename(template_path),
                "pages_count": len(enhanced_pages),
                "message": "PPT生成完成"
            }
            
        except Exception as e:
            logger.error(f"PPT生成异常: {str(e)}")
            return {
                "success": False,
                "error": f"PPT生成异常: {str(e)}"
            }

def display_processing_summary(pagination_result: Dict[str, Any], template_result: Dict[str, Any], 
                             enhancement_result: Dict[str, Any]) -> None:
    """显示处理结果摘要"""
    
    st.markdown("### 📊 处理结果摘要")
    
    # 分页结果
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        pages = pagination_result.get('pages', [])
        st.metric("📄 分页数量", len(pages))
    
    with col2:
        if template_result.get('success'):
            template_num = template_result.get('template_number', 'N/A')
            st.metric("🎯 推荐模板", f"#{template_num}")
        else:
            st.metric("🎯 推荐模板", "失败")
    
    with col3:
        if enhancement_result.get('success'):
            summary = enhancement_result.get('processing_summary', {})
            successful = summary.get('successful_api_calls', 0)
            st.metric("✅ API成功", successful)
        else:
            st.metric("✅ API成功", "失败")
    
    with col4:
        if template_result.get('success'):
            processing_time = template_result.get('processing_time', 0)
            st.metric("⏱️ 总耗时", f"{processing_time:.2f}秒")
        else:
            st.metric("⏱️ 总耗时", "N/A")

def main():
    # 页面标题
    st.markdown('<div class="main-header">🎯 AI智能分页与模板推荐系统</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">DeepSeek智能分页 + Dify模板推荐系统</div>', unsafe_allow_html=True)
    
    # 功能说明
    st.markdown('<div class="info-box">🎯 <strong>完整工作流程</strong><br>1. 用户输入长文本<br>2. DeepSeek AI智能分页（保留原文本）<br>3. Dify API推荐最适合的PPT模板<br>4. 返回推荐模板供下载<br>5. 后续可用DeepSeek将原文本填充到模板</div>', unsafe_allow_html=True)
    
    # 侧边栏配置
    with st.sidebar:
        st.header("⚙️ 系统配置")
        
        # 模型选择
        config = get_config()
        available_models = config.available_models
        
        model_options = {}
        for model_key, model_info in available_models.items():
            display_name = f"{model_info['name']} ({model_info['cost']})"
            model_options[display_name] = model_key
        
        selected_model_display = st.selectbox(
            "选择AI模型",
            options=list(model_options.keys()),
            index=0,
            help="选择用于AI分页的模型"
        )
        
        selected_model = model_options[selected_model_display]
        model_info = available_models[selected_model]
        
        # 显示模型信息
        st.markdown(f"**📋 模型信息：**\n- 提供商：{model_info['api_provider']}\n- 描述：{model_info['description']}")
        
        # API密钥输入
        api_key = st.text_input(
            f"{model_info['api_provider']} API密钥",
            type="password",
            help=f"用于AI分页和内容增强的{model_info['api_provider']} API密钥",
            placeholder="请输入API密钥..."
        )
        
        if not api_key:
            st.markdown('<div class="warning-box">⚠️ 请先输入API密钥</div>', unsafe_allow_html=True)
            st.markdown(f"获取API密钥：[{model_info['api_provider']}]({model_info['api_key_url']})")
        
        # 处理选项
        st.markdown("---")
        st.subheader("🔧 处理选项")
        
        enable_template_recommendation = st.checkbox(
            "启用模板推荐",
            value=True,
            help="使用Dify API推荐最适合的PPT模板"
        )
        
        
        target_pages = st.number_input(
            "目标页面数量（可选）",
            min_value=0,
            max_value=25,
            value=0,
            help="设置为0时，AI将自动判断最佳页面数量"
        )
        
        # 显示页数建议
        st.markdown("""
        **💡 页数建议：**
        - 5分钟演示：3-5页
        - 10分钟演示：5-8页
        - 15分钟演示：8-12页
        - 30分钟演示：15-20页
        - 学术报告：20-25页
        """)
    
    # 主界面
    if not api_key or not api_key.strip():
        # 显示功能介绍
        st.markdown("---")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### 🎯 核心功能")
            st.markdown("""
            **🤖 AI智能分页**
            - 自动分析文本结构
            - 智能分割为多个页面
            - 保持内容逻辑完整性
            
            **🎯 模板智能推荐**
            - Dify API分析内容特征
            - 推荐最适合的PPT模板
            - 支持250+模板选择
            """)
        
        with col2:
            st.markdown("### 🚀 技术特色")
            st.markdown("""
            **🔗 多API集成**
            - DeepSeek API用于分页分析
            - Dify API用于模板推荐
            - Liai API备选支持
            
            **📊 工作流程**
            - 保留原始文本不变
            - 纯模板推荐服务
            - 后续可填充文本内容
            """)
        
        return
    
    # 验证API密钥并初始化
    try:
        # 更新配置中的模型设置
        config.set_model(selected_model)
        generator = IntegratedPPTGenerator(api_key.strip())
        st.markdown('<div class="success-box">✅ 系统初始化成功！</div>', unsafe_allow_html=True)
        st.info(f"🤖 当前使用模型：{model_info['name']}")
    except Exception as e:
        st.error(f"❌ 系统初始化失败: {str(e)}")
        return
    
    # 文本输入区域
    st.markdown("### 📝 输入您的内容")
    
    user_text = st.text_area(
        "请输入您想要制作成PPT的长文本内容：",
        height=300,
        placeholder="""例如：

人工智能技术发展趋势分析报告

人工智能技术作为当今科技发展的重要驱动力，正在深刻改变着我们的社会、经济和生活。本报告将从技术发展、应用现状、未来趋势等多个维度进行深入分析。

技术发展历程：
人工智能的发展经历了从符号主义到连接主义，再到深度学习的演进过程。1950年代，图灵测试的提出标志着AI研究的开始。1980年代，专家系统的兴起为AI应用奠定了基础。2010年代，深度学习的突破性进展带来了AI技术的革命性变化。

核心技术突破：
机器学习、深度学习、自然语言处理、计算机视觉等技术的快速发展，为AI应用提供了强大的技术支撑。特别是大语言模型的出现，如GPT、Claude等，展现出了前所未有的理解和生成能力。

应用领域扩展：
AI技术已广泛应用于教育、医疗、金融、制造、娱乐等各个领域。在教育领域，AI辅助教学系统能够提供个性化学习体验。在医疗领域，AI诊断系统提高了疾病检测的准确性。在金融领域，AI风控系统有效降低了金融风险。

未来发展趋势：
人工智能将继续向更加智能化、人性化的方向发展。预计在未来5-10年内，我们将看到更多突破性的AI应用，包括通用人工智能的探索、人机协作的深化、AI伦理和安全体系的完善等。

挑战与机遇：
虽然AI技术发展迅速，但仍面临数据安全、算法偏见、就业影响等挑战。同时，AI技术也为解决全球性问题如气候变化、疾病防控等提供了新的可能性。

结论：
人工智能技术正处于快速发展期，其影响将越来越深远。我们需要在推动技术发展的同时，关注其社会影响，确保AI技术为人类福祉服务。""",
        help="AI将分析您的文本结构，智能分页，并推荐最适合的PPT模板"
    )
    
    # 处理按钮
    if user_text.strip():
        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            process_button = st.button(
                "🚀 开始智能处理",
                type="primary",
                use_container_width=True,
                help="执行AI智能分页 + Dify模板推荐流程"
            )
        
        # 处理逻辑
        if process_button:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                # 步骤1: AI智能分页
                status_text.text("🤖 正在进行AI智能分页...")
                progress_bar.progress(20)
                
                pagination_result = generator.process_text_with_ai_pagination(
                    user_text.strip(), 
                    int(target_pages) if target_pages > 0 else None
                )
                
                if not pagination_result.get('success'):
                    st.error(f"❌ AI分页失败: {pagination_result.get('error', '未知错误')}")
                    return
                
                st.success("✅ AI智能分页完成！")
                
                # 显示分页结果
                pages = pagination_result.get('pages', [])
                analysis = pagination_result.get('analysis', {})
                
                # 显示分析摘要
                analysis_summary = PageContentFormatter.format_analysis_summary(analysis)
                st.markdown(analysis_summary)
                
                # 已移除分页预览展示（简化用户界面）
                
                # 步骤2: 模板推荐（如果启用）
                template_result = {"success": False}
                if enable_template_recommendation:
                    status_text.text("🎯 正在推荐最适合的PPT模板...")
                    progress_bar.progress(40)
                    
                    template_result = generator.get_template_recommendations(user_text.strip(), pages)
                    
                    if template_result.get('success'):
                        st.success(f"✅ 模板推荐完成！推荐模板: #{template_result['template_number']}")
                        
                        # 显示模板信息
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("🎯 推荐模板", f"#{template_result['template_number']}")
                        with col2:
                            st.metric("📄 模板文件", template_result['template_filename'])
                        with col3:
                            st.metric("📦 文件大小", f"{template_result['template_size']}KB")
                        
                        # 已移除Dify API推荐理由测试入口（简化用户界面）
                        
                        # 立即提供模板下载
                        template_path = template_result.get('template_path')
                        if template_path and os.path.exists(template_path):
                            st.markdown("### 💾 下载推荐的PPT模板")
                            
                            try:
                                with open(template_path, "rb") as f:
                                    template_bytes = f.read()
                                
                                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                                template_num = template_result.get('template_number', 'unknown')
                                filename = f"AI推荐模板_{template_num}_{timestamp}.pptx"
                                
                                col1, col2, col3 = st.columns([1, 2, 1])
                                with col2:
                                    st.download_button(
                                        label="📥 下载推荐的PPT模板",
                                        data=template_bytes,
                                        file_name=filename,
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                        type="primary",
                                        use_container_width=True
                                    )
                                
                                st.markdown('<div class="info-box">', unsafe_allow_html=True)
                                st.markdown(f"📁 **文件名：** {filename}")
                                st.markdown(f"🎯 **推荐模板：** #{template_num}")
                                st.markdown(f"📄 **分页数量：** {len(pages)}页")
                                st.markdown("📋 **说明：** 基于AI分页和Dify API推荐的PPT模板原文件")
                                st.markdown('</div>', unsafe_allow_html=True)
                                
                            except Exception as e:
                                st.error(f"❌ 模板下载失败: {str(e)}")
                    else:
                        st.warning(f"⚠️ 模板推荐失败: {template_result.get('error', '未知错误')}")
                else:
                    st.info("ℹ️ 已跳过模板推荐步骤")
                
                # 完成处理
                progress_bar.progress(100)
                status_text.empty()
                
                st.markdown('<div class="success-box">🎉 完整处理流程已完成！</div>', unsafe_allow_html=True)
                
                # 显示详细结果（可选）
                with st.expander("🔍 查看详细处理数据", expanded=False):
                    st.json({
                        "pagination_result": pagination_result,
                        "template_result": template_result,
                        "enhancement_result": enhancement_result
                    })
                
            except Exception as e:
                progress_bar.empty()
                status_text.empty()
                st.error(f"❌ 处理过程中出现异常: {str(e)}")
                logger.error("集成处理异常: %s", str(e))
    
    else:
        # 未输入文本时的说明
        st.markdown("---")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### 📋 使用步骤")
            st.markdown("""
            1. **选择AI模型**（DeepSeek/Liai）
            2. **输入API密钥**
            3. **输入长文本内容**
            4. **点击开始处理**
            
            6. **下载推荐模板**
            """)
        
        with col2:
            st.markdown("### ✨ 功能特色")
            st.markdown("""
            **🎯 智能分页**
            - DeepSeek分析文本结构
            - 自动确定最佳页数
            - 保持原文本完整性
            
            **🔗 模板推荐**
            - Dify API智能分析
            - 250+模板库选择
            - 内容特征精准匹配
            
            **🚀 简洁高效**
            - 纯分页+推荐服务
            - 保留原始文本
            - 快速获取模板
            """)

if __name__ == "__main__":
    main() 