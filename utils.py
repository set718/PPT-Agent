#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
工具函数模块
包含项目中的共用工具函数
"""

import os
import re
import json
import time
from datetime import datetime
from typing import Dict, List, Any, Optional, Tuple
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from config import get_config
from ppt_beautifier import PPTBeautifier
from ppt_advanced_analyzer import PPTStructureAnalyzer, PositionExtractor, SmartLayoutAdjuster, create_advanced_ppt_analysis
from ppt_visual_analyzer import PPTVisualAnalyzer, VisualLayoutOptimizer

class PPTAnalyzer:
    """PPT分析器"""
    
    @staticmethod
    def analyze_ppt_structure(presentation: Presentation) -> Dict[str, Any]:
        """
        分析PPT结构，提取占位符和文本信息
        
        Args:
            presentation: PPT演示文稿对象
            
        Returns:
            Dict: PPT结构信息
        """
        slides_info = []
        
        for i, slide in enumerate(presentation.slides):
            slide_info = {
                "slide_index": i,
                "title": "",
                "placeholders": {},
                "text_shapes": [],
                "has_content": False
            }
            
            # 分析幻灯片中的文本框、表格和占位符
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    current_text = shape.text.strip()
                    if current_text:
                        # 检查是否包含占位符 - 识别所有{}格式的占位符
                        placeholder_pattern = r'\{([^}]+)\}'
                        placeholders = re.findall(placeholder_pattern, current_text)
                        
                        if placeholders:
                            # 这个文本框包含占位符
                            # 为了避免多个占位符在同一个文本框中的冲突，
                            # 我们记录这个文本框的所有占位符
                            for placeholder in placeholders:
                                slide_info["placeholders"][placeholder] = {
                                    "shape": shape,
                                    "original_text": current_text,
                                    "placeholder": placeholder,
                                    "all_placeholders": placeholders,  # 记录同一文本框中的所有占位符
                                    "type": "text_shape",  # 标识为文本框
                                    "text": current_text  # 添加文本内容用于调试
                                }
                        
                        # 如果是简短文本且没有占位符，可能是标题
                        if not placeholders and len(current_text) < 100:
                            if slide_info["title"] == "":
                                slide_info["title"] = current_text
                        
                        slide_info["has_content"] = True
                    
                    # 记录所有可编辑的文本形状
                    if hasattr(shape, "text_frame"):
                        slide_info["text_shapes"].append({
                            "shape_id": shape.shape_id if hasattr(shape, "shape_id") else len(slide_info["text_shapes"]),
                            "current_text": shape.text,
                            "shape": shape,
                            "has_placeholder": bool(re.findall(r'\{([^}]+)\}', shape.text)) if shape.text else False
                        })
                
                # 处理表格中的占位符
                elif shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
                    # 这是一个表格
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if cell_text:
                                # 检查表格单元格中的占位符
                                placeholder_pattern = r'\{([^}]+)\}'
                                placeholders = re.findall(placeholder_pattern, cell_text)
                                
                                if placeholders:
                                    # 表格单元格包含占位符
                                    for placeholder in placeholders:
                                        slide_info["placeholders"][placeholder] = {
                                            "shape": shape,
                                            "table": table,
                                            "cell": cell,
                                            "row_idx": row_idx,
                                            "col_idx": col_idx,
                                            "original_text": cell_text,
                                            "placeholder": placeholder,
                                            "all_placeholders": placeholders,
                                            "type": "table_cell"  # 标识为表格单元格
                                        }
                                        
                                # 记录表格单元格为文本形状（用于调试）
                                slide_info["text_shapes"].append({
                                    "shape_id": f"table_{row_idx}_{col_idx}",
                                    "current_text": cell_text,
                                    "shape": cell,  # 单元格对象
                                    "table_info": {
                                        "table": table,
                                        "row_idx": row_idx,
                                        "col_idx": col_idx
                                    },
                                    "has_placeholder": bool(placeholders),
                                    "type": "table_cell"
                                })
                                
                                slide_info["has_content"] = True
            
            slides_info.append(slide_info)
        
        return {
            "total_slides": len(presentation.slides),
            "slides": slides_info
        }

class AIProcessor:
    """AI处理器"""
    
    def __init__(self, api_key: str = None):
        """初始化AI处理器"""
        config = get_config()
        model_info = config.get_model_info()
        
        # 处理API密钥设置
        if api_key:
            self.api_key = api_key
        else:
            # 检查模型配置中的环境变量设置
            api_key_env = model_info.get('api_key_env')
            if api_key_env:
                import os
                # 如果是火山引擎且支持多密钥，优先使用第一个密钥
                if model_info.get('api_provider') == 'Volces' and model_info.get('use_multiple_keys'):
                    # 尝试获取多个密钥，使用第一个可用的
                    for i in range(1, 6):
                        key = os.getenv(f'{api_key_env}_{i}')
                        if key:
                            self.api_key = key
                            break
                    else:
                        # 如果没找到编号密钥，尝试单个密钥
                        self.api_key = os.getenv(api_key_env) or config.openai_api_key or ""
                else:
                    self.api_key = os.getenv(api_key_env) or config.openai_api_key or ""
            else:
                self.api_key = config.openai_api_key
        
        if not self.api_key:
            raise ValueError("请设置API密钥")
        
        # 根据当前选择的模型获取对应的base_url
        self.base_url = model_info.get('base_url', config.openai_base_url)
        
        # 延迟初始化client，避免在创建时就验证API密钥
        self.client = None
        self.config = config
    
    def _ensure_client(self):
        """确保client已初始化"""
        if self.client is None:
            try:
                self.client = OpenAI(
                    api_key=self.api_key,
                    base_url=self.base_url
                )
            except Exception as e:
                raise ValueError(f"API密钥验证失败: {str(e)}")
    
    def analyze_text_for_ppt(self, user_text: str, ppt_structure: Dict[str, Any], enhanced_info: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        使用AI分析文本并生成PPT填充方案
        
        Args:
            user_text: 用户输入的文本
            ppt_structure: PPT结构信息
            enhanced_info: 增强的结构信息（可选）
            
        Returns:
            Dict: 文本分配方案
        """
        # 确保client已初始化
        self._ensure_client()
        
        # 创建PPT结构描述
        if enhanced_info:
            ppt_description = self._create_enhanced_ppt_description(enhanced_info)
        else:
            ppt_description = self._create_ppt_description(ppt_structure)
        
        # 构建系统提示
        system_prompt = self._build_system_prompt(ppt_description)
        
        # 检查是否为Liai API
        model_info = self.config.get_model_info()
        if model_info.get('request_format') == 'dify_compatible':
            # 使用Liai API格式，单页的所有占位符用一次API调用处理
            content = self._call_liai_api(system_prompt, user_text)
        else:
            # 使用OpenAI格式
            try:
                # 使用actual_model而不是ai_model配置名
                actual_model = model_info.get('actual_model', self.config.ai_model)
                response = self.client.chat.completions.create(
                    model=actual_model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_text}
                    ],
                    temperature=self.config.ai_temperature,
                    max_tokens=self.config.ai_max_tokens,
                    stream=True
                )
                
                # 收集流式响应内容
                content = ""
                for chunk in response:
                    if chunk.choices and chunk.choices[0].delta and chunk.choices[0].delta.content:
                        content += chunk.choices[0].delta.content
                
                content = content.strip() if content else ""
                
            except Exception as e:
                raise e
        
        try:
            # 提取JSON内容
            return self._extract_json_from_response(content, user_text)
            
        except Exception as e:
            print("调用AI API时出错: %s", str(e))
            error_msg = str(e)
            
            # 检查是否是OpenAI API的特定错误
            if hasattr(e, 'status_code'):
                status_code = e.status_code
                if status_code == 401:
                    return self._create_fallback_assignment(user_text, f"❌ GPT API认证失败 (401): API密钥无效，请检查密钥是否正确")
                elif status_code == 402:
                    return self._create_fallback_assignment(user_text, f"❌ GPT API付费限制 (402): 账户余额不足，请充值后重试")
                elif status_code == 403:
                    return self._create_fallback_assignment(user_text, f"❌ GPT API权限拒绝 (403): 当前API密钥没有访问权限")
                elif status_code == 404:
                    return self._create_fallback_assignment(user_text, f"❌ GPT API模型不存在 (404): 请检查模型名称是否正确")
                elif status_code == 429:
                    return self._create_fallback_assignment(user_text, f"❌ GPT API请求频率限制 (429): 请求过于频繁，请稍后重试")
                elif status_code == 500:
                    return self._create_fallback_assignment(user_text, f"⚠️ GPT API服务器错误 (500): OpenAI服务器内部错误，请稍后重试")
                elif status_code == 502:
                    return self._create_fallback_assignment(user_text, f"⚠️ GPT API网关错误 (502): 服务暂时不可用，请稍后重试")
                elif status_code == 503:
                    return self._create_fallback_assignment(user_text, f"⚠️ GPT API服务不可用 (503): 服务暂时维护中，请稍后重试")
                else:
                    return self._create_fallback_assignment(user_text, f"❌ GPT API错误 ({status_code}): {error_msg}，这不是文本填充功能的问题")
            
            # 检查其他常见错误类型
            elif "connection" in error_msg.lower() or "network" in error_msg.lower():
                return self._create_fallback_assignment(user_text, f"⚠️ GPT API网络连接失败: 请检查网络连接或稍后重试")
            elif "timeout" in error_msg.lower():
                return self._create_fallback_assignment(user_text, f"⚠️ GPT API请求超时: 请稍后重试")
            elif "authentication" in error_msg.lower() or "unauthorized" in error_msg.lower():
                return self._create_fallback_assignment(user_text, f"❌ GPT API密钥认证失败: 请检查API密钥是否正确")
            elif "quota" in error_msg.lower() or "limit" in error_msg.lower():
                return self._create_fallback_assignment(user_text, f"❌ GPT API配额不足: 请检查账户余额或使用限制")
            else:
                return self._create_fallback_assignment(user_text, f"❌ GPT API调用失败: {error_msg}，这不是文本填充功能的问题")
    
    def _call_liai_api(self, system_prompt: str, user_text: str) -> str:
        """调用Liai API"""
        import requests
        import json
        import os
        import random
        
        # 从环境变量获取API keys用于负载均衡
        liai_api_keys = []
        for i in range(1, 6):
            key = os.getenv(f"LIAI_API_KEY_{i}")
            if key:
                liai_api_keys.append(key)
        
        # 如果有多个API key，随机选择一个
        if liai_api_keys:
            selected_key = random.choice(liai_api_keys)
        else:
            selected_key = self.api_key
        
        model_info = self.config.get_model_info()
        base_url = model_info.get('base_url', '')
        endpoint = model_info.get('chat_endpoint', '/chat-messages')
        
        url = base_url + endpoint
        
        # 构建Liai API请求格式
        combined_query = f"{system_prompt}\n\n用户输入：{user_text}"
        
        payload = {
            "inputs": {},
            "query": combined_query,
            "response_mode": "streaming",
            "conversation_id": "",
            "user": "ai-ppt-user",
            "files": []
        }
        
        headers = {
            'Authorization': f'Bearer {selected_key}',
            'Content-Type': 'application/json',
            'Connection': 'keep-alive'
        }
        
        try:
            response = requests.post(url, headers=headers, json=payload, timeout=120, stream=True)
            response.raise_for_status()
            
            # 处理streaming响应
            content = ""
            for line in response.iter_lines():
                if line:
                    try:
                        line_text = line.decode('utf-8').strip()
                        # 忽略keep-alive注释
                        if line_text == ': keep-alive' or line_text == '':
                            continue
                        if line_text.startswith('data: '):
                            json_str = line_text[6:]  # 去掉'data: '前缀
                            if json_str.strip() == '[DONE]':
                                break
                            data = json.loads(json_str)
                            if 'answer' in data:
                                content += data['answer']
                            elif 'data' in data and 'answer' in data['data']:
                                content += data['data']['answer']
                    except (json.JSONDecodeError, UnicodeDecodeError):
                        continue
            
            return content.strip()
            
        except requests.exceptions.RequestException as e:
            print(f"Liai API调用失败: {str(e)}")
            raise e
        except Exception as e:
            print(f"Liai API处理失败: {str(e)}")
            raise e
    
    def batch_process_liai_requests(self, requests_data: List[Dict], batch_size: int = 5) -> List[Dict]:
        """
        批处理 Liai API 请求，每批5个请求
        
        Args:
            requests_data: 请求数据列表，每个元素包含 system_prompt 和 user_text
            batch_size: 每批处理的请求数量，默认5个
            
        Returns:
            List[Dict]: 处理结果列表
        """
        import time
        
        results = []
        total_requests = len(requests_data)
        
        # 分批处理
        for i in range(0, total_requests, batch_size):
            batch_end = min(i + batch_size, total_requests)
            batch_requests = requests_data[i:batch_end]
            
            print(f"Processing batch {i//batch_size + 1}/{(total_requests + batch_size - 1)//batch_size}: {len(batch_requests)} requests")
            
            batch_results = []
            for j, request_data in enumerate(batch_requests):
                try:
                    print(f"  Processing request {j+1}/{len(batch_requests)}")
                    
                    system_prompt = request_data.get('system_prompt', '')
                    user_text = request_data.get('user_text', '')
                    
                    # 调用 Liai API
                    response = self._call_liai_api(system_prompt, user_text)
                    
                    batch_results.append({
                        'success': True,
                        'response': response,
                        'request_index': i + j,
                        'request_data': request_data
                    })
                    
                    # 请求间延迟，避免频率过高
                    if j < len(batch_requests) - 1:
                        time.sleep(0.5)
                        
                except Exception as e:
                    print(f"  Request {j+1} failed: {str(e)}")
                    batch_results.append({
                        'success': False,
                        'error': str(e),
                        'request_index': i + j,
                        'request_data': request_data
                    })
            
            results.extend(batch_results)
            
            # 批次间延迟，避免API频率限制
            if batch_end < total_requests:
                print(f"  Batch completed. Waiting 2 seconds before next batch...")
                time.sleep(2.0)
        
        print(f"All batches completed. Processed {total_requests} requests.")
        return results
    
    def batch_analyze_pages_for_liai(self, pages_data: List[Dict], batch_size: int = 5) -> List[Dict]:
        """
        为Liai API批处理多页内容分析，每批5页
        
        Args:
            pages_data: 页面数据列表，每个元素包含页面内容和结构信息
            batch_size: 每批处理的页面数量，默认5页
            
        Returns:
            List[Dict]: 每页的分析结果
        """
        import time
        
        results = []
        total_pages = len(pages_data)
        
        print(f"开始Liai批处理分析 {total_pages} 页内容，每批 {batch_size} 页")
        
        # 分批处理
        for batch_idx in range(0, total_pages, batch_size):
            batch_end = min(batch_idx + batch_size, total_pages)
            batch_pages = pages_data[batch_idx:batch_end]
            batch_num = (batch_idx // batch_size) + 1
            total_batches = (total_pages + batch_size - 1) // batch_size
            
            print(f"处理第 {batch_num}/{total_batches} 批，包含 {len(batch_pages)} 页")
            
            batch_results = []
            
            # 处理当前批次的每一页
            for page_idx, page_data in enumerate(batch_pages):
                try:
                    page_number = page_data.get('page_number', batch_idx + page_idx + 1)
                    user_text = page_data.get('content', '')
                    ppt_structure = page_data.get('ppt_structure', {})
                    
                    print(f"  分析第{page_number}页...")
                    
                    # 调用Liai API进行分析
                    analysis_result = self.analyze_text_for_ppt(user_text, ppt_structure)
                    
                    batch_results.append({
                        'page_number': page_number,
                        'content': user_text,
                        'analysis_result': analysis_result,
                        'success': True,
                        'processing_time': 0,  # 实际会在分析中计算
                        'batch_index': batch_num
                    })
                    
                    print(f"  第{page_number}页分析完成")
                    
                    # 页面间延迟
                    if page_idx < len(batch_pages) - 1:
                        time.sleep(0.5)
                        
                except Exception as e:
                    print(f"  第{page_number}页分析失败: {str(e)}")
                    batch_results.append({
                        'page_number': page_number,
                        'content': user_text,
                        'analysis_result': None,
                        'success': False,
                        'error': str(e),
                        'processing_time': 0,
                        'batch_index': batch_num
                    })
            
            results.extend(batch_results)
            
            # 批次间延迟
            if batch_end < total_pages:
                print(f"  第{batch_num}批完成，等待2秒后处理下一批...")
                time.sleep(2.0)
        
        print(f"Liai批处理完成，共处理 {total_pages} 页")
        return results
    
    def _create_ppt_description(self, ppt_structure: Dict[str, Any]) -> str:
        """创建PPT结构描述"""
        description = f"现有PPT共有{ppt_structure['total_slides']}张幻灯片，模板设计意图分析:\n"
        
        # 分析整体结构
        total_placeholders = sum(len(slide.get('placeholders', {})) for slide in ppt_structure['slides'])
        description += f"总占位符数量: {total_placeholders}个，需要智能分配用户文本\n"
        
        # 分析各类占位符分布 - 智能识别所有类型
        placeholder_types = {}
        for slide in ppt_structure['slides']:
            for placeholder_name in slide.get('placeholders', {}).keys():
                placeholder_type = self._analyze_placeholder_type(placeholder_name)
                placeholder_key = placeholder_type.split('-')[0]  # 提取类型部分，如"标题类"
                if placeholder_key in placeholder_types:
                    placeholder_types[placeholder_key] += 1
                else:
                    placeholder_types[placeholder_key] = 1
        
        description += f"占位符类型分布: {dict(placeholder_types)}\n"
        
        # 详细描述每张幻灯片
        for slide in ppt_structure['slides']:
            description += f"\n第{slide['slide_index']+1}页:"
            
            # 幻灯片标题分析
            if slide['title']:
                description += f" 现有标题「{slide['title']}」"
            else:
                description += f" (无现有标题)"
            
            # 占位符详细分析
            if slide['placeholders']:
                description += f"\n  占位符详情:"
                
                # 按重要性排序显示占位符
                sorted_placeholders = sorted(
                    slide['placeholders'].items(),
                    key=lambda x: self._get_placeholder_priority(x[0])
                )
                
                for placeholder_name, placeholder_info in sorted_placeholders:
                    placeholder_type = self._analyze_placeholder_type(placeholder_name)
                    description += f"\n    - {{{placeholder_name}}} [{placeholder_type}]"
                    
                description += f"\n  设计意图: {self._analyze_slide_design_intent(slide)}\n"
            else:
                description += f" (无占位符)\n"
        
        return description
    
    def _get_placeholder_priority(self, placeholder_name: str) -> int:
        """智能获取占位符优先级（数字越小优先级越高），支持复合占位符"""
        name_lower = placeholder_name.lower()
        
        # 分析复合占位符的所有组件
        import re
        components = re.split(r'[_\-\s]+', name_lower)
        all_components = [name_lower] + components
        
        # 标题类：最高优先级
        title_keywords = ['title', 'heading', '主题', 'topic', '标题', 'header']
        if any(keyword in comp for comp in all_components for keyword in title_keywords if keyword in comp):
            return 1
            
        # 副标题类：高优先级
        subtitle_keywords = ['subtitle', 'sub', '副标题', 'secondary']
        if any(keyword in comp for comp in all_components for keyword in subtitle_keywords if keyword in comp):
            return 2
            
        # 人物和时间类：重要信息，优先填充
        person_keywords = ['author', 'name', 'speaker', '演讲者', '作者', '姓名', 'presenter', 'who']
        time_keywords = ['time', 'date', '年份', '日期', '时间', 'year', 'month', 'day', 'schedule', 'when']
        if any(keyword in comp for comp in all_components for keyword in person_keywords + time_keywords if keyword in comp):
            return 2
            
        # 主题/话题类：重要内容标识
        topic_keywords = ['topic', 'subject', '主题', '话题', '议题']
        if any(keyword in comp for comp in all_components for keyword in topic_keywords if keyword in comp):
            return 3
            
        # 要点类：中等优先级（bullet类型通常很重要）
        bullet_keywords = ['bullet', 'point', 'list', 'item', '要点', '列表', '项目']
        if any(keyword in comp for comp in all_components for keyword in bullet_keywords if keyword in comp):
            return 3
            
        # 内容类：中高优先级
        content_keywords = ['content', 'text', 'description', '介绍', '内容', '描述', 'detail', 'info']
        if any(keyword in comp for comp in all_components for keyword in content_keywords if keyword in comp):
            return 4
            
        # 数据类：中等优先级
        data_keywords = ['number', 'data', 'percentage', 'statistic', '统计', '数字', '数据', '百分比', 'count']
        if any(keyword in comp for comp in all_components for keyword in data_keywords if keyword in comp):
            return 4
            
        # 结论类：较高优先级
        conclusion_keywords = ['conclusion', 'summary', '结论', '总结', 'result', '结果']
        if any(keyword in comp for comp in all_components for keyword in conclusion_keywords if keyword in comp):
            return 3
            
        # 未知类型：较低优先级，但仍会处理
        return 5
    
    def _analyze_placeholder_type(self, placeholder_name: str) -> str:
        """智能分析占位符类型，根据名称语义自动判断，支持复合命名格式"""
        name_lower = placeholder_name.lower()
        
        # 分析复合占位符的所有组件
        # 使用下划线、连字符等分隔符分割占位符名称
        import re
        components = re.split(r'[_\-\s]+', name_lower)
        all_components = [name_lower] + components  # 包含完整名称和所有组件
        
        # 计算各类型的匹配权重
        type_scores = {}
        
        # 标题类占位符检测
        title_keywords = ['title', 'heading', '主题', 'topic', '标题', 'header']
        title_score = sum(1 for comp in all_components for keyword in title_keywords if keyword in comp)
        if title_score > 0:
            type_scores['标题类-高视觉权重'] = title_score
        
        # 副标题类占位符检测  
        subtitle_keywords = ['subtitle', 'sub', '副标题', 'secondary']
        subtitle_score = sum(1 for comp in all_components for keyword in subtitle_keywords if keyword in comp)
        if subtitle_score > 0:
            type_scores['副标题类-中高视觉权重'] = subtitle_score
        
        # 要点类占位符检测（优先于内容类检测）
        bullet_keywords = ['bullet', 'point', 'list', 'item', '要点', '列表', '项目']
        bullet_score = sum(1 for comp in all_components for keyword in bullet_keywords if keyword in comp)
        if bullet_score > 0:
            type_scores['要点类-核心信息'] = bullet_score + 1  # 给要点类额外权重
        
        # 时间类占位符检测（检测包含time等的复合词）
        time_keywords = ['time', 'date', '年份', '日期', '时间', 'year', 'month', 'day', 'schedule', 'when']
        time_score = sum(1 for comp in all_components for keyword in time_keywords if keyword in comp)
        if time_score > 0:
            type_scores['时间类-日期信息'] = time_score
            
        # 主题/话题类占位符检测（适合topic等）
        topic_keywords = ['topic', 'subject', '主题', '话题', '议题']  
        topic_score = sum(1 for comp in all_components for keyword in topic_keywords if keyword in comp)
        if topic_score > 0:
            type_scores['主题类-内容标识'] = topic_score
        
        # 人物类占位符检测
        person_keywords = ['author', 'name', 'speaker', '演讲者', '作者', '姓名', 'presenter', 'who']
        person_score = sum(1 for comp in all_components for keyword in person_keywords if keyword in comp)
        if person_score > 0:
            type_scores['人物类-身份信息'] = person_score
        
        # 内容类占位符检测
        content_keywords = ['content', 'text', 'description', '介绍', '内容', '描述', 'detail', 'info']
        content_score = sum(1 for comp in all_components for keyword in content_keywords if keyword in comp)
        if content_score > 0:
            type_scores['内容类-框架构建'] = content_score
            
        # 数据类占位符检测
        data_keywords = ['number', 'data', 'percentage', 'statistic', '统计', '数字', '数据', '百分比', 'count']
        data_score = sum(1 for comp in all_components for keyword in data_keywords if keyword in comp)
        if data_score > 0:
            type_scores['数据类-数值信息'] = data_score
        
        # 结论类占位符检测
        conclusion_keywords = ['conclusion', 'summary', '结论', '总结', 'result', '结果']
        conclusion_score = sum(1 for comp in all_components for keyword in conclusion_keywords if keyword in comp)
        if conclusion_score > 0:
            type_scores['结论类-总结升华'] = conclusion_score
        
        # 返回得分最高的类型
        if type_scores:
            best_type = max(type_scores.items(), key=lambda x: x[1])
            return f"{best_type[0]}(复合:{'+'.join(components)})"
        
        # 如果都不匹配，返回通用类型，但提供组件分析
        return f"通用类-复合占位符({'+'.join(components)})"
    
    def _analyze_slide_design_intent(self, slide: Dict[str, Any]) -> str:
        """智能分析幻灯片设计意图，根据占位符类型自动判断页面用途"""
        placeholders = slide.get('placeholders', {})
        if not placeholders:
            return "纯展示页面，无需填充"
        
        placeholder_names = [name.lower() for name in placeholders.keys()]
        
        # 智能检测各类占位符
        title_keywords = ['title', 'heading', '主题', 'topic', '标题', 'header']
        content_keywords = ['content', 'text', 'description', '介绍', '内容', '描述', 'detail']
        bullet_keywords = ['bullet', 'point', 'list', 'item', '要点', '列表', '项目']
        person_keywords = ['author', 'name', 'speaker', '演讲者', '作者', '姓名', 'presenter']
        time_keywords = ['date', 'time', '年份', '日期', '时间', 'year', 'month', 'day']
        data_keywords = ['number', 'data', 'percentage', 'statistic', '统计', '数字', '数据', '百分比']
        
        has_title = any(any(keyword in name for keyword in title_keywords) for name in placeholder_names)
        has_content = any(any(keyword in name for keyword in content_keywords) for name in placeholder_names)
        has_bullets = any(any(keyword in name for keyword in bullet_keywords) for name in placeholder_names)
        has_person = any(any(keyword in name for keyword in person_keywords) for name in placeholder_names)
        has_time = any(any(keyword in name for keyword in time_keywords) for name in placeholder_names)
        has_data = any(any(keyword in name for keyword in data_keywords) for name in placeholder_names)
        
        # 根据占位符组合判断页面类型
        if has_person and has_time:
            return "封面型页面，适合标题展示和基本信息"
        elif has_title and has_bullets:
            return "标题要点型页面，适合概要展示和要点列举"
        elif has_content and has_bullets:
            return "内容详解型页面，适合分点阐述和详细说明"
        elif has_title and has_content:
            return "标题内容型页面，适合主题阐述和内容展开"
        elif has_data:
            return "数据展示型页面，适合统计信息和数字展示"
        elif len(placeholders) > 3:
            return "复合型页面，包含多种信息类型，需要平衡布局"
        else:
            return f"灵活型页面，根据实际占位符({list(placeholders.keys())})智能安排内容"
    
    def _create_enhanced_ppt_description(self, enhanced_info: Dict[str, Any]) -> str:
        """创建增强的PPT结构描述"""
        basic_structure = enhanced_info.get('basic_structure', {})
        advanced_analysis = enhanced_info.get('advanced_analysis', {})
        position_analysis = enhanced_info.get('position_analysis', {})
        layout_suggestions = enhanced_info.get('layout_suggestions', [])
        
        # 基础信息
        total_slides = basic_structure.get('total_slides', 0)
        description = f"现有PPT共有{total_slides}张幻灯片，高级结构分析如下:\n"
        
        # 添加整体结构分析
        if advanced_analysis:
            overall_structure = advanced_analysis.get('overall_structure', {})
            if overall_structure:
                description += f"\n【整体设计分析】\n"
                description += f"• 整体风格：{overall_structure.get('overall_style', '未知')}\n"
                description += f"• 设计一致性：{overall_structure.get('design_consistency', 0):.2f}/1.0\n"
                
                avg_metrics = overall_structure.get('average_metrics', {})
                if avg_metrics:
                    description += f"• 平均内容密度：{avg_metrics.get('content_density', 0):.2f}/1.0\n"
                    description += f"• 平均视觉平衡度：{avg_metrics.get('visual_balance', 0):.2f}/1.0\n"
                    description += f"• 平均层次清晰度：{avg_metrics.get('hierarchy_clarity', 0):.2f}/1.0\n"
                
                layout_dist = overall_structure.get('layout_distribution', {})
                if layout_dist:
                    description += f"• 布局类型分布：{layout_dist}\n"
        
        # 添加详细的幻灯片分析
        slide_layouts = advanced_analysis.get('slide_layouts', [])
        for i, slide_layout in enumerate(slide_layouts):
            description += f"\n第{i+1}页详细分析：\n"
            description += f"• 布局类型：{slide_layout.layout_type}\n"
            description += f"• 设计意图：{slide_layout.design_intent}\n"
            description += f"• 内容密度：{slide_layout.content_density:.2f}/1.0\n"
            description += f"• 视觉平衡度：{slide_layout.visual_balance:.2f}/1.0\n"
            description += f"• 层次清晰度：{slide_layout.hierarchy_clarity:.2f}/1.0\n"
            
            # 添加元素信息
            elements = slide_layout.elements
            if elements:
                description += f"• 包含{len(elements)}个元素：\n"
                for element in elements:
                    if element.placeholder_name:
                        description += f"  - {{{element.placeholder_name}}} [{element.element_type}] 视觉权重:{element.visual_weight}/5\n"
                        description += f"    位置:(x:{element.position.left:.0f}, y:{element.position.top:.0f}, w:{element.position.width:.0f}, h:{element.position.height:.0f})\n"
            
            # 添加视觉区域分析
            visual_regions = slide_layout.visual_regions
            if visual_regions:
                description += f"• 视觉区域分布：\n"
                for region_name, region_elements in visual_regions.items():
                    if region_elements:
                        description += f"  - {region_name}区域：{len(region_elements)}个元素\n"
        
        # 添加布局建议
        if layout_suggestions:
            description += f"\n【布局优化建议】\n"
            for suggestion in layout_suggestions:
                slide_idx = suggestion.get('slide_index', 0)
                suggestions = suggestion.get('suggestions', {})
                
                layout_sugg = suggestions.get('layout_suggestions', [])
                if layout_sugg:
                    description += f"第{slide_idx+1}页建议：\n"
                    for sugg in layout_sugg:
                        description += f"• {sugg.get('description', '')}\n"
        
        # 添加位置分析摘要
        if position_analysis:
            description += f"\n【空间布局分析】\n"
            spatial_relationships = position_analysis.get('spatial_relationships', {})
            if spatial_relationships:
                description += f"• 幻灯片间布局一致性分析已完成\n"
                # 可以添加更多空间关系的描述
        
        return description
    
    def _build_system_prompt(self, ppt_description: str) -> str:
        """构建系统提示"""
        return """你是一个专业的PPT内容分析专家，具备强大的PPT文件识别能力。你的任务是将用户提供的文本内容智能分配到PPT模板的合适占位符中。

**你的PPT识别能力包括：**
- **文本内容**：议程页 (Agenda with icons)、欢迎页 (Welcome 10:00 am)、团队介绍 (The team 11:00 am)、服务介绍 (Our services 12:00 pm)、愿景展示 (Vision 1:00 pm) 等
- **结构信息**：每一页的标题、正文、占位符（所有{}格式，AI需根据占位符名称理解其含义）  
- **布局元素**：能读取每张幻灯片的布局类型 (标题+正文、两栏布局、带图标的议程、带图文的组合页等)
- **样式信息**：字体名称、字号、是否加粗/斜体、颜色等
- **对象元素**：图标、图片、形状、表格等 (能知道它们存在、类型、位置和大小参数)

**重要原则：**
1. 充分利用你的PPT识别能力，深度理解模板结构和设计意图
2. 只使用用户提供的文本内容，不生成新内容
3. 可以对文本进行适当的优化、精简或重新组织
4. 根据占位符的语义含义和布局位置选择最合适的内容片段
5. 不是所有占位符都必须填充，只填充有合适内容的占位符

现有PPT深度结构分析：
%s""" % ppt_description + """

**核心任务（基于深度识别）：**
**占位符处理规则：**
- 识别并处理所有{}格式占位符，支持文本框和表格
- **占位符命名语义规则**：找到语义词汇确定内容类型，忽略技术参数词汇
  * `{bullet_2}` - 重点是bullet（要点内容），数字2表示第2个要点
  * `{bullet_2_time_1}` - 重点是time（时间段），表示第2个要点下的第1个时间段
  * `{title_max_token_50}` - 重点是title（标题），max_token_50是技术参数，忽略
  * `{section_3_title}` - 重点是title（标题），表示第3个部分的标题
- **内容类型匹配**：根据语义词汇确定填充内容类型，忽略数字和技术参数
  * **time**: 填入时间段或时间点（如"9:45-10:00"、"上午9点"）
  * **title**: 填入标题性内容
  * **bullet/content**: 填入具体要点或详细内容
  * **author**: 填入人名或机构名
  * **date**: 填入日期信息

**约束条件：**
- 只使用用户提供的信息，不生成新内容
- 保持原文核心含义不变
- 优先填充重要占位符

**处理原则：**
- 根据占位符名称的语义含义智能匹配内容类型
- 控制内容长度，标题简洁、要点明确、描述适中
- 保持语言风格一致，避免重复和冗长

**输出格式：**
只返回JSON，包含assignments数组：
- slide_index: 幻灯片索引（从0开始）
- action: "replace_placeholder"
- placeholder: 占位符名称
- content: 优化后的填充内容
- reason: 选择理由

只返回JSON，不要其他文字。"""
    
    def _extract_json_from_response(self, content: str, user_text: str) -> Dict[str, Any]:
        """从AI响应中提取JSON"""
        # 提取JSON内容（如果有代码块包围）
        json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', content, re.DOTALL)
        if json_match:
            content = json_match.group(1)
        
        try:
            return json.loads(content)
        except json.JSONDecodeError as e:
            print(f"AI返回的JSON格式有误，错误：{str(e)}")
            return self._create_fallback_assignment(user_text, f"JSON解析失败: {str(e)}")
    
    def _create_fallback_assignment(self, user_text: str, error_msg: str) -> Dict[str, Any]:
        """创建备用分配方案"""
        return {
            "assignments": [
                {
                    "slide_index": 0,
                    "action": "replace_placeholder",
                    "placeholder": "content",
                    "content": user_text,
                    "reason": "API调用失败或解析错误，默认填入content占位符。支持所有{}格式占位符。错误: " + str(error_msg)
                }
            ]
        }

class PPTProcessor:
    """PPT处理器"""
    
    def __init__(self, presentation: Presentation):
        """初始化PPT处理器"""
        self.presentation = presentation
        self.ppt_structure = PPTAnalyzer.analyze_ppt_structure(presentation)
        self.beautifier = PPTBeautifier(presentation)
        self.filled_placeholders = {}  # 记录已填充的占位符
        
        # 初始化高级分析器
        self.advanced_analysis = create_advanced_ppt_analysis(presentation)
        self.structure_analyzer = self.advanced_analysis['analyzers']['structure_analyzer'] if 'analyzers' in self.advanced_analysis else None
        self.position_extractor = self.advanced_analysis['analyzers']['position_extractor'] if 'analyzers' in self.advanced_analysis else None
        self.layout_adjuster = self.advanced_analysis['analyzers']['layout_adjuster'] if 'analyzers' in self.advanced_analysis else None
        
        # 视觉分析器（需要API密钥时才初始化）
        self.visual_analyzer = None
        self.visual_optimizer = None
    
    def get_enhanced_structure_info(self) -> Dict[str, Any]:
        """获取增强的PPT结构信息"""
        if not self.structure_analyzer:
            return self.ppt_structure
        
        # 合并基础分析和高级分析结果
        enhanced_info = {
            'basic_structure': self.ppt_structure,
            'advanced_analysis': self.advanced_analysis.get('structure_analysis', {}),
            'position_analysis': self.advanced_analysis.get('position_analysis', {}),
            'layout_suggestions': []
        }
        
        # 为每张幻灯片生成布局建议
        if self.layout_adjuster and 'structure_analysis' in self.advanced_analysis:
            slide_layouts = self.advanced_analysis['structure_analysis'].get('slide_layouts', [])
            for i, layout in enumerate(slide_layouts):
                # 模拟一些内容来生成建议
                mock_content = {}
                if i < len(self.ppt_structure['slides']):
                    slide_info = self.ppt_structure['slides'][i]
                    for placeholder in slide_info.get('placeholders', {}).keys():
                        mock_content[placeholder] = f"示例内容_{placeholder}"
                
                if mock_content:
                    suggestions = self.layout_adjuster.suggest_optimal_layout(i, mock_content)
                    enhanced_info['layout_suggestions'].append({
                        'slide_index': i,
                        'suggestions': suggestions
                    })
        
        return enhanced_info
    
    def initialize_visual_analyzer(self, api_key: str) -> bool:
        """
        初始化视觉分析器（仅在启用视觉分析时）
        
        Args:
            api_key: OpenAI API密钥
            
        Returns:
            bool: 初始化是否成功
        """
        # 检查配置是否启用视觉分析
        config = get_config()
        if not config.enable_visual_analysis:
            print(f"[INFO] 当前模型 {config.ai_model} 不支持视觉分析，跳过视觉分析器初始化")
            self.visual_analyzer = None
            self.visual_optimizer = None
            return True  # 返回True表示按配置正确初始化
        
        try:
            self.visual_analyzer = PPTVisualAnalyzer(api_key)
            self.visual_optimizer = VisualLayoutOptimizer(self.visual_analyzer)
            print(f"[INFO] 视觉分析器初始化成功，使用模型: {config.ai_model}")
            return True
        except Exception as e:
            print("视觉分析器初始化失败: %s", str(e))
            return False
    
    def analyze_visual_quality(self, ppt_path: str) -> Dict[str, Any]:
        """
        分析PPT视觉质量（如果启用了视觉分析功能）
        
        Args:
            ppt_path: PPT文件路径
            
        Returns:
            Dict: 视觉分析结果
        """
        config = get_config()
        
        if not config.enable_visual_analysis:
            # 视觉分析被禁用，返回简单的默认分析结果
            return {
                "analysis_skipped": True,
                "reason": f"当前使用的模型 {config.ai_model} 不支持视觉分析功能",
                "slides_analysis": [],
                "overall_quality": {
                    "visual_appeal": 0.5,
                    "content_balance": 0.5,
                    "consistency": 0.5
                }
            }
        
        if not self.visual_analyzer:
            return {"error": "视觉分析器未初始化，请先提供API密钥"}
        
        try:
            return self.visual_analyzer.analyze_presentation_visual_quality(ppt_path)
        except Exception as e:
            return {"error": f"视觉分析失败: {e}"}
    
    def apply_visual_optimizations(self, visual_analysis: Dict[str, Any]) -> Dict[str, Any]:
        """
        应用视觉优化建议
        
        Args:
            visual_analysis: 视觉分析结果
            
        Returns:
            Dict: 优化结果
        """
        if not self.visual_optimizer:
            return {"error": "视觉优化器未初始化"}
        
        try:
            slide_analyses = visual_analysis.get("slide_analyses", [])
            optimization_results = []
            
            for slide_analysis in slide_analyses:
                slide_index = slide_analysis.get("slide_index", 0)
                result = self.visual_optimizer.optimize_slide_layout(
                    self.presentation, slide_index, slide_analysis
                )
                optimization_results.append(result)
            
            return {
                "success": True,
                "optimization_results": optimization_results,
                "total_optimizations": sum(
                    len(r.get("optimizations_applied", [])) 
                    for r in optimization_results 
                    if r.get("success")
                )
            }
            
        except Exception as e:
            return {"error": f"视觉优化失败: {e}"}
    
    def apply_assignments(self, assignments: Dict[str, Any], user_text: str = "") -> List[str]:
        """
        应用文本分配方案
        
        Args:
            assignments: 分配方案
            user_text: 用户原始文本（可选，用于添加到幻灯片备注）
            
        Returns:
            List[str]: 处理结果列表
        """
        assignments_list = assignments.get('assignments', [])
        results = []
        
        # 清理旧缓存并提取最新的格式信息
        print("清理旧格式缓存...")
        self._clear_format_cache()
        print("预先提取所有占位符格式信息...")
        self._cache_all_placeholder_formats(assignments_list)
        
        # 如果提供了用户原始文本，则为幻灯片添加备注
        if user_text.strip():
            notes_results = self._add_notes_to_slides(assignments_list, user_text)
            results.extend(notes_results)
        
        # 按文本框分组处理分配，避免多次刷新同一文本框
        shape_assignments = {}  # {shape_id: [assignments]}
        other_assignments = []  # 非占位符替换的其他操作
        
        # 先将分配按文本框分组
        for assignment in assignments_list:
            action = assignment.get('action')
            content = assignment.get('content', '')
            slide_index = assignment.get('slide_index', 0)
            
            if action == 'replace_placeholder':
                placeholder = assignment.get('placeholder', '')
                if 0 <= slide_index < len(self.presentation.slides):
                    slide_info = self.ppt_structure['slides'][slide_index]
                    if placeholder in slide_info['placeholders']:
                        placeholder_info = slide_info['placeholders'][placeholder]
                        shape = placeholder_info['shape']
                        shape_type = placeholder_info.get('type', 'text_box')
                        
                        # 创建唯一的shape标识
                        if shape_type == 'table_cell':
                            shape_id = f"slide_{slide_index}_table_{id(shape)}_cell_{placeholder_info['row_idx']}_{placeholder_info['col_idx']}"
                        else:
                            shape_id = f"slide_{slide_index}_shape_{id(shape)}"
                        
                        if shape_id not in shape_assignments:
                            shape_assignments[shape_id] = []
                        shape_assignments[shape_id].append(assignment)
                    else:
                        results.append(f"ERROR: 第{slide_index+1}页不存在 {{{placeholder}}} 占位符")
                else:
                    results.append(f"ERROR: 幻灯片索引 {slide_index+1} 超出范围")
            else:
                other_assignments.append(assignment)
        
        # 批量处理每个文本框的所有占位符
        for shape_id, assignments in shape_assignments.items():
            batch_success = self._replace_placeholders_in_shape_batch(assignments)
            
            # 记录结果
            for assignment in assignments:
                slide_index = assignment.get('slide_index', 0)
                placeholder = assignment.get('placeholder', '')
                
                # 记录为已处理
                if slide_index not in self.filled_placeholders:
                    self.filled_placeholders[slide_index] = set()
                self.filled_placeholders[slide_index].add(placeholder)
                
                if batch_success:
                    results.append(f"SUCCESS: 已替换第{slide_index+1}页的 {{{placeholder}}} 占位符: {assignment.get('reason', '')}")
                else:
                    results.append(f"WARNING: 第{slide_index+1}页的 {{{placeholder}}} 占位符替换失败，但已标记为已处理")
        
        # 处理其他类型的操作
        for assignment in other_assignments:
            action = assignment.get('action')
            content = assignment.get('content', '')
            slide_index = assignment.get('slide_index', 0)
            
            if action == 'update':
                if 0 <= slide_index < len(self.presentation.slides):
                    slide = self.presentation.slides[slide_index]
                    self._update_slide_content(slide, content)
                    results.append(f"SUCCESS: 已更新第{slide_index+1}页: {assignment.get('reason', '')}")
                
            elif action == 'add_new':
                title = assignment.get('title', '新增内容')
                self._add_new_slide(title, content)
                results.append(f"SUCCESS: 已新增幻灯片「{title}」: {assignment.get('reason', '')}")
        
        return results
    
    def _replace_placeholders_in_shape_batch(self, assignments: List[Dict]) -> bool:
        """逐个处理占位符，每个占位符单独替换并应用格式"""
        if not assignments:
            return False
            
        try:
            # 逐个处理每个占位符
            for assignment in assignments:
                slide_index = assignment.get('slide_index', 0)
                placeholder = assignment.get('placeholder', '')
                content = assignment.get('content', '')
                
                # 获取占位符信息
                slide_info = self.ppt_structure['slides'][slide_index]
                if placeholder not in slide_info['placeholders']:
                    continue
                    
                placeholder_info = slide_info['placeholders'][placeholder]
                
                # 使用统一的占位符替换方法（已经支持表格和文本框）
                success = self._replace_placeholder_in_slide_with_cached_format(placeholder_info, content)
                if not success:
                    print(f"占位符{placeholder}替换失败")
                    return False
                    
            return True
                
        except Exception as e:
            print(f"逐个替换占位符时出错: {e}")
            return False
    
    def _clear_format_cache(self):
        """清理所有占位符的格式缓存，确保使用最新格式"""
        cleared_count = 0
        for slide_info in self.ppt_structure['slides']:
            for placeholder_name, placeholder_info in slide_info.get('placeholders', {}).items():
                if 'cached_format' in placeholder_info:
                    del placeholder_info['cached_format']
                    cleared_count += 1
        
        if cleared_count > 0:
            print(f"   已清理{cleared_count}个占位符的旧格式缓存")
        else:
            print("   无需清理，首次使用")
    
    def _cache_all_placeholder_formats(self, assignments_list: List[Dict]):
        """预先提取所有占位符的格式信息，避免替换过程中格式丢失"""
        cached_count = 0
        for assignment in assignments_list:
            if assignment.get('action') == 'replace_placeholder':
                slide_index = assignment.get('slide_index', 0)
                placeholder = assignment.get('placeholder', '')
                
                if 0 <= slide_index < len(self.presentation.slides):
                    slide_info = self.ppt_structure['slides'][slide_index]
                    
                    if placeholder in slide_info['placeholders']:
                        placeholder_info = slide_info['placeholders'][placeholder]
                        
                        # 只有在还没有缓存格式时才提取
                        if 'cached_format' not in placeholder_info:
                            # 根据占位符类型选择正确的容器
                            container = placeholder_info.get('cell') if placeholder_info.get('type') == 'table_cell' else placeholder_info.get('shape')
                            format_info = self._extract_placeholder_format(container, placeholder)
                            placeholder_info['cached_format'] = format_info
                            cached_count += 1
                            font_size = format_info.get('font_size')
                            if font_size is not None:
                                font_size = float(font_size.pt) if hasattr(font_size, 'pt') else font_size
                            font_color = format_info.get('font_color', 'None')
                            print(f"   缓存格式: 第{slide_index+1}页 {{{placeholder}}} - 字体:{format_info.get('font_name', 'None')}, 大小:{font_size}, 颜色:{font_color}")
        
        print(f"格式缓存完成，共缓存{cached_count}个占位符的格式信息")
    
    def _add_notes_to_slides(self, assignments_list: List[Dict], user_text: str) -> List[str]:
        """
        为幻灯片添加用户原始文本备注
        
        Args:
            assignments_list: 分配方案列表
            user_text: 用户原始文本
            
        Returns:
            List[str]: 备注添加结果
        """
        results = []
        
        # 获取涉及的幻灯片索引
        involved_slides = set()
        for assignment in assignments_list:
            slide_index = assignment.get('slide_index', 0)
            if 0 <= slide_index < len(self.presentation.slides):
                involved_slides.add(slide_index)
        
        # 如果只有一张幻灯片被涉及，将完整的用户文本添加到该幻灯片
        if len(involved_slides) == 1:
            slide_index = list(involved_slides)[0]
            success = self._add_note_to_slide(slide_index, user_text)
            if success:
                results.append(f"NOTES: 已将原始文本添加到第{slide_index+1}页备注")
            else:
                results.append(f"NOTES ERROR: 添加备注到第{slide_index+1}页失败")
        
        # 如果涉及多张幻灯片，智能分割用户文本
        elif len(involved_slides) > 1:
            text_segments = self._split_text_for_slides(user_text, involved_slides, assignments_list)
            for slide_index, text_segment in text_segments.items():
                if text_segment.strip():
                    success = self._add_note_to_slide(slide_index, text_segment)
                    if success:
                        results.append(f"NOTES: 已将相关文本添加到第{slide_index+1}页备注")
                    else:
                        results.append(f"NOTES ERROR: 添加备注到第{slide_index+1}页失败")
        
        return results
    
    def _add_note_to_slide(self, slide_index: int, note_text: str) -> bool:
        """
        为指定幻灯片添加备注
        
        Args:
            slide_index: 幻灯片索引
            note_text: 备注文本
            
        Returns:
            bool: 是否成功添加备注
        """
        try:
            slide = self.presentation.slides[slide_index]
            
            # 获取或创建备注页
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
            else:
                notes_slide = slide.notes_slide  # 这会自动创建notes_slide
            
            # 获取备注文本框
            notes_text_frame = notes_slide.notes_text_frame
            
            # 设置备注内容
            if notes_text_frame.text.strip():
                # 如果已有备注，添加分隔符和新内容
                notes_text_frame.text += f"\n\n【原始文本】\n{note_text}"
            else:
                # 如果没有备注，直接添加
                notes_text_frame.text = f"【原始文本】\n{note_text}"
            
            return True
            
        except Exception as e:
            print(f"添加备注失败: {e}")
            return False
    
    def _split_text_for_slides(self, user_text: str, involved_slides: set, assignments_list: List[Dict]) -> Dict[int, str]:
        """
        智能分割用户文本，为不同幻灯片分配相关的文本段落
        
        Args:
            user_text: 用户原始文本
            involved_slides: 涉及的幻灯片索引集合
            assignments_list: 分配方案列表
            
        Returns:
            Dict[int, str]: 每张幻灯片对应的文本段落
        """
        # 按段落分割用户文本
        paragraphs = [p.strip() for p in user_text.split('\n\n') if p.strip()]
        if not paragraphs:
            paragraphs = [user_text]
        
        # 为每张幻灯片分配文本段落
        slide_texts = {}
        sorted_slides = sorted(involved_slides)
        
        # 如果段落数量 >= 幻灯片数量，平均分配
        if len(paragraphs) >= len(sorted_slides):
            paragraphs_per_slide = len(paragraphs) // len(sorted_slides)
            remainder = len(paragraphs) % len(sorted_slides)
            
            start_idx = 0
            for i, slide_index in enumerate(sorted_slides):
                end_idx = start_idx + paragraphs_per_slide
                if i < remainder:
                    end_idx += 1
                
                slide_paragraphs = paragraphs[start_idx:end_idx]
                slide_texts[slide_index] = '\n\n'.join(slide_paragraphs)
                start_idx = end_idx
        else:
            # 如果段落少于幻灯片，优先为前几张幻灯片分配
            for i, slide_index in enumerate(sorted_slides):
                if i < len(paragraphs):
                    slide_texts[slide_index] = paragraphs[i]
                else:
                    # 剩余幻灯片分享最后一个段落或完整文本
                    slide_texts[slide_index] = user_text if len(paragraphs) == 1 else paragraphs[-1]
        
        return slide_texts
    
    def beautify_presentation(self, enable_visual_optimization: bool = False, ppt_path: str = None) -> Dict[str, Any]:
        """
        美化演示文稿，清理未填充的占位符并重新排版
        
        Args:
            enable_visual_optimization: 是否启用视觉优化
            ppt_path: PPT文件路径（视觉分析需要）
            
        Returns:
            Dict: 美化结果
        """
        beautify_results = self.beautifier.cleanup_and_beautify(self.filled_placeholders)
        optimization_results = self.beautifier.optimize_slide_sequence()
        
        # 基础美化结果
        result = {
            'beautify_results': beautify_results,
            'optimization_results': optimization_results,
            'summary': {
                'removed_placeholders_count': sum(
                    item['removed_count'] for item in beautify_results['removed_placeholders']
                ),
                'reorganized_slides_count': len(beautify_results['reorganized_slides']),
                'removed_empty_slides_count': len(optimization_results['removed_empty_slides']),
                'final_slide_count': optimization_results['final_slide_count']
            }
        }
        
        # 如果启用视觉优化且视觉分析器可用
        if enable_visual_optimization and self.visual_analyzer and ppt_path:
            try:
                print("🎨 执行视觉质量分析...")
                visual_analysis = self.analyze_visual_quality(ppt_path)
                
                if "error" not in visual_analysis:
                    print("🔧 应用视觉优化建议...")
                    visual_optimization = self.apply_visual_optimizations(visual_analysis)
                    
                    result['visual_analysis'] = visual_analysis
                    result['visual_optimization'] = visual_optimization
                    result['summary']['visual_optimizations_applied'] = visual_optimization.get('total_optimizations', 0)
                    
                    overall_score = visual_analysis.get('overall_analysis', {}).get('weighted_score', 0)
                    result['summary']['visual_quality_score'] = overall_score
                else:
                    result['visual_analysis'] = {"error": visual_analysis.get("error")}
                    
            except Exception as e:
                result['visual_analysis'] = {"error": f"视觉分析过程中出错: {e}"}
        
        return result
    
    def _replace_placeholder_in_slide(self, placeholder_info: Dict[str, Any], new_content: str) -> bool:
        """在文本框或表格单元格中替换占位符，保持原有格式"""
        try:
            placeholder_name = placeholder_info['placeholder']
            placeholder_pattern = f"{{{placeholder_name}}}"
            
            # 判断是表格单元格还是普通文本框
            if placeholder_info.get('type') == 'table_cell':
                # 处理表格单元格中的占位符
                return self._replace_placeholder_in_table_cell(placeholder_info, new_content)
            
            # 处理普通文本框中的占位符
            shape = placeholder_info['shape']
            
            # 检查当前文本框的实际内容
            current_text = shape.text if hasattr(shape, 'text') else ""
            
            if placeholder_pattern not in current_text:
                print(f"占位符 {placeholder_pattern} 在文本 '{current_text}' 中未找到")
                return False
            
            # 执行文本替换
            updated_text = current_text.replace(placeholder_pattern, new_content, 1)
            
            print(f"替换占位符: {placeholder_pattern} -> '{new_content}'")
            
            # 保持格式的文本替换
            if hasattr(shape, "text_frame") and shape.text_frame:
                # 提取原始格式信息
                original_format = self._extract_text_format(shape)
                
                # 应用新文本并保持格式
                return self._apply_text_with_format(shape, updated_text, original_format)
            else:
                # 直接设置text属性（备用方案）
                shape.text = updated_text
                return True
                
        except Exception as e:
            print("替换占位符时出错: %s", str(e))
            return False
    
    def _replace_placeholder_in_table_cell(self, placeholder_info: Dict[str, Any], new_content: str) -> bool:
        """在表格单元格中替换占位符"""
        try:
            cell = placeholder_info['cell']
            placeholder_name = placeholder_info['placeholder']
            row_idx = placeholder_info['row_idx']
            col_idx = placeholder_info['col_idx']
            
            # 获取单元格当前文本
            current_text = cell.text
            placeholder_pattern = f"{{{placeholder_name}}}"
            
            if placeholder_pattern not in current_text:
                print(f"表格占位符 {placeholder_pattern} 在单元格[{row_idx},{col_idx}]文本 '{current_text}' 中未找到")
                return False
            
            # 执行文本替换
            updated_text = current_text.replace(placeholder_pattern, new_content, 1)
            
            print(f"替换占位符: {placeholder_pattern} -> '{new_content}'")
            
            # 直接替换单元格文本
            cell.text = updated_text
            
            return True
                
        except Exception as e:
            print(f"替换表格占位符时出错: {str(e)}")
            return False
    
    def _replace_placeholder_in_slide_with_cached_format(self, placeholder_info: Dict[str, Any], new_content: str) -> bool:
        """使用预先缓存的格式信息替换占位符 - 统一使用shape级别处理"""
        try:
            placeholder_name = placeholder_info['placeholder']
            
            # 统一使用shape级别处理所有占位符类型
            # 根据用户反馈：正确填充的级别都是shape，不正确的都是run
            if placeholder_info.get('type') == 'table_cell':
                # 表格单元格也使用shape级别的处理逻辑
                return self._replace_single_placeholder_in_table_cell(placeholder_info, new_content)
            else:
                # 文本框占位符统一使用shape级别处理
                return self._replace_single_placeholder_in_shape(placeholder_info, new_content)
                
        except Exception as e:
            print("替换占位符时出错: %s", str(e))
            return False
    
    def _replace_single_placeholder_in_table_cell(self, placeholder_info: Dict[str, Any], new_content: str) -> bool:
        """在表格单元格中替换单个占位符，应用缓存的格式"""
        try:
            cell = placeholder_info['cell']
            placeholder_pattern = f"{{{placeholder_info['placeholder']}}}"
            cached_format = placeholder_info.get('cached_format', {})
            
            # 直接在单元格文本中替换
            current_text = cell.text
            if placeholder_pattern not in current_text:
                return False
            
            updated_text = current_text.replace(placeholder_pattern, new_content, 1)
            cell.text = updated_text
            
            # 应用格式到单元格的文本框
            if cached_format and hasattr(cell, 'text_frame') and cell.text_frame:
                self._apply_format_to_cell(cell, cached_format)
                # 输出替换和格式应用信息
                font_size = cached_format.get('font_size', 'None')
                font_color = cached_format.get('font_color', 'None')
                print(f"替换占位符: {placeholder_pattern} -> '{new_content}' - 字体:{cached_format.get('font_name', 'None')}, 大小:{font_size}, 颜色:{font_color}")
            else:
                print(f"替换占位符: {placeholder_pattern} -> '{new_content}'")
            
            return True
            
        except Exception as e:
            print(f"表格占位符替换失败: {e}")
            return False
    
    def _replace_single_placeholder_in_run(self, placeholder_info: Dict[str, Any], new_content: str) -> bool:
        """在run中替换单个占位符，保持run的格式"""
        try:
            run = placeholder_info['run']
            placeholder_pattern = f"{{{placeholder_info['placeholder']}}}"
            cached_format = placeholder_info.get('cached_format', {})
            
            # 在run文本中替换
            if placeholder_pattern not in run.text:
                return False
            
            run.text = run.text.replace(placeholder_pattern, new_content, 1)
            
            # 应用格式到run
            if cached_format:
                self._apply_format_to_run(run, cached_format)
            
            return True
            
        except Exception as e:
            print(f"Run级占位符替换失败: {e}")
            return False
    
    def _replace_single_placeholder_in_shape(self, placeholder_info: Dict[str, Any], new_content: str) -> bool:
        """在文本框形状中替换单个占位符，保持其他占位符不变"""
        print(f"DEBUG: 进入_replace_single_placeholder_in_shape方法")
        try:
            shape = placeholder_info['shape']
            placeholder_pattern = f"{{{placeholder_info['placeholder']}}}"
            cached_format = placeholder_info.get('cached_format', {})
            
            # 在runs级别进行替换，保持格式
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder_pattern in run.text:
                            # 保存原始格式，然后替换文本
                            original_font = run.font
                            original_name = original_font.name
                            original_size = original_font.size
                            original_bold = original_font.bold
                            original_italic = original_font.italic
                            original_color = original_font.color
                            
                            # 在包含占位符的run中替换文本
                            run.text = run.text.replace(placeholder_pattern, new_content, 1)
                            
                            # 立即重新应用该占位符的缓存格式
                            if cached_format:
                                self._apply_format_to_run(run, cached_format)
                                # 验证格式是否真的被应用  
                                print(f"    格式已应用到run")
                                
                            # 输出替换和格式应用信息
                            font_size = cached_format.get('font_size', 'None') if cached_format else 'None'
                            font_color = cached_format.get('font_color', 'None') if cached_format else 'None'
                            font_name = cached_format.get('font_name', 'None') if cached_format else 'None'
                            print(f"替换占位符: {placeholder_pattern} -> '{new_content}' - 字体:{font_name}, 大小:{font_size}, 颜色:{font_color}")
                            
                            return True
            
            # 如果runs级别替换失败，使用shape级别替换作为备用
            current_text = shape.text if hasattr(shape, 'text') else ""
            if placeholder_pattern in current_text:
                updated_text = current_text.replace(placeholder_pattern, new_content, 1)
                shape.text = updated_text
                
                # 对整个文本框应用格式
                if cached_format and hasattr(shape, 'text_frame') and shape.text_frame:
                    self._apply_format_to_shape_text(shape, cached_format, new_content)
                    # 输出替换和格式应用信息
                    font_size = cached_format.get('font_size', 'None')
                    font_color = cached_format.get('font_color', 'None')
                    print(f"替换占位符: {placeholder_pattern} -> '{new_content}' - 字体:{cached_format.get('font_name', 'None')}, 大小:{font_size}, 颜色:{font_color}")
                else:
                    print(f"替换占位符: {placeholder_pattern} -> '{new_content}'")
                
                return True
            
            return False
            
        except Exception as e:
            print(f"Shape级占位符替换失败: {e}")
            return False
    
    def _apply_format_to_cell(self, cell, format_info: Dict[str, Any]):
        """应用格式到表格单元格"""
        try:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            if hasattr(cell, 'text_frame') and cell.text_frame:
                for paragraph in cell.text_frame.paragraphs:
                    font = paragraph.font
                    if format_info.get('font_name'):
                        font.name = format_info['font_name']
                    if format_info.get('font_size'):
                        font.size = Pt(format_info['font_size'])
                    if format_info.get('font_bold') is not None:
                        font.bold = format_info['font_bold']
                    if format_info.get('font_italic') is not None:
                        font.italic = format_info['font_italic']
                    if format_info.get('font_color'):
                        try:
                            color_str = format_info['font_color']
                            if color_str.startswith('theme_'):
                                theme_color_id = int(color_str.replace('theme_', ''))
                                font.color.theme_color = theme_color_id
                            elif color_str and len(color_str) == 6 and all(c in '0123456789ABCDEFabcdef' for c in color_str):
                                # 应用十六进制颜色
                                r = int(color_str[0:2], 16)
                                g = int(color_str[2:4], 16)
                                b = int(color_str[4:6], 16)
                                font.color.rgb = RGBColor(r, g, b)
                            elif 'RGB(' in color_str:
                                rgb_values = color_str.replace('RGB(', '').replace(')', '').split(', ')
                                if len(rgb_values) == 3:
                                    r, g, b = map(int, rgb_values)
                                    font.color.rgb = RGBColor(r, g, b)
                        except Exception:
                            pass
        except Exception as e:
            print(f"应用单元格格式失败: {e}")
    
    def _apply_format_to_run(self, run, format_info: Dict[str, Any]):
        """应用格式到run"""
        try:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            font = run.font
            print(f"      应用格式 - 字体:{format_info.get('font_name')}, 大小:{format_info.get('font_size')}, 颜色:{format_info.get('font_color')}")
            
            if format_info.get('font_name'):
                font.name = format_info['font_name']
                print(f"      设置字体名称: {format_info['font_name']}")
            if format_info.get('font_size'):
                font.size = Pt(format_info['font_size'])
                print(f"      设置字体大小: {format_info['font_size']}pt")
            if format_info.get('font_bold') is not None:
                font.bold = format_info['font_bold']
            if format_info.get('font_italic') is not None:
                font.italic = format_info['font_italic']
            if format_info.get('font_color'):
                try:
                    color_str = format_info['font_color']
                    print(f"      尝试应用颜色: {color_str}")
                    if color_str.startswith('theme_'):
                        # 应用主题颜色
                        theme_color_id = int(color_str.replace('theme_', ''))
                        font.color.theme_color = theme_color_id
                        print(f"      应用主题颜色: {theme_color_id}")
                    elif color_str and len(color_str) == 6 and all(c in '0123456789ABCDEFabcdef' for c in color_str):
                        # 应用十六进制颜色
                        r = int(color_str[0:2], 16)
                        g = int(color_str[2:4], 16)
                        b = int(color_str[4:6], 16)
                        font.color.rgb = RGBColor(r, g, b)
                        print(f"      应用十六进制颜色: #{color_str} = RGB({r},{g},{b})")
                    elif 'RGB(' in color_str:
                        # 应用RGB颜色
                        rgb_values = color_str.replace('RGB(', '').replace(')', '').split(', ')
                        if len(rgb_values) == 3:
                            r, g, b = map(int, rgb_values)
                            font.color.rgb = RGBColor(r, g, b)
                            print(f"      应用RGB颜色: RGB({r},{g},{b})")
                except Exception as e:
                    print(f"      颜色应用失败: {e}")
        except Exception as e:
            print(f"应用run格式失败: {e}")
    
    def _apply_format_to_shape_text(self, shape, format_info: Dict[str, Any], new_content: str):
        """应用格式到文本框中替换的内容"""
        try:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # 直接对整个文本框的所有runs应用格式（因为shape.text替换会重建runs结构）
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 对所有runs应用格式，确保替换后的内容有正确格式
                        font = run.font
                        if format_info.get('font_name'):
                            font.name = format_info['font_name']
                        if format_info.get('font_size'):
                            font.size = Pt(format_info['font_size'])
                        if format_info.get('font_bold') is not None:
                            font.bold = format_info['font_bold']
                        if format_info.get('font_italic') is not None:
                            font.italic = format_info['font_italic']
                        if format_info.get('font_color'):
                            try:
                                color_str = format_info['font_color']
                                if color_str.startswith('theme_'):
                                    theme_color_id = int(color_str.replace('theme_', ''))
                                    font.color.theme_color = theme_color_id
                                elif color_str and len(color_str) == 6 and all(c in '0123456789ABCDEFabcdef' for c in color_str):
                                    # 应用十六进制颜色
                                    r = int(color_str[0:2], 16)
                                    g = int(color_str[2:4], 16)
                                    b = int(color_str[4:6], 16)
                                    font.color.rgb = RGBColor(r, g, b)
                                elif 'RGB(' in color_str:
                                    rgb_values = color_str.replace('RGB(', '').replace(')', '').split(', ')
                                    if len(rgb_values) == 3:
                                        r, g, b = map(int, rgb_values)
                                        font.color.rgb = RGBColor(r, g, b)
                            except Exception:
                                pass
        except Exception as e:
            print(f"应用文本框格式失败: {e}")

    def _apply_cached_format_to_shape(self, shape, cached_format: Dict[str, Any]):
        """将缓存的格式应用到shape的所有文本"""
        try:
            from pptx.util import Pt
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    font = paragraph.font
                    if cached_format.get('font_name'):
                        font.name = cached_format['font_name']
                    if cached_format.get('font_size'):
                        font.size = Pt(cached_format['font_size'])
                        print(f"      应用字体大小: {cached_format['font_size']}pt")
                    if cached_format.get('font_bold') is not None:
                        font.bold = cached_format['font_bold']
                    if cached_format.get('font_italic') is not None:
                        font.italic = cached_format['font_italic']
        except Exception as e:
            print(f"应用缓存格式失败: {e}")
    
    def _apply_text_with_cached_format(self, shape, text: str, format_info: Dict[str, Any]) -> bool:
        """使用缓存的格式信息应用文本"""
        try:
            text_frame = shape.text_frame
            
            # 保持文本框边距设置
            if format_info.get('margin_left') is not None:
                text_frame.margin_left = format_info['margin_left']
            if format_info.get('margin_right') is not None:
                text_frame.margin_right = format_info['margin_right']
            if format_info.get('margin_top') is not None:
                text_frame.margin_top = format_info['margin_top']
            if format_info.get('margin_bottom') is not None:
                text_frame.margin_bottom = format_info['margin_bottom']
            if format_info.get('vertical_anchor') is not None:
                text_frame.vertical_anchor = format_info['vertical_anchor']
            
            # 不清空整个text_frame，而是直接替换文本来更好地保持格式
            if len(text_frame.paragraphs) > 0:
                # 直接替换第一个段落的文本
                paragraph = text_frame.paragraphs[0]
                paragraph.text = text
            else:
                # 如果没有段落，则创建一个
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.text = text
            
            # 应用段落格式
            if format_info.get('paragraph_alignment') is not None:
                paragraph.alignment = format_info['paragraph_alignment']
            
            # 应用字体格式到段落的font对象
            font = paragraph.font
            
            print(f"   🎨 应用缓存格式 - 原有Runs: {format_info.get('has_runs', False)}, Runs数: {format_info.get('runs_count', 0)}")
            
            applied_changes = []
            
            if format_info.get('font_name'):
                font.name = format_info['font_name']
                applied_changes.append(f"字体:{format_info['font_name']}")
            
            if format_info.get('font_size') is not None:
                font.size = format_info['font_size']
                applied_changes.append(f"大小:{format_info['font_size']}")
            elif font.size is None:
                # 如果原来没有大小设置，给个默认值
                font.size = Pt(16)
                applied_changes.append("大小:默认16pt")
            
            if format_info.get('font_bold') is not None:
                font.bold = format_info['font_bold']
                applied_changes.append(f"粗体:{format_info['font_bold']}")
                
            if format_info.get('font_italic') is not None:
                font.italic = format_info['font_italic']
                applied_changes.append(f"斜体:{format_info['font_italic']}")
            
            if format_info.get('font_color') is not None:
                try:
                    font.color.rgb = format_info['font_color']
                    applied_changes.append("颜色:已应用")
                except Exception:
                    applied_changes.append("颜色:应用失败")
            
            print(f"   缓存格式应用完成 - {', '.join(applied_changes) if applied_changes else '无格式变更'}")
            
            # 确保run级别的格式也正确
            if paragraph.runs:
                for run in paragraph.runs:
                    run_font = run.font
                    if format_info.get('font_name'):
                        run_font.name = format_info['font_name']
                    if format_info.get('font_size') is not None:
                        run_font.size = format_info['font_size']
                    if format_info.get('font_bold') is not None:
                        run_font.bold = format_info['font_bold']
                    if format_info.get('font_italic') is not None:
                        run_font.italic = format_info['font_italic']
                    if format_info.get('font_color') is not None:
                        try:
                            run_font.color.rgb = format_info['font_color']
                        except Exception:
                            # 如果设置颜色失败，忽略颜色设置
                            pass
            
            return True
            
        except Exception as e:
            print(f"应用缓存格式时出错: {str(e)}")
            return False
    
    def _extract_placeholder_format(self, container, placeholder_name: str) -> Dict[str, Any]:
        """以占位符为单元提取格式信息 - 通用文本扫描"""
        format_info = {
            'font_name': None,
            'font_size': None,
            'font_bold': False,
            'font_italic': False,
            'font_color': None
        }
        
        def scan_text_frame(text_frame):
            """扫描任何text_frame中的占位符格式"""
            placeholder_pattern = f"{{{placeholder_name}}}"
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder_pattern in run.text:
                        # 找到包含占位符的run，提取其格式
                        font = run.font
                        font_size = font.size.pt if font.size else None
                        font_color = None
                        if font.color:
                            try:
                                if hasattr(font.color, 'rgb') and font.color.rgb:
                                    font_color = str(font.color.rgb)
                                elif hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                                    font_color = f"theme_{font.color.theme_color}"
                            except Exception:
                                pass
                        return {
                            'font_name': font.name,
                            'font_size': font_size,
                            'font_bold': font.bold,
                            'font_italic': font.italic,
                            'font_color': font_color
                        }
            # 如果没找到，返回第一个run的格式作为默认
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                first_run = text_frame.paragraphs[0].runs[0]
                font = first_run.font
                font_size = font.size.pt if font.size else None
                font_color = None
                if font.color:
                    try:
                        if hasattr(font.color, 'rgb') and font.color.rgb:
                            font_color = str(font.color.rgb)
                        elif hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                            font_color = f"theme_{font.color.theme_color}"
                    except Exception:
                        pass
                return {
                    'font_name': font.name,
                    'font_size': font_size,
                    'font_bold': font.bold,
                    'font_italic': font.italic,
                    'font_color': font_color
                }
            return format_info
        
        try:
            # 通用扫描：任何有text_frame的对象
            if hasattr(container, 'text_frame') and container.text_frame:
                format_info.update(scan_text_frame(container.text_frame))
            # 如果是表格等复杂对象，递归扫描
            elif hasattr(container, 'table'):
                # 表格：扫描所有单元格
                for row in container.table.rows:
                    for cell in row.cells:
                        if hasattr(cell, 'text_frame') and cell.text_frame:
                            result = scan_text_frame(cell.text_frame)
                            if result['font_name'] or result['font_size']:
                                format_info.update(result)
                                break
            
        except Exception as e:
            print(f"提取占位符{placeholder_name}格式失败: {e}")
        
        return format_info

    def _extract_text_format(self, shape) -> Dict[str, Any]:
        """提取文本框的格式信息"""
        format_info = {
            'font_name': None,
            'font_size': None,
            'font_bold': False,
            'font_italic': False,
            'font_color': None,
            'paragraph_alignment': None,
            'vertical_anchor': None,
            'margin_left': None,
            'margin_right': None,
            'margin_top': None,
            'margin_bottom': None,
            'shape_type': None,  # 新增：形状类型
            'has_runs': False,   # 新增：是否有runs
            'runs_count': 0      # 新增：runs数量
        }
        
        try:
            # 记录形状类型用于调试
            if hasattr(shape, 'shape_type'):
                format_info['shape_type'] = str(shape.shape_type)
            
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text_frame = shape.text_frame
                
                # 提取文本框边距和垂直对齐
                format_info['margin_left'] = text_frame.margin_left
                format_info['margin_right'] = text_frame.margin_right  
                format_info['margin_top'] = text_frame.margin_top
                format_info['margin_bottom'] = text_frame.margin_bottom
                format_info['vertical_anchor'] = text_frame.vertical_anchor
                
                
                # 从第一个段落提取格式
                if text_frame.paragraphs:
                    first_paragraph = text_frame.paragraphs[0]
                    format_info['paragraph_alignment'] = first_paragraph.alignment
                    
                    # 尝试获取段落字体信息作为备用
                    try:
                        paragraph_font = first_paragraph.font
                        if not format_info['font_name'] and paragraph_font.name:
                            format_info['font_name'] = paragraph_font.name
                        if not format_info['font_size'] and paragraph_font.size:
                            format_info['font_size'] = paragraph_font.size
                        if format_info['font_bold'] is False and paragraph_font.bold is not None:
                            format_info['font_bold'] = paragraph_font.bold
                        if format_info['font_italic'] is False and paragraph_font.italic is not None:
                            format_info['font_italic'] = paragraph_font.italic
                    except Exception:
                        pass
                    
                    # 从第一个运行提取字体格式
                    if first_paragraph.runs:
                        format_info['has_runs'] = True
                        format_info['runs_count'] = len(first_paragraph.runs)
                        
                        first_run = first_paragraph.runs[0]
                        font = first_run.font
                        
                        format_info['font_name'] = font.name
                        format_info['font_size'] = font.size
                        format_info['font_bold'] = font.bold
                        format_info['font_italic'] = font.italic
                        
                        
                        # 特殊处理：如果runs中没有字体信息，尝试从其他runs获取
                        if not font.name or not font.size:
                            for run in first_paragraph.runs[1:]:
                                if not font.name and run.font.name:
                                    format_info['font_name'] = run.font.name
                                if not font.size and run.font.size:
                                    format_info['font_size'] = run.font.size
                        
                        # 提取字体颜色
                        if font.color:
                            try:
                                if hasattr(font.color, 'rgb') and font.color.rgb:
                                    format_info['font_color'] = font.color.rgb
                                elif hasattr(font.color, 'theme_color'):
                                    # 主题颜色，保持None让系统使用默认颜色
                                    format_info['font_color'] = None
                            except Exception:
                                format_info['font_color'] = None
                    else:
                        # 如果没有runs，从段落字体获取
                        format_info['has_runs'] = False
                        print(f"   ⚠️ 无Runs，使用段落格式")
                        
                        font = first_paragraph.font
                        format_info['font_name'] = font.name
                        format_info['font_size'] = font.size
                        format_info['font_bold'] = font.bold
                        format_info['font_italic'] = font.italic
                        
                        print(f"   📄 段落格式 - 字体: {font.name}, 大小: {font.size}, 粗体: {font.bold}, 斜体: {font.italic}")
                        
                        if font.color:
                            try:
                                if hasattr(font.color, 'rgb') and font.color.rgb:
                                    format_info['font_color'] = font.color.rgb
                                elif hasattr(font.color, 'theme_color'):
                                    # 主题颜色，保持None让系统使用默认颜色
                                    format_info['font_color'] = None
                            except Exception:
                                format_info['font_color'] = None
                            
        except Exception as e:
            print(f"提取文本格式时出错: {str(e)}")
        
        return format_info
    
    def _apply_text_with_format(self, shape, text: str, format_info: Dict[str, Any]) -> bool:
        """应用文本并保持格式"""
        try:
            text_frame = shape.text_frame
            
            # 保持文本框边距设置
            if format_info['margin_left'] is not None:
                text_frame.margin_left = format_info['margin_left']
            if format_info['margin_right'] is not None:
                text_frame.margin_right = format_info['margin_right']
            if format_info['margin_top'] is not None:
                text_frame.margin_top = format_info['margin_top']
            if format_info['margin_bottom'] is not None:
                text_frame.margin_bottom = format_info['margin_bottom']
            if format_info['vertical_anchor'] is not None:
                text_frame.vertical_anchor = format_info['vertical_anchor']
            
            # 不清空整个text_frame，而是直接替换文本来更好地保持格式
            if len(text_frame.paragraphs) > 0:
                # 直接替换第一个段落的文本
                paragraph = text_frame.paragraphs[0]
                paragraph.text = text
            else:
                # 如果没有段落，则创建一个
                text_frame.clear()
                paragraph = text_frame.paragraphs[0]
                paragraph.text = text
            
            # 应用段落格式
            if format_info['paragraph_alignment'] is not None:
                paragraph.alignment = format_info['paragraph_alignment']
            
            # 应用字体格式到段落的font对象
            font = paragraph.font
            
            print(f"   🎨 应用格式 - 原有Runs: {format_info['has_runs']}, Runs数: {format_info['runs_count']}")
            
            applied_changes = []
            
            if format_info['font_name']:
                font.name = format_info['font_name']
                applied_changes.append(f"字体:{format_info['font_name']}")
            
            if format_info['font_size'] is not None:
                font.size = format_info['font_size']
                applied_changes.append(f"大小:{format_info['font_size']}")
            elif font.size is None:
                # 如果原来没有大小设置，给个默认值
                font.size = Pt(16)
                applied_changes.append("大小:默认16pt")
            
            if format_info['font_bold'] is not None:
                font.bold = format_info['font_bold']
                applied_changes.append(f"粗体:{format_info['font_bold']}")
                
            if format_info['font_italic'] is not None:
                font.italic = format_info['font_italic']
                applied_changes.append(f"斜体:{format_info['font_italic']}")
            
            if format_info['font_color'] is not None:
                try:
                    font.color.rgb = format_info['font_color']
                    applied_changes.append("颜色:已应用")
                except Exception:
                    applied_changes.append("颜色:应用失败")
            
            print(f"   应用完成 - {', '.join(applied_changes) if applied_changes else '无格式变更'}")
            
            # 确保run级别的格式也正确
            if paragraph.runs:
                for run in paragraph.runs:
                    run_font = run.font
                    if format_info['font_name']:
                        run_font.name = format_info['font_name']
                    if format_info['font_size'] is not None:
                        run_font.size = format_info['font_size']
                    if format_info['font_bold'] is not None:
                        run_font.bold = format_info['font_bold']
                    if format_info['font_italic'] is not None:
                        run_font.italic = format_info['font_italic']
                    if format_info['font_color'] is not None:
                        try:
                            run_font.color.rgb = format_info['font_color']
                        except Exception:
                            # 如果设置颜色失败，忽略颜色设置
                            pass
            
            return True
            
        except Exception as e:
            print(f"应用格式时出错: {str(e)}")
            return False
    
    def _update_slide_content(self, slide, content: str):
        """更新幻灯片内容"""
        # 查找可用的文本框
        text_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text_shapes.append(shape)
        
        if text_shapes:
            # 使用最后一个可用的文本框（通常是主要内容区域）
            target_shape = text_shapes[-1] if len(text_shapes) > 1 else text_shapes[0]
            
            # 清空现有内容并添加新内容
            tf = target_shape.text_frame
            tf.clear()
            
            # 添加内容
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(16)
    
    def _add_new_slide(self, title: str, content: str):
        """添加新幻灯片"""
        # 使用标题和内容布局
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # 设置内容
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()
            
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(16)

class FileManager:
    """文件管理器"""
    
    @staticmethod
    def save_ppt_to_bytes(presentation: Presentation) -> bytes:
        """
        将PPT保存为字节数据
        
        Args:
            presentation: PPT演示文稿对象
            
        Returns:
            bytes: PPT文件的字节数据
        """
        config = get_config()
        
        # 创建临时文件
        timestamp = str(int(time.time() * 1000))
        temp_filename = f"temp_ppt_{timestamp}.pptx"
        temp_filepath = os.path.join(config.temp_output_dir, temp_filename)
        
        try:
            # 保存文件
            presentation.save(temp_filepath)
            
            # 读取字节数据
            with open(temp_filepath, 'rb') as f:
                ppt_bytes = f.read()
            
            return ppt_bytes
        finally:
            # 清理临时文件
            try:
                if os.path.exists(temp_filepath):
                    os.remove(temp_filepath)
            except Exception:
                pass
    
    @staticmethod
    def save_ppt_to_file(presentation: Presentation, filename: str = None) -> str:
        """
        将PPT保存到文件
        
        Args:
            presentation: PPT演示文稿对象
            filename: 文件名（可选）
            
        Returns:
            str: 保存的文件路径
        """
        config = get_config()
        
        if not filename:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"updated_ppt_{timestamp}.pptx"
        
        filepath = os.path.join(config.output_dir, filename)
        presentation.save(filepath)
        return filepath
    
    @staticmethod
    def validate_ppt_file(file_path: str) -> Tuple[bool, str]:
        """
        验证PPT文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            Tuple[bool, str]: (是否有效, 错误信息)
        """
        if not os.path.exists(file_path):
            return False, f"文件不存在: {file_path}"
        
        if not file_path.lower().endswith('.pptx'):
            return False, "文件格式不支持，请使用.pptx格式"
        
        try:
            # 尝试打开文件
            presentation = Presentation(file_path)
            if len(presentation.slides) == 0:
                return False, "PPT文件为空"
            return True, ""
        except Exception as e:
            return False, f"文件损坏或格式错误: {e}"

def format_timestamp(timestamp: float = None) -> str:
    """
    格式化时间戳
    
    Args:
        timestamp: 时间戳（可选，默认当前时间）
        
    Returns:
        str: 格式化的时间字符串
    """
    if timestamp is None:
        timestamp = time.time()
    return datetime.fromtimestamp(timestamp).strftime('%Y-%m-%d %H:%M:%S')

def sanitize_filename(filename: str) -> str:
    """
    清理文件名，移除非法字符
    
    Args:
        filename: 原始文件名
        
    Returns:
        str: 清理后的文件名
    """
    # 移除或替换非法字符
    filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
    # 移除开头和结尾的空白字符和点
    filename = filename.strip('. ')
    # 如果文件名为空，使用默认名称
    if not filename:
        filename = 'untitled'
    return filename

def is_valid_api_key(api_key: str) -> bool:
    """
    验证API密钥格式
    
    Args:
        api_key: API密钥
        
    Returns:
        bool: 是否有效
    """
    if not api_key:
        return False
    
    # 简单验证：支持OpenAI (sk-) 和OpenRouter (sk-or-) 格式
    return (api_key.startswith('sk-or-') or api_key.startswith('sk-')) and len(api_key) > 20